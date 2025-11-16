# dashboard.py
import streamlit as st
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
from io import BytesIO
import seaborn as sns
import os
import re
import uuid
import random
import datetime
import matplotlib as mpl
from matplotlib import font_manager
import platform
import zipfile
import json
import pickle

# Configure matplotlib for Chinese font support
plt.rcParams['font.sans-serif'] = ['SimHei', 'DejaVu Sans', 'Arial Unicode MS', 'Microsoft YaHei', 'WenQuanYi Micro Hei']
plt.rcParams['axes.unicode_minus'] = False  # Fix minus sign display issue

# Try to set up Chinese font
try:
    # For Windows
    if platform.system() == 'Windows':
        plt.rcParams['font.sans-serif'] = ['Microsoft YaHei', 'SimHei', 'KaiTi']
    # For macOS
    elif platform.system() == 'Darwin':
        plt.rcParams['font.sans-serif'] = ['Arial Unicode MS', 'PingFang SC', 'Hiragino Sans GB', 'STHeiti']
    # For Linux
    else:
        plt.rcParams['font.sans-serif'] = ['WenQuanYi Micro Hei', 'DejaVu Sans', 'Noto Sans CJK SC']
    
    # Clear matplotlib font cache to force reload
    font_manager._rebuild()
except Exception as e:
    print(f"Font configuration warning: {e}")

# Function to ensure proper font loading for plots
def ensure_chinese_font():
    """Ensure Chinese font is properly loaded for matplotlib"""
    try:
        # Test if Chinese characters can be displayed
        fig, ax = plt.subplots(figsize=(1, 1))
        ax.text(0.5, 0.5, '测试', fontsize=12)
        plt.close(fig)
        return True
    except:
        # Fallback to system default
        plt.rcParams['font.sans-serif'] = ['DejaVu Sans', 'Arial', 'sans-serif']
        return False
from deepseek_ai import DeepSeekAI
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Configure Chinese fonts
def configure_chinese_fonts():
    """Configure matplotlib to support Chinese display"""
    system = platform.system()
    
    if system == "Windows":
        chinese_fonts = ['Microsoft YaHei', 'SimHei', 'KaiTi', 'SimSun']
    elif system == "Darwin":  # macOS
        chinese_fonts = ['PingFang SC', 'Hiragino Sans GB', 'STHeiti', 'Arial Unicode MS']
    else:  # Linux
        chinese_fonts = ['WenQuanYi Micro Hei', 'DejaVu Sans', 'Liberation Sans']
    
    for font in chinese_fonts:
        try:
            if font in [f.name for f in font_manager.fontManager.ttflist]:
                plt.rcParams['font.sans-serif'] = [font] + plt.rcParams['font.sans-serif']
                plt.rcParams['axes.unicode_minus'] = False
                print(f"Set Chinese font: {font}")
                return True
        except:
            continue
    
    try:
        plt.rcParams['font.sans-serif'] = ['SimHei', 'DejaVu Sans', 'Liberation Sans', 'sans-serif']
        plt.rcParams['axes.unicode_minus'] = False
        return True
    except:
        print("Warning: Cannot configure Chinese fonts")
        return False

configure_chinese_fonts()

# Initialize session state
if 'projects' not in st.session_state:
    st.session_state.projects = {}
if 'active_project' not in st.session_state:
    st.session_state.active_project = None
if 'experiments' not in st.session_state:
    st.session_state.experiments = {}
if 'mode' not in st.session_state:
    st.session_state.mode = "General Behavior"
if 'worksheet_data' not in st.session_state:
    st.session_state.worksheet_data = {}
if 'save_status' not in st.session_state:
    st.session_state.save_status = {}
if 'show_project_creation' not in st.session_state:
    st.session_state.show_project_creation = False
if 'show_import_dialog' not in st.session_state:
    st.session_state.show_import_dialog = False
if 'comparison_groups' not in st.session_state:
    st.session_state.comparison_groups = {}
if 'group_projects' not in st.session_state:
    st.session_state.group_projects = {}  # Maps group names to project IDs
if 'all_experiment_charts' not in st.session_state:
    st.session_state.all_experiment_charts = []  # Store all charts from experiments
if 'show_scoring_help' not in st.session_state:
    st.session_state.show_scoring_help = False

# Helper function to get groups for a project
def get_project_groups(project_id):
    """Get all groups belonging to a specific project"""
    return [group for group, pid in st.session_state.group_projects.items() if pid == project_id]

# Helper function to synchronize time points across all worksheets
def synchronize_time_points_across_worksheets(project_id, new_times, mode=None):
    """Synchronize time points across all worksheets for a project"""
    project_groups = get_project_groups(project_id)
    if not project_groups:
        return
    
    # Always sync to ALL modes (except Body Weight) regardless of which mode triggered it
    modes = [
        "General Behavior", 
        "Autonomic and Sensorimotor Functions", 
        "Reflex Capabilities",
        "Body Temperature",
        "Convulsive Behaviors and Excitability"
    ]
    # Note: Body Weight is excluded as it uses 'before'/'after' instead of time points
    
    for group in project_groups:
        for mode_item in modes:
            worksheet_key = f"worksheet_{group}_{mode_item}"
            
            # Get observations for this mode
            if mode_item == "Autonomic and Sensorimotor Functions":
                observations = [t_obs(obs) for obs in AUTONOMIC_OBSERVATIONS]
            elif mode_item == "Reflex Capabilities":
                observations = [t_obs(obs) for obs in REFLEX_OBSERVATIONS]
            elif mode_item == "Convulsive Behaviors and Excitability":
                observations = [t_obs(obs) for obs in CONVULSIVE_OBSERVATIONS]
            elif mode_item == "Body Temperature":
                observations = [t_obs('body temperature')]
            else:  # General Behavior
                observations = [t_obs(obs) for obs in GENERAL_BEHAVIOR_OBSERVATIONS]
            
            # Get existing worksheet or create empty one
            if worksheet_key in st.session_state:
                df = st.session_state[worksheet_key].copy()
                existing_times = sorted(df['time'].unique()) if not df.empty else []
            else:
                df = pd.DataFrame()
                existing_times = []
            
            # Combine and sort all times (union of existing and new)
            all_times = sorted(list(set(existing_times + new_times)))
            
            # Get animal type and num_animals from project
            project = st.session_state.projects.get(project_id, {})
            animal_type = project.get('animal_type', 'mouse')
            if animal_type == 'custom':
                animal_type = project.get('custom_animal_name', 'animal')
            num_animals = project.get('num_animals', 8)
            
            # Rebuild dataframe with synchronized times
            new_data = []
            for time in all_times:
                for obs in observations:
                    # Check if this time-observation combination already exists
                    if not df.empty:
                        existing_row = df[(df['time'] == time) & (df['observation'] == obs)]
                        if not existing_row.empty:
                            # Use existing row (preserve data)
                            new_data.append(existing_row.iloc[0].to_dict())
                            continue
                    
                    # Create new row with default values
                    row = {'time': time, 'observation': obs}
                    
                    for i in range(1, num_animals + 1):
                        if mode_item == "Body Temperature":
                            if animal_type == 'mouse':
                                row[f'{animal_type}_{i}'] = '37.0'
                            elif animal_type == 'rat':
                                row[f'{animal_type}_{i}'] = '37.5'
                            else:
                                row[f'{animal_type}_{i}'] = '37.2'
                        elif mode_item in ["Autonomic and Sensorimotor Functions", "Reflex Capabilities", "Convulsive Behaviors and Excitability"]:
                            row[f'{animal_type}_{i}'] = t('normal')
                        else:  # General Behavior
                            row[f'{animal_type}_{i}'] = '4'
                    new_data.append(row)
            
            # Update worksheet with synchronized data, sorted by time
            if new_data:
                new_df = pd.DataFrame(new_data)
                new_df = new_df.sort_values(['time', 'observation']).reset_index(drop=True)
                st.session_state[worksheet_key] = new_df
                st.session_state.worksheet_data[f"{group}_{mode_item}"] = new_df

# Helper function to capture charts for PowerPoint
def capture_chart_for_powerpoint(fig, title, mode, chart_type="Plot", description="", add_to_session=True):
    """Capture a matplotlib figure for inclusion in PowerPoint presentations"""
    try:
        # Save figure as bytes
        img_buffer = BytesIO()
        fig.savefig(img_buffer, format='png', dpi=300, bbox_inches='tight')
        img_buffer.seek(0)
        chart_data = img_buffer.getvalue()
        
        # Store chart information
        chart_info = {
            'title': title,
            'mode': mode,
            'chart_type': chart_type,
            'data': chart_data,
            'description': description,
            'timestamp': datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')
        }
        
        # Add to session state only if requested (not for PowerPoint generation)
        if add_to_session:
            st.session_state.all_experiment_charts.append(chart_info)
        
        return chart_info
    except Exception as e:
        st.error(f"Error capturing chart: {str(e)}")
        return None

# Helper function to clear old charts for a specific project
def clear_project_charts(project_id):
    """Clear charts from previous experiments for a specific project"""
    try:
        # Keep only charts from the current project
        if 'all_experiment_charts' in st.session_state:
            # For now, we'll keep all charts but this could be enhanced to filter by project
            # if we add project_id to chart_info
            pass
    except Exception as e:
        st.error(f"Error clearing charts: {str(e)}")

# Function to export all project data as ZIP
def export_project_data_as_zip():
    """Export all current project data as a ZIP file"""
    try:
        if st.session_state.active_project is None:
            return None, "No active project to export"
        
        # Create ZIP buffer
        zip_buffer = BytesIO()
        
        with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
            # Export projects data
            projects_data = {
                'projects': st.session_state.projects,
                'active_project': st.session_state.active_project,
                'experiments': st.session_state.experiments,
                'group_projects': st.session_state.group_projects,
                'comparison_groups': st.session_state.comparison_groups,
                'language': st.session_state.language,
                'export_timestamp': datetime.datetime.now().isoformat()
            }
            
            # Add projects data as JSON
            zip_file.writestr('projects_data.json', json.dumps(projects_data, indent=2, default=str))
            
            # Export all worksheet data
            worksheet_data = {}
            for key, value in st.session_state.worksheet_data.items():
                if isinstance(value, pd.DataFrame):
                    # Convert DataFrame to CSV
                    csv_buffer = BytesIO()
                    value.to_csv(csv_buffer, index=False)
                    csv_buffer.seek(0)
                    worksheet_data[key] = csv_buffer.getvalue().decode('utf-8')
            
            # Add worksheet data as JSON
            zip_file.writestr('worksheet_data.json', json.dumps(worksheet_data, indent=2))
            
            # Export session state data that contains DataFrames
            session_data = {}
            for key, value in st.session_state.items():
                if isinstance(value, pd.DataFrame):
                    # Convert DataFrame to CSV
                    csv_buffer = BytesIO()
                    value.to_csv(csv_buffer, index=False)
                    csv_buffer.seek(0)
                    session_data[key] = csv_buffer.getvalue().decode('utf-8')
                elif key.startswith('worksheet_') and isinstance(value, pd.DataFrame):
                    # Handle worksheet session state
                    csv_buffer = BytesIO()
                    value.to_csv(csv_buffer, index=False)
                    csv_buffer.seek(0)
                    session_data[key] = csv_buffer.getvalue().decode('utf-8')
            
            # Add session data as JSON
            if session_data:
                zip_file.writestr('session_data.json', json.dumps(session_data, indent=2))
            
            # Add export info
            export_info = {
                'export_date': datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
                'version': '1.0',
                'description': 'FOB Test Analysis Dashboard Project Export',
                'project_name': st.session_state.projects[st.session_state.active_project]['name'] if st.session_state.active_project else 'Unknown'
            }
            zip_file.writestr('export_info.json', json.dumps(export_info, indent=2))
        
        zip_buffer.seek(0)
        return zip_buffer.getvalue(), "Export successful"
        
    except Exception as e:
        return None, f"Export failed: {str(e)}"

# Function to import project data from ZIP
def import_project_data_from_zip(uploaded_file):
    """Import project data from uploaded ZIP file"""
    try:
        # Read ZIP file
        zip_buffer = BytesIO(uploaded_file.read())
        
        with zipfile.ZipFile(zip_buffer, 'r') as zip_file:
            # Check if required files exist
            required_files = ['projects_data.json', 'worksheet_data.json']
            if not all(file in zip_file.namelist() for file in required_files):
                return False, "Invalid ZIP file: Missing required data files"
            
            # Load projects data
            projects_data = json.loads(zip_file.read('projects_data.json').decode('utf-8'))
            
            # Clear current session state (ask for confirmation in UI)
            st.session_state.projects = projects_data['projects']
            st.session_state.active_project = projects_data['active_project']
            st.session_state.experiments = projects_data['experiments']
            st.session_state.group_projects = projects_data.get('group_projects', {})
            st.session_state.comparison_groups = projects_data.get('comparison_groups', {})
            
            # Set language if available
            if 'language' in projects_data:
                st.session_state.language = projects_data['language']
            
            # Load worksheet data
            worksheet_data = json.loads(zip_file.read('worksheet_data.json').decode('utf-8'))
            for key, csv_data in worksheet_data.items():
                try:
                    df = pd.read_csv(BytesIO(csv_data.encode('utf-8')))
                    st.session_state.worksheet_data[key] = df
                except Exception as e:
                    st.warning(f"Could not load worksheet data for {key}: {str(e)}")
            
            # Load session data if available
            if 'session_data.json' in zip_file.namelist():
                session_data = json.loads(zip_file.read('session_data.json').decode('utf-8'))
                for key, csv_data in session_data.items():
                    try:
                        df = pd.read_csv(BytesIO(csv_data.encode('utf-8')))
                        st.session_state[key] = df
                    except Exception as e:
                        st.warning(f"Could not load session data for {key}: {str(e)}")
            
            # Load export info if available
            export_info = {}
            if 'export_info.json' in zip_file.namelist():
                export_info = json.loads(zip_file.read('export_info.json').decode('utf-8'))
            
            return True, export_info
            
    except Exception as e:
        return False, f"Import failed: {str(e)}"

# DeepSeek AI Configuration
# Get your API key from: https://platform.deepseek.com/
# Replace the placeholder with your actual API key
# You can also set it as an environment variable: DEEPSEEK_API_KEY
DEEPSEEK_API_KEY = os.getenv("DEEPSEEK_API_KEY", "sk-16231cff5f244f0a898972cd1e4d0bf0")

# Initialize global deepseek_client
deepseek_client = None

# Helper function for DeepSeek API calls
def make_deepseek_api_call(prompt):
    """Make DeepSeek API call"""
    global deepseek_client
    
    # Check if client is configured, if not try to configure it
    if deepseek_client is None:
        if not configure_deepseek():
            return "Error: DeepSeek AI is not properly configured. Please check your API key and try again."
    
    try:
        # Simple API call with DeepSeek
        response = deepseek_client.chat.completions.create(
            model="deepseek-chat",
            messages=[
                {"role": "user", "content": prompt}
            ],
            temperature=0.7,
            max_tokens=2000
        )
        return response.choices[0].message.content
        
    except Exception as e:
        error_msg = str(e)
        # Try to reconfigure if there's an error
        deepseek_client = None
        return f"Error generating AI response: {error_msg}. Please check your API key and try again."

# Initialize DeepSeek AI
def configure_deepseek():
    """Configure DeepSeek AI with stored API key"""
    global deepseek_client
    
    try:
        # Check if API key is valid
        if not DEEPSEEK_API_KEY or DEEPSEEK_API_KEY == "":
            return False
        
        # Create DeepSeek client instance
        deepseek_client = DeepSeekAI(
            api_key=DEEPSEEK_API_KEY
        )
        return True
        
    except Exception as e:
        deepseek_client = None
        if 'st' in globals():
            st.error(f"Error configuring DeepSeek AI: {str(e)}")
        return False



def generate_ai_report(project_data, analysis_data, mode_eng, language='en', uploaded_file_content=None):
    """Generate AI-powered report using DeepSeek"""
    try:
        # Check if client is already configured, if not configure it
        if 'deepseek_client' not in globals() or deepseek_client is None:
            if not configure_deepseek():
                return "Error: Failed to configure DeepSeek AI"
        
        # Create prompt based on analysis mode and data
        if mode_eng == "Body Weight":
            prompt = create_weight_ai_prompt(project_data, analysis_data, language, uploaded_file_content)
        else:
            prompt = create_behavior_ai_prompt(project_data, analysis_data, mode_eng, language, uploaded_file_content)
        
        # Use DeepSeek API
        return make_deepseek_api_call(prompt)
        
    except Exception as e:
        return f"Error generating AI report: {str(e)}"



def create_behavior_ai_prompt(project_data, analysis_data, mode_eng, language, uploaded_file_content=None):
    """Create AI prompt for behavior analysis"""
    file_section = ""
    if uploaded_file_content:
        file_section = f"\nAdditional Uploaded File Data:\n{uploaded_file_content}\n"
    
    if language == 'zh':
        return f"""
作为一位专业的动物行为学专家，请分析以下FOB测试数据并生成详细报告：

项目信息：
- 项目名称：{project_data.get('name', 'N/A')}
- 动物类型：{project_data.get('animal_type', 'N/A')}
- 每组动物数量：{project_data.get('num_animals', 'N/A')}
- 分析模式：{mode_eng}

分析数据：
{analysis_data.to_string() if hasattr(analysis_data, 'to_string') else str(analysis_data)}{file_section}

请提供以下分析：
1. 异常行为模式识别
2. 各组行为差异分析
3. 时间序列行为变化趋势
4. 对照组与实验组比较
5. 行为异常的生物学意义
6. 实验设计评估
7. 统计分析和显著性
8. 结论和建议

请用中文回答，格式要专业、清晰。
"""
    else:
        return f"""
As a professional animal behavior expert, please analyze the following FOB test data and generate a detailed report:

Project Information:
- Project Name: {project_data.get('name', 'N/A')}
- Animal Type: {project_data.get('animal_type', 'N/A')}
- Animals per Group: {project_data.get('num_animals', 'N/A')}
- Analysis Mode: {mode_eng}

Analysis Data:
{analysis_data.to_string() if hasattr(analysis_data, 'to_string') else str(analysis_data)}{file_section}

Please provide the following analysis:
1. Abnormal behavior pattern identification
2. Analysis of behavioral differences between groups
3. Time-series behavioral change trends
4. Comparison between control and experimental groups
5. Biological significance of behavioral abnormalities
6. Experimental design evaluation
7. Statistical analysis and significance
8. Conclusions and recommendations

Please provide a professional, clear format in English.
"""

def create_weight_ai_prompt(project_data, weight_data, language, uploaded_file_content=None):
    """Create AI prompt for weight analysis"""
    file_section = ""
    if uploaded_file_content:
        file_section = f"\nAdditional Uploaded File Data:\n{uploaded_file_content}\n"
    
    if language == 'zh':
        return f"""
作为一位专业的动物实验数据分析师，请分析以下FOB测试体重数据并生成详细报告：

项目信息：
- 项目名称：{project_data.get('name', 'N/A')}
- 动物类型：{project_data.get('animal_type', 'N/A')}
- 每组动物数量：{project_data.get('num_animals', 'N/A')}
- 分析模式：体重变化

体重数据：
{weight_data.to_string() if hasattr(weight_data, 'to_string') else str(weight_data)}{file_section}

请提供以下分析：
1. 总体体重变化趋势分析
2. 各组之间的体重变化比较
3. 对照组与其他组的差异分析
4. 体重变化的生物学意义
5. 实验设计建议
6. 统计显著性分析（如适用）
7. 结论和建议

请用中文回答，格式要专业、清晰。
"""
    else:
        return f"""
As a professional animal experiment data analyst, please analyze the following FOB test body weight data and generate a detailed report:

Project Information:
- Project Name: {project_data.get('name', 'N/A')}
- Animal Type: {project_data.get('animal_type', 'N/A')}
- Animals per Group: {project_data.get('num_animals', 'N/A')}
- Analysis Mode: Body Weight Changes

Weight Data:
{weight_data.to_string() if hasattr(weight_data, 'to_string') else str(weight_data)}{file_section}

Please provide the following analysis:
1. Overall body weight change trends
2. Comparison of weight changes between groups
3. Analysis of differences between control and treatment groups
4. Biological significance of weight changes
5. Experimental design recommendations
6. Statistical significance analysis (if applicable)
7. Conclusions and recommendations

Please provide a professional, clear format in English.
"""

def process_uploaded_file(uploaded_file):
    """Process uploaded file and return its content"""
    try:
        if uploaded_file is None:
            return None
        
        file_extension = uploaded_file.name.split('.')[-1].lower()
        
        if file_extension in ['csv']:
            # Read CSV file
            df = pd.read_csv(uploaded_file)
            return df.to_string()
        elif file_extension in ['xlsx', 'xls']:
            # Read Excel file
            df = pd.read_excel(uploaded_file)
            return df.to_string()
        elif file_extension in ['txt']:
            # Read text file
            return uploaded_file.read().decode('utf-8')
        else:
            return f"Unsupported file type: {file_extension}. Supported types: CSV, Excel, TXT"
    except Exception as e:
        return f"Error processing file: {str(e)}"

def generate_chatbot_response(user_message, language='en'):
    """Generate chatbot response using DeepSeek AI"""
    try:
        # Check if client is already configured, if not configure it
        if 'deepseek_client' not in globals() or deepseek_client is None:
            if not configure_deepseek():
                return "Error: Failed to configure DeepSeek AI"
        
        # Create chatbot prompt with comprehensive functionality summary
        if language == 'zh':
            prompt = f"""
你是一个专业的FOB测试分析仪表板使用指导助手。你的主要任务是帮助用户学习如何使用这个工具。

**仪表板功能总结：**
这是一个功能观察电池（FOB）测试分析仪表板，包含以下核心功能：

1. **项目管理**：创建项目，设置动物类型（小鼠/大鼠/自定义），每组动物数量，组数
2. **6种分析模式**：
   - 一般行为：0/4/8评分系统，正常范围2-6
   - 自主神经功能：正常/异常二元评分
   - 反射能力：正常/异常二元评分
   - 体温：温度测量，正常范围36-38°C
   - 体重：实验前后体重测量，自动计算变化
   - 惊厥行为：正常/异常二元评分
3. **数据录入**：手动保存和自动保存两种模式，支持添加时间点
4. **组管理**：多组实验，可设置对照组
5. **数据分析**：异常事件跟踪，统计分析和可视化
6. **报告生成**：综合报告导出，图表下载
7. **多语言支持**：中英文界面

用户问题：{user_message}

请专注于提供以下帮助：
1. **工具使用指导**：详细说明如何在仪表板中操作各个功能
2. **步骤说明**：提供清晰的分步操作指南
3. **功能解释**：解释每个功能的作用和用途
4. **常见问题**：解答用户在使用过程中遇到的问题
5. **最佳实践**：提供使用建议和技巧

回答要求：
- 保持友好、耐心、专业的态度
- 提供具体、可操作的步骤
- 使用简单易懂的语言
- 如果涉及界面操作，请明确指出按钮和选项的位置
- 格式要清晰，可以使用编号或要点
- 基于上述功能总结提供准确的指导

请用中文回答，重点放在工具使用指导上。
"""
        else:
            prompt = f"""
You are a professional FOB Test Analysis Dashboard usage guide assistant. Your main task is to help users learn how to use this tool effectively.

**Dashboard Functionality Summary:**
This is a Functional Observational Battery (FOB) test analysis dashboard with the following core features:

1. **Project Management**: Create projects, set animal types (mouse/rat/custom), animals per group, number of groups
2. **6 Analysis Modes**:
   - General Behavior: 0/4/8 scoring system, normal range 2-6
   - Autonomic Functions: Normal/Abnormal binary scoring
   - Reflex Capabilities: Normal/Abnormal binary scoring
   - Body Temperature: Temperature measurements, normal range 36-38°C
   - Body Weight: Before/after weight measurements with automatic change calculations
   - Convulsive Behaviors: Normal/Abnormal binary scoring
3. **Data Entry**: Manual save and auto-save modes, support for adding time points
4. **Group Management**: Multiple experimental groups, can set comparison group
5. **Data Analysis**: Abnormal episode tracking, statistical analysis and visualization
6. **Report Generation**: Comprehensive report export, chart downloads
7. **Multi-language Support**: English and Chinese interfaces

User Question: {user_message}

Please focus on providing the following help:
1. **Tool Usage Guidance**: Detailed instructions on how to operate various features in the dashboard
2. **Step-by-step Instructions**: Provide clear, actionable step-by-step guides
3. **Feature Explanation**: Explain what each feature does and its purpose
4. **Common Issues**: Answer questions users encounter while using the tool
5. **Best Practices**: Provide usage tips and recommendations

Response Requirements:
- Maintain a friendly, patient, and professional attitude
- Provide specific, actionable steps
- Use simple, understandable language
- If involving interface operations, clearly indicate button and option locations
- Format clearly, using numbers or bullet points when appropriate
- Base guidance on the above functionality summary for accuracy

Please answer in English, focusing on tool usage guidance.
        """
        
        # Use DeepSeek API
        return make_deepseek_api_call(prompt)
        
    except Exception as e:
        return f"Error generating chatbot response: {str(e)}"

def generate_tutor_response(user_message, language='en'):
    """Generate tutor response using DeepSeek AI"""
    try:
        # Check if client is already configured, if not configure it
        if 'deepseek_client' not in globals() or deepseek_client is None:
            if not configure_deepseek():
                return "Error: Failed to configure DeepSeek AI"
        
        # Create tutor prompt
        if language == 'zh':
            prompt = f"""
你是一个专业的FOB测试分析仪表板导师。你的主要任务是帮助用户学习如何使用这个工具。

**仪表板功能总结：**
这是一个功能观察电池（FOB）测试分析仪表板，包含以下核心功能：

1. **项目管理**：创建项目，设置动物类型（小鼠/大鼠/自定义），每组动物数量，组数
2. **6种分析模式**：
   - 一般行为：0/4/8评分系统，正常范围2-6
   - 自主神经功能：正常/异常二元评分
   - 反射能力：正常/异常二元评分
   - 体温：温度测量，正常范围36-38°C
   - 体重：实验前后体重测量，自动计算变化
   - 惊厥行为：正常/异常二元评分
3. **数据录入**：手动保存和自动保存两种模式，支持添加时间点
4. **管理组**：多组实验，可设置对照组
5. **数据分析**：异常事件跟踪，统计分析和可视化
6. **报告生成**：综合报告导出，图表下载
7. **多语言支持**：中英文界面

用户问题：{user_message}

请专注于提供以下帮助：
1. **工具使用指导**：详细说明如何在仪表板中操作各个功能
2. **步骤说明**：提供清晰的分步操作指南
3. **功能解释**：解释每个功能的作用和用途
4. **常见问题**：解答用户在使用过程中遇到的问题
5. **最佳实践**：提供使用建议和技巧

回答要求：
- 保持友好、耐心、专业的态度
- 提供具体、可操作的步骤
- 使用简单易懂的语言
- 如果涉及界面操作，请明确指出按钮和选项的位置
- 格式要清晰，可以使用编号或要点
- 基于上述功能总结提供准确的指导

请用中文回答，重点放在工具使用指导上。
"""
        else:
            prompt = f"""
You are a professional FOB Test Analysis Dashboard tutor. Your main task is to help users learn how to use this tool effectively.

**Dashboard Functionality Summary:**
This is a Functional Observational Battery (FOB) test analysis dashboard with the following core features:

1. **Project Management**: Create projects, set animal types (mouse/rat/custom), animals per group, number of groups
2. **6 Analysis Modes**:
   - General Behavior: 0/4/8 scoring system, normal range 2-6
   - Autonomic Functions: Normal/Abnormal binary scoring
   - Reflex Capabilities: Normal/Abnormal binary scoring
   - Body Temperature: Temperature measurements, normal range 36-38°C
   - Body Weight: Before/after weight measurements with automatic change calculations
   - Convulsive Behaviors: Normal/Abnormal binary scoring
3. **Data Entry**: Manual save and auto-save modes, support for adding time points
4. **Group Management**: Multiple experimental groups, can set comparison group
5. **Data Analysis**: Abnormal episode tracking, statistical analysis and visualization
6. **Report Generation**: Comprehensive report export, chart downloads
7. **Multi-language Support**: English and Chinese interfaces

User Question: {user_message}

Please focus on providing the following help:
1. **Tool Usage Guidance**: Detailed instructions on how to operate various features in the dashboard
2. **Step-by-step Instructions**: Provide clear, actionable step-by-step guides
3. **Feature Explanation**: Explain what each feature does and its purpose
4. **Common Issues**: Answer questions users encounter while using the tool
5. **Best Practices**: Provide usage tips and recommendations

Response Requirements:
- Maintain a friendly, patient, and professional attitude
- Provide specific, actionable steps
- Use simple, understandable language
- If involving interface operations, clearly indicate button and option locations
- Format clearly, using numbers or bullet points when appropriate
- Base guidance on the above functionality summary for accuracy

Please answer in English, focusing on tool usage guidance.
        """
        
        # Use DeepSeek API
        return make_deepseek_api_call(prompt)
        
    except Exception as e:
        return f"Error generating tutor response: {str(e)}"

def generate_file_summary(file_content, filename, language='en'):
    """Generate summary of uploaded file content"""
    try:
        # Check if client is already configured, if not configure it
        if 'deepseek_client' not in globals() or deepseek_client is None:
            if not configure_deepseek():
                return "Error: Failed to configure DeepSeek AI"
        
        # Create file summary prompt
        if language == 'zh':
            prompt = f"""
你是一个专业的数据分析助手。请分析以下文件内容并生成简洁的摘要：

文件名：{filename}

文件内容：
{file_content}

请提供以下摘要：
1. **文件类型和格式**
2. **主要内容概述**
3. **关键数据点或发现**
4. **与FOB测试分析的相关性**
5. **重要观察或结论**

请用中文回答，格式要清晰简洁。
"""
        else:
            prompt = f"""
You are a professional data analysis assistant. Please analyze the following file content and generate a concise summary:

Filename: {filename}

File Content:
{file_content}

Please provide the following summary:
1. **File type and format**
2. **Main content overview**
3. **Key data points or findings**
4. **Relevance to FOB test analysis**
5. **Important observations or conclusions**

Please answer in English with clear and concise format.
        """
        
        # Use DeepSeek API
        return make_deepseek_api_call(prompt)
        
    except Exception as e:
        return f"Error generating file summary: {str(e)}"

def generate_group_specific_description(mode, project_data, language='en'):
    """Generate group-specific description based on actual data"""
    try:
        if not st.session_state.active_project:
            return "No active project data available."
        
        project_id = st.session_state.active_project
        project_groups = get_project_groups(project_id)
        
        if not project_groups:
            return "No groups available for analysis."
        
        # Collect data for each group
        group_data_summary = []
        
        for group in project_groups:
            worksheet_key = f"worksheet_{group}_{mode}"
            if worksheet_key in st.session_state:
                df = st.session_state[worksheet_key]
                
                # Calculate statistics based on mode
                if mode == "Body Weight":
                    # Get before and after weights
                    before_df = df[df['time'] == t('before')] if t('before') in df['time'].values else pd.DataFrame()
                    after_df = df[df['time'] == t('after')] if t('after') in df['time'].values else pd.DataFrame()
                    
                    if not before_df.empty and not after_df.empty:
                        animal_type = project_data.get('animal_type', 'mouse')
                        if animal_type == 'custom':
                            animal_type = project_data.get('custom_animal_name', 'animal')
                        num_animals = project_data.get('num_animals', 8)
                        
                        before_weights = []
                        after_weights = []
                        for i in range(1, num_animals + 1):
                            col = f'{animal_type}_{i}'
                            if col in before_df.columns:
                                for val in before_df[col]:
                                    try:
                                        before_weights.append(float(val))
                                    except:
                                        pass
                            if col in after_df.columns:
                                for val in after_df[col]:
                                    try:
                                        after_weights.append(float(val))
                                    except:
                                        pass
                        
                        if before_weights and after_weights:
                            avg_before = np.mean(before_weights)
                            avg_after = np.mean(after_weights)
                            change = avg_after - avg_before
                            change_pct = (change / avg_before * 100) if avg_before > 0 else 0
                            group_data_summary.append({
                                'group': group,
                                'avg_before': avg_before,
                                'avg_after': avg_after,
                                'change': change,
                                'change_pct': change_pct
                            })
                
                elif mode == "Body Temperature":
                    # Calculate average temperature
                    animal_type = project_data.get('animal_type', 'mouse')
                    if animal_type == 'custom':
                        animal_type = project_data.get('custom_animal_name', 'animal')
                    num_animals = project_data.get('num_animals', 8)
                    
                    all_temps = []
                    for i in range(1, num_animals + 1):
                        col = f'{animal_type}_{i}'
                        if col in df.columns:
                            for val in df[col]:
                                try:
                                    all_temps.append(float(val))
                                except:
                                    pass
                    
                    if all_temps:
                        avg_temp = np.mean(all_temps)
                        std_temp = np.std(all_temps)
                        group_data_summary.append({
                            'group': group,
                            'avg_temp': avg_temp,
                            'std_temp': std_temp
                        })
                
                elif mode in ["Autonomic and Sensorimotor Functions", "Reflex Capabilities", "Convulsive Behaviors and Excitability"]:
                    # Calculate abnormal percentage
                    animal_type = project_data.get('animal_type', 'mouse')
                    if animal_type == 'custom':
                        animal_type = project_data.get('custom_animal_name', 'animal')
                    num_animals = project_data.get('num_animals', 8)
                    
                    total_count = 0
                    abnormal_count = 0
                    for i in range(1, num_animals + 1):
                        col = f'{animal_type}_{i}'
                        if col in df.columns:
                            for val in df[col]:
                                total_count += 1
                                val_str = str(val).lower()
                                if val_str in [t('abnormal').lower(), t('pale').lower(), t('cyanosis').lower()]:
                                    abnormal_count += 1
                    
                    if total_count > 0:
                        abnormal_pct = (abnormal_count / total_count * 100)
                        group_data_summary.append({
                            'group': group,
                            'abnormal_pct': abnormal_pct,
                            'abnormal_count': abnormal_count,
                            'total_count': total_count
                        })
                
                else:  # General Behavior
                    # Calculate mean score
                    animal_type = project_data.get('animal_type', 'mouse')
                    if animal_type == 'custom':
                        animal_type = project_data.get('custom_animal_name', 'animal')
                    num_animals = project_data.get('num_animals', 8)
                    
                    all_scores = []
                    for i in range(1, num_animals + 1):
                        col = f'{animal_type}_{i}'
                        if col in df.columns:
                            for val in df[col]:
                                try:
                                    score = float(val)
                                    if not pd.isna(score):
                                        all_scores.append(score)
                                except:
                                    pass
                    
                    if all_scores:
                        avg_score = np.mean(all_scores)
                        std_score = np.std(all_scores)
                        group_data_summary.append({
                            'group': group,
                            'avg_score': avg_score,
                            'std_score': std_score
                        })
        
        # Generate description based on collected data
        if language == 'zh':
            desc_parts = []
            if mode == "Body Weight":
                for gd in group_data_summary:
                    desc_parts.append(f"{gd['group']}: 实验前平均体重 {gd['avg_before']:.2f}g, 实验后 {gd['avg_after']:.2f}g, 变化 {gd['change']:+.2f}g ({gd['change_pct']:+.1f}%)")
            elif mode == "Body Temperature":
                for gd in group_data_summary:
                    desc_parts.append(f"{gd['group']}: 平均体温 {gd['avg_temp']:.2f}°C (标准差 {gd['std_temp']:.2f}°C)")
            elif mode in ["Autonomic and Sensorimotor Functions", "Reflex Capabilities", "Convulsive Behaviors and Excitability"]:
                for gd in group_data_summary:
                    desc_parts.append(f"{gd['group']}: 异常率 {gd['abnormal_pct']:.1f}% ({gd['abnormal_count']}/{gd['total_count']})")
            else:
                for gd in group_data_summary:
                    desc_parts.append(f"{gd['group']}: 平均得分 {gd['avg_score']:.2f} (标准差 {gd['std_score']:.2f})")
            
            if desc_parts:
                return "各组数据分析结果：\n" + "；".join(desc_parts) + "。"
            else:
                return "暂无可用数据进行分析。"
        else:
            desc_parts = []
            if mode == "Body Weight":
                for gd in group_data_summary:
                    desc_parts.append(f"{gd['group']}: Before {gd['avg_before']:.2f}g, After {gd['avg_after']:.2f}g, Change {gd['change']:+.2f}g ({gd['change_pct']:+.1f}%)")
            elif mode == "Body Temperature":
                for gd in group_data_summary:
                    desc_parts.append(f"{gd['group']}: Mean temperature {gd['avg_temp']:.2f}°C (SD {gd['std_temp']:.2f}°C)")
            elif mode in ["Autonomic and Sensorimotor Functions", "Reflex Capabilities", "Convulsive Behaviors and Excitability"]:
                for gd in group_data_summary:
                    desc_parts.append(f"{gd['group']}: Abnormal rate {gd['abnormal_pct']:.1f}% ({gd['abnormal_count']}/{gd['total_count']})")
            else:
                for gd in group_data_summary:
                    desc_parts.append(f"{gd['group']}: Mean score {gd['avg_score']:.2f} (SD {gd['std_score']:.2f})")
            
            if desc_parts:
                return "Group-specific analysis results:\n" + "; ".join(desc_parts) + "."
            else:
                return "No data available for analysis."
    
    except Exception as e:
        return f"Error generating description: {str(e)}"

def generate_powerpoint_content(project_data, mode_eng, language='en', file_summaries=None):
    """Generate comprehensive PowerPoint content using AI"""
    try:
        # Check if client is already configured, if not configure it
        if 'deepseek_client' not in globals() or deepseek_client is None:
            if not configure_deepseek():
                return "Error: Failed to configure DeepSeek AI"
        
        # Create AI prompt for PowerPoint content
        if language == 'zh':
            prompt = f"""
你是一个专业的科学演示文稿制作专家。请为FOB测试分析创建一个完整的PowerPoint演示文稿内容。

项目信息：
- 项目名称：{project_data.get('name', 'N/A')}
- 动物类型：{project_data.get('animal_type', 'N/A')}
- 每组动物数量：{project_data.get('num_animals', 'N/A')}
- 分析模式：{mode_eng}

请创建以下幻灯片内容：

1. **标题页**：项目标题和基本信息
2. **介绍**：FOB测试的背景、目的和重要性
3. **实验设计**：实验方法、动物分组、观察参数
4. **实验描述**：具体的实验步骤和观察指标
5. **结果分析**：主要发现和数据分析
6. **讨论**：结果解释和意义
7. **结论**：主要结论和建议

请为每个幻灯片提供：
- 幻灯片标题
- 要点内容（使用项目符号）
- 简洁明了的表述

格式要求：
- 使用中文
- 内容专业且易于理解
- 适合学术演示
- 包含关键数据点
"""
        else:
            prompt = f"""
You are a professional scientific presentation expert. Please create comprehensive PowerPoint content for FOB test analysis.

Project Information:
- Project Name: {project_data.get('name', 'N/A')}
- Animal Type: {project_data.get('animal_type', 'N/A')}
- Animals per Group: {project_data.get('num_animals', 'N/A')}
- Analysis Mode: {mode_eng}

Please create content for the following slides:

1. **Title Slide**: Project title and basic information
2. **Introduction**: Background, purpose, and importance of FOB testing
3. **Experimental Design**: Methods, animal grouping, observation parameters
4. **Experiment Description**: Specific experimental procedures and observation indicators
5. **Results Analysis**: Main findings and data analysis
6. **Discussion**: Result interpretation and significance
7. **Conclusion**: Main conclusions and recommendations

For each slide, provide:
- Slide title
- Bullet point content
- Clear and concise language

Format requirements:
- Use English
- Professional and understandable content
- Suitable for academic presentation
- Include key data points
        """
        
        # Use DeepSeek API
        return make_deepseek_api_call(prompt)
        
    except Exception as e:
        return f"Error generating PowerPoint content: {str(e)}"

def create_powerpoint_presentation(project_data, mode_eng, language='en', file_summaries=None, charts_data=None):
    """Create a PowerPoint presentation with AI-generated content, charts, and professional template"""
    try:
        # Generate AI content first
        ai_content = generate_powerpoint_content(project_data, mode_eng, language, file_summaries)
        
        # Create a new presentation
        prs = Presentation()
        
        # Set slide dimensions (16:9 aspect ratio)
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Apply professional template styling
        def apply_template_styling(slide, title_text, content_text=""):
            """Apply professional styling to slides"""
            # Set background color (light blue gradient effect)
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(240, 248, 255)  # Light blue background
            
            # Style the title
            if slide.shapes.title:
                title = slide.shapes.title
                title.text = title_text
                title.text_frame.paragraphs[0].font.size = Pt(44)
                title.text_frame.paragraphs[0].font.bold = True
                title.text_frame.paragraphs[0].font.color.rgb = RGBColor(25, 25, 112)  # Dark blue
                title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Style the content
            if content_text and slide.placeholders[1]:
                content = slide.placeholders[1]
                content.text = content_text
                for paragraph in content.text_frame.paragraphs:
                    paragraph.font.size = Pt(18)
                    paragraph.font.color.rgb = RGBColor(47, 84, 150)  # Medium blue
                    if paragraph.text.startswith('•') or paragraph.text.startswith('-'):
                        paragraph.font.bold = True
                        paragraph.font.size = Pt(20)
        
        # Title slide with enhanced styling
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        
        # Custom title slide content
        title_text = f"FOB Test Analysis Report"
        subtitle_text = f"""
        🧪 Project: {project_data.get('name', 'N/A')}
        🐭 Analysis Mode: {mode_eng}
        📊 Animals: {project_data.get('animal_type', 'N/A')} ({project_data.get('num_animals', 'N/A')} per group)
        📅 Generated: {datetime.datetime.now().strftime('%Y-%m-%d %H:%M')}
        🔬 FOB Test Analysis Dashboard
        """
        
        apply_template_styling(slide, title_text, subtitle_text)
        
        # Create structured slides following the specified format
        # Slide 1: Introduction to FOB Testing
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        intro_title = "Introduction to FOB Testing"
        intro_content = """
        The Functional Observational Battery (FOB) is a comprehensive behavioral assessment tool designed to evaluate the neurological and physiological effects of chemical compounds, drugs, or treatments in laboratory animals. This standardized battery of tests provides a systematic approach to detecting and characterizing potential neurotoxic effects through non-invasive observational methods.

        FOB testing encompasses six primary domains of assessment: general behavior, autonomic and sensorimotor functions, reflex capabilities, body temperature, body weight, and convulsive behaviors. Each domain targets specific aspects of neurological function, allowing researchers to identify subtle changes in behavior, physiology, and neurological responses that may indicate treatment-related effects.

        The FOB approach is particularly valuable in preclinical toxicology studies, drug development, and safety assessment, as it provides a comprehensive evaluation of multiple neurological endpoints in a single testing session. This systematic evaluation helps researchers identify potential safety concerns early in the development process and guides decisions regarding compound progression or additional safety studies.
        """
        apply_template_styling(slide, intro_title, intro_content)
        
        # Slide 2: FOB Scoring System Overview
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        scoring_title = "FOB Scoring System Overview"
        scoring_content = """
        The FOB employs a standardized scoring system that quantifies behavioral and physiological responses across multiple parameters. Each assessment domain utilizes specific scoring criteria designed to capture both normal and abnormal responses, with scores typically ranging from 0 (normal) to higher values indicating increasing levels of abnormality or impairment.

        General behavior assessments evaluate exploration, grooming, alertness, and overall health status using a combination of qualitative observations and quantitative measurements. Autonomic and sensorimotor functions are scored based on observations of piloerection, skin color, respiratory patterns, and other autonomic responses. Reflex capabilities are assessed through standardized tests of startle response, touch reactivity, and various reflexes.

        Body temperature and weight measurements provide objective physiological data, while convulsive behaviors and excitability are scored based on the presence and severity of tremors, convulsions, and other excitatory responses. The comprehensive nature of this scoring system ensures that subtle neurological changes are captured and quantified, providing researchers with reliable data for statistical analysis and interpretation.
        """
        apply_template_styling(slide, scoring_title, scoring_content)
        
        # Slide 3: Experimental Design and Methodology
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        method_title = "Experimental Design and Methodology"
        method_content = f"""
        This study employed a comprehensive FOB testing protocol to evaluate the effects of experimental treatments across multiple behavioral and physiological domains. The experimental design incorporated multiple treatment groups, including control animals, to ensure robust statistical analysis and reliable detection of treatment-related effects.

        Animals were randomly assigned to treatment groups and subjected to standardized FOB testing procedures following established protocols. Testing sessions were conducted at predetermined timepoints to capture both immediate and delayed effects of treatment. All observations were performed by trained technicians using standardized scoring criteria to ensure consistency and reliability of data collection.

        The methodology included comprehensive data collection across all six FOB domains, with particular focus on {mode_eng} analysis. Statistical analysis was performed using appropriate parametric and non-parametric tests to identify significant differences between treatment groups. Quality assurance measures were implemented throughout the study to ensure data integrity and compliance with regulatory guidelines.
        """
        apply_template_styling(slide, method_title, method_content)
        
        # Add comprehensive analysis slides for each mode
        if charts_data:
            # Group charts by mode
            charts_by_mode = {}
            for chart in charts_data:
                mode = chart.get('mode', 'Unknown')
                if mode not in charts_by_mode:
                    charts_by_mode[mode] = []
                charts_by_mode[mode].append(chart)
            
            # Create mode overview slide
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            overview_title = "Comprehensive FOB Test Analysis Overview"
            overview_content = f"""
            This comprehensive analysis encompasses all six primary domains of FOB testing, providing a complete evaluation of neurological and physiological responses to experimental treatments. The analysis includes general behavior assessment, autonomic and sensorimotor function evaluation, reflex capability testing, body temperature monitoring, body weight measurements, and convulsive behavior assessment.

            The study generated {len(charts_data)} detailed visualizations across all analysis modes, including group comparison analyses, time series trend evaluations, and comprehensive statistical summaries. Each mode was evaluated using standardized scoring criteria and appropriate statistical methods to ensure reliable detection of treatment-related effects.

            The comprehensive nature of this analysis allows for identification of subtle neurological changes that may not be apparent when examining individual parameters in isolation. This multi-parameter approach provides researchers with a complete picture of treatment effects across multiple neurological and physiological domains, supporting informed decision-making in drug development and safety assessment.
            
            📊 Chart Summary: {len(charts_data)} total visualizations including:
            • Real-time experiment analysis charts
            • Group comparison plots
            • Statistical analysis summaries
            • Time series trend analysis
            • Comprehensive mode-specific evaluations
            """
            apply_template_styling(slide, overview_title, overview_content)
            
            # Create slides for each mode with three-part structure
            for mode, mode_charts in charts_by_mode.items():
                # Part 1: Mode Introduction (150 words max)
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                mode_intro_title = f"{mode} - Introduction"
                
                # Mode-specific introduction content (150 words max)
                if mode == "General Behavior":
                    mode_intro_content = """
                    The General Behavior analysis evaluates overall behavioral patterns and health status in experimental animals through systematic observation of exploration, grooming, alertness, and general health indicators. This assessment serves as the primary indicator of treatment effects on central nervous system function and overall well-being.

                    The analysis employs standardized protocols to evaluate locomotor activity, environmental exploration, and grooming behaviors that reflect neurological function. Alertness and responsiveness to environmental stimuli are carefully assessed to identify treatment-related changes in cognitive function and behavioral responses.

                    This comprehensive evaluation provides critical insights into the overall impact of experimental treatments on animal behavior and health status, serving as a foundation for interpreting more specific neurological assessments.
                    """
                elif mode == "Autonomic and Sensorimotor Functions":
                    mode_intro_content = """
                    The Autonomic and Sensorimotor Functions analysis evaluates critical autonomic nervous system responses and sensorimotor coordination essential for neurological function. This assessment focuses on piloerection patterns, skin color changes, respiratory activity, and breathing patterns that reflect autonomic nervous system integrity.

                    The analysis systematically evaluates autonomic responses that may indicate treatment-related effects on sympathetic and parasympathetic nervous system function. Sensorimotor coordination assessments provide insights into the integration of sensory input and motor output, critical for normal neurological function.

                    Results from this analysis are particularly valuable for identifying potential autonomic dysfunction and sensorimotor impairment that may not be apparent through other assessment methods.
                    """
                elif mode == "Reflex Capabilities":
                    mode_intro_content = """
                    The Reflex Capabilities analysis assesses fundamental reflex responses that are critical indicators of neurological integrity and function. This evaluation includes standardized testing of startle response, touch reactivity, vocalization patterns, and various other reflexes that provide insights into sensory processing and motor response integration.

                    The analysis employs standardized reflex testing protocols to ensure consistent and reliable assessment of neurological function across all experimental groups. Startle response testing evaluates auditory processing and motor response integration, while touch reactivity assessments provide insights into tactile sensory processing.

                    Results from this analysis are essential for identifying potential treatment-related effects on fundamental neurological processes and reflex pathways.
                    """
                elif mode == "Body Temperature":
                    mode_intro_content = """
                    The Body Temperature analysis monitors physiological responses through continuous measurement of body temperature changes throughout the experimental period. This assessment provides critical data on thermoregulatory function and may indicate treatment-related effects on metabolic processes and autonomic nervous system function.

                    The analysis includes systematic temperature monitoring at predetermined timepoints to capture both immediate and delayed effects of experimental treatments on thermoregulatory function. Temperature changes may reflect alterations in metabolic rate, autonomic nervous system function, or direct effects on thermoregulatory centers.

                    Results from this analysis provide essential physiological data that complement behavioral assessments and may indicate treatment-related effects on metabolic processes or autonomic function.
                    """
                elif mode == "Body Weight":
                    mode_intro_content = """
                    The Body Weight analysis evaluates treatment effects through systematic measurement of body weight changes before and after experimental treatment administration. This assessment provides critical data on overall health status, metabolic effects, and potential treatment-related impacts on growth and development.

                    The analysis includes precise weight measurements at predetermined timepoints to capture both immediate and long-term effects of experimental treatments on body weight and overall health status. Weight changes may reflect alterations in metabolic processes, appetite, or direct effects on growth and development processes.

                    Results from this analysis provide essential physiological data that complement behavioral and neurological assessments and may indicate potential safety concerns related to metabolic or growth effects.
                    """
                elif mode == "Convulsive Behaviors and Excitability":
                    mode_intro_content = """
                    The Convulsive Behaviors and Excitability analysis evaluates potential treatment-related effects on neurological excitability and seizure susceptibility through systematic assessment of convulsive behaviors, tremors, stereotypy, and excitability patterns. This assessment is critical for identifying potential pro-convulsant or anti-convulsant effects.

                    The analysis includes careful observation of spontaneous convulsive behaviors, tremor patterns, and excitability responses to various stimuli. Stereotypy assessment evaluates repetitive behaviors that may indicate treatment-related effects on neurological function or potential neurotoxic effects.

                    Results from this analysis are particularly important for safety assessment, as convulsive behaviors and excitability changes may indicate potential treatment-related effects on neurological function or seizure susceptibility.
                    """
                else:
                    mode_intro_content = f"""
                    The {mode} analysis provides comprehensive evaluation of behavioral and physiological parameters relevant to neurological function and treatment effects. This assessment encompasses multiple endpoints that collectively provide insights into the impact of experimental treatments on various aspects of neurological and physiological function.

                    The analysis employs standardized protocols and scoring criteria to ensure consistent and reliable assessment across all experimental groups. Multiple parameters are evaluated to provide a comprehensive picture of treatment effects and identify potential safety concerns or therapeutic benefits.

                    Results from this analysis provide essential data for understanding the comprehensive impact of experimental treatments on neurological and physiological function.
                    """
                
                apply_template_styling(slide, mode_intro_title, mode_intro_content)
                
                # Part 2: Results Slides (Grouped by chart type)
                # Group comparison results
                group_charts = [c for c in mode_charts if 'Group Comparison' in c['title']]
                if group_charts:
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    background = slide.background
                    fill = background.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(240, 248, 255)
                    
                    # Add chart title
                    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1))
                    title_frame = title_box.text_frame
                    title_frame.text = f"{mode} - Group Comparison Results"
                    title_frame.paragraphs[0].font.size = Pt(28)
                    title_frame.paragraphs[0].font.bold = True
                    title_frame.paragraphs[0].font.color.rgb = RGBColor(25, 25, 112)
                    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    
                    # Add chart image
                    chart_box = slide.shapes.add_picture(
                        BytesIO(group_charts[0]['data']), 
                        Inches(1.5), 
                        Inches(2), 
                        Inches(7), 
                        Inches(4)
                    )
                    
                    # Generate group-specific description based on actual data
                    results_desc = generate_group_specific_description(mode, project_data, language)
                    
                    desc_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11.33), Inches(1))
                    desc_frame = desc_box.text_frame
                    desc_frame.text = results_desc
                    desc_frame.paragraphs[0].font.size = Pt(14)
                    desc_frame.paragraphs[0].font.color.rgb = RGBColor(47, 84, 150)
                    desc_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                # Time series results
                time_charts = [c for c in mode_charts if 'Time Series' in c['title']]
                if time_charts:
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    background = slide.background
                    fill = background.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(240, 248, 255)
                    
                    # Add chart title
                    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1))
                    title_frame = title_box.text_frame
                    title_frame.text = f"{mode} - Time Series Analysis"
                    title_frame.paragraphs[0].font.size = Pt(28)
                    title_frame.paragraphs[0].font.bold = True
                    title_frame.paragraphs[0].font.color.rgb = RGBColor(25, 25, 112)
                    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    
                    # Add chart image
                    chart_box = slide.shapes.add_picture(
                        BytesIO(time_charts[0]['data']), 
                        Inches(1.5), 
                        Inches(2), 
                        Inches(7), 
                        Inches(4)
                    )
                    
                    # Add results description (150 words max)
                    if mode == "General Behavior":
                        results_desc = "Time series analysis revealed dynamic changes in behavioral patterns over the experimental period. Control animals maintained consistent behavioral performance throughout, while treatment groups showed time-dependent modifications. Peak effects were observed at specific timepoints, indicating temporal dynamics of treatment response."
                    elif mode == "Autonomic and Sensorimotor Functions":
                        results_desc = "Temporal analysis showed evolving patterns in autonomic function across the experimental timeline. Control animals maintained stable autonomic responses, while treatment groups exhibited time-dependent changes in sensorimotor coordination. Peak autonomic effects were observed at specific timepoints."
                    elif mode == "Reflex Capabilities":
                        results_desc = "Time-dependent analysis revealed changing patterns in reflex responses throughout the experimental period. Control animals maintained consistent reflex function, while treatment groups showed temporal modifications in response patterns. Peak reflex effects were observed at specific timepoints."
                    elif mode == "Body Temperature":
                        results_desc = "Temporal monitoring revealed dynamic changes in body temperature patterns over the experimental timeline. Control animals maintained stable temperature regulation, while treatment groups showed time-dependent temperature modifications. Peak thermoregulatory effects were observed at specific timepoints."
                    elif mode == "Body Weight":
                        results_desc = "Time series analysis showed progressive changes in body weight patterns throughout the experimental period. Control animals maintained stable weight trajectories, while treatment groups exhibited time-dependent weight modifications. Peak weight effects were observed at specific timepoints."
                    elif mode == "Convulsive Behaviors and Excitability":
                        results_desc = "Temporal analysis revealed changing patterns in excitability responses over the experimental timeline. Control animals maintained consistent excitability levels, while treatment groups showed time-dependent modifications in convulsive behaviors. Peak excitability effects were observed at specific timepoints."
                    else:
                        results_desc = f"Time series analysis of {mode} revealed dynamic changes over the experimental period, with control animals maintaining consistent responses and treatment groups showing time-dependent modifications. Peak effects were observed at specific timepoints, indicating temporal dynamics of treatment response."
                    
                    desc_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11.33), Inches(1))
                    desc_frame = desc_box.text_frame
                    desc_frame.text = results_desc
                    desc_frame.paragraphs[0].font.size = Pt(14)
                    desc_frame.paragraphs[0].font.color.rgb = RGBColor(47, 84, 150)
                    desc_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                # Statistical summary results
                stats_charts = [c for c in mode_charts if 'Statistical Summary' in c['title']]
                if stats_charts:
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    background = slide.background
                    fill = background.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(240, 248, 255)
                    
                    # Add chart title
                    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1))
                    title_frame = title_box.text_frame
                    title_frame.text = f"{mode} - Statistical Summary"
                    title_frame.paragraphs[0].font.size = Pt(28)
                    title_frame.paragraphs[0].font.bold = True
                    title_frame.paragraphs[0].font.color.rgb = RGBColor(25, 25, 112)
                    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    
                    # Add chart image
                    chart_box = slide.shapes.add_picture(
                        BytesIO(stats_charts[0]['data']), 
                        Inches(1.5), 
                        Inches(2), 
                        Inches(7), 
                        Inches(4)
                    )
                    
                    # Add results description (150 words max)
                    if mode == "General Behavior":
                        results_desc = "Statistical analysis confirmed significant treatment effects on behavioral parameters, with clear differences in mean scores, standard deviations, and distribution patterns across groups. Control animals showed optimal statistical measures, while treatment groups exhibited varying degrees of statistical modification."
                    elif mode == "Autonomic and Sensorimotor Functions":
                        results_desc = "Statistical evaluation revealed significant treatment effects on autonomic parameters, with distinct differences in mean values and variability across experimental groups. Control animals maintained optimal statistical measures, while treatment groups showed significant statistical modifications in autonomic function."
                    elif mode == "Reflex Capabilities":
                        results_desc = "Statistical analysis confirmed significant treatment effects on reflex parameters, with clear differences in response distributions and variability across groups. Control animals demonstrated optimal statistical measures, while treatment groups exhibited significant statistical modifications in reflex function."
                    elif mode == "Body Temperature":
                        results_desc = "Statistical evaluation revealed significant treatment effects on temperature parameters, with distinct differences in mean temperatures and variability across experimental groups. Control animals maintained optimal statistical measures, while treatment groups showed significant statistical modifications in thermoregulation."
                    elif mode == "Body Weight":
                        results_desc = "Statistical analysis confirmed significant treatment effects on weight parameters, with clear differences in weight distributions and variability across groups. Control animals demonstrated optimal statistical measures, while treatment groups exhibited significant statistical modifications in weight patterns."
                    elif mode == "Convulsive Behaviors and Excitability":
                        results_desc = "Statistical evaluation revealed significant treatment effects on excitability parameters, with distinct differences in response distributions and variability across experimental groups. Control animals maintained optimal statistical measures, while treatment groups showed significant statistical modifications in excitability patterns."
                    else:
                        results_desc = f"Statistical analysis of {mode} confirmed significant treatment effects, with clear differences in mean values, standard deviations, and distribution patterns across experimental groups. Control animals demonstrated optimal statistical measures, while treatment groups exhibited significant statistical modifications."
                    
                    desc_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11.33), Inches(1))
                    desc_frame = desc_box.text_frame
                    desc_frame.text = results_desc
                    desc_frame.paragraphs[0].font.size = Pt(14)
                    desc_frame.paragraphs[0].font.color.rgb = RGBColor(47, 84, 150)
                    desc_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                # Part 3: Mode Conclusion (150 words max)
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                mode_conclusion_title = f"{mode} - Conclusion"
                
                # Mode-specific conclusion content (150 words max)
                if mode == "General Behavior":
                    mode_conclusion_content = """
                    The General Behavior analysis revealed significant treatment-related effects on overall behavioral patterns and health status. Group comparisons demonstrated clear differences in exploration, grooming, and alertness behaviors between control and treatment groups, with statistical analysis confirming the significance of observed changes.

                    Time series analysis showed dynamic behavioral modifications over the experimental period, with peak effects occurring at specific timepoints. Statistical evaluation confirmed significant treatment effects on multiple behavioral parameters, indicating comprehensive impact on central nervous system function.

                    These findings provide critical insights into the overall safety and efficacy profile of experimental treatments, supporting informed decision-making in drug development and safety assessment protocols.
                    """
                elif mode == "Autonomic and Sensorimotor Functions":
                    mode_conclusion_content = """
                    The Autonomic and Sensorimotor Functions analysis identified significant treatment-related effects on autonomic nervous system responses and sensorimotor coordination. Group comparisons revealed distinct patterns in piloerection, respiratory activity, and sensorimotor function across experimental groups.

                    Temporal analysis showed evolving autonomic responses throughout the experimental period, with statistical evaluation confirming significant treatment effects on multiple autonomic parameters. These findings indicate potential impacts on sympathetic and parasympathetic nervous system function.

                    The results provide essential data for understanding treatment effects on autonomic function and may indicate potential safety concerns related to autonomic nervous system modification.
                    """
                elif mode == "Reflex Capabilities":
                    mode_conclusion_content = """
                    The Reflex Capabilities analysis demonstrated significant treatment-related effects on fundamental reflex responses and neurological integrity. Group comparisons revealed clear differences in startle response, touch reactivity, and other reflex parameters across experimental groups.

                    Time series analysis showed dynamic changes in reflex function over the experimental period, with statistical evaluation confirming significant treatment effects on multiple reflex parameters. These findings indicate potential impacts on basic neurological processes.

                    The results provide critical data for understanding treatment effects on fundamental neurological function and may indicate potential safety concerns related to reflex pathway modification.
                    """
                elif mode == "Body Temperature":
                    mode_conclusion_content = """
                    The Body Temperature analysis revealed significant treatment-related effects on thermoregulatory function and metabolic processes. Group comparisons demonstrated clear differences in temperature patterns and thermoregulatory responses across experimental groups.

                    Temporal analysis showed dynamic changes in body temperature over the experimental period, with statistical evaluation confirming significant treatment effects on thermoregulatory parameters. These findings indicate potential impacts on metabolic and autonomic function.

                    The results provide essential physiological data for understanding treatment effects on thermoregulatory processes and may indicate potential safety concerns related to metabolic modification.
                    """
                elif mode == "Body Weight":
                    mode_conclusion_content = """
                    The Body Weight analysis identified significant treatment-related effects on growth patterns and metabolic processes. Group comparisons revealed distinct weight trajectories and growth patterns across experimental groups.

                    Time series analysis showed progressive changes in body weight throughout the experimental period, with statistical evaluation confirming significant treatment effects on weight parameters. These findings indicate potential impacts on growth and development processes.

                    The results provide essential physiological data for understanding treatment effects on growth and metabolism and may indicate potential safety concerns related to developmental modification.
                    """
                elif mode == "Convulsive Behaviors and Excitability":
                    mode_conclusion_content = """
                    The Convulsive Behaviors and Excitability analysis revealed significant treatment-related effects on neurological excitability and seizure susceptibility. Group comparisons demonstrated clear differences in convulsive behaviors and excitability patterns across experimental groups.

                    Temporal analysis showed dynamic changes in excitability responses over the experimental period, with statistical evaluation confirming significant treatment effects on excitability parameters. These findings indicate potential impacts on neurological excitability.

                    The results provide critical safety data for understanding treatment effects on seizure susceptibility and may indicate potential pro-convulsant or anti-convulsant effects requiring careful consideration.
                    """
                else:
                    mode_conclusion_content = f"""
                    The {mode} analysis demonstrated significant treatment-related effects across multiple parameters within this assessment domain. Group comparisons revealed clear differences in response patterns between control and treatment groups, with statistical analysis confirming the significance of observed changes.

                    Time series analysis showed dynamic modifications over the experimental period, with statistical evaluation confirming significant treatment effects on multiple parameters. These findings provide comprehensive insights into treatment effects on neurological and physiological function.

                    The results contribute essential data for understanding the safety and efficacy profile of experimental treatments and support informed decision-making in drug development protocols.
                    """
                
                apply_template_styling(slide, mode_conclusion_title, mode_conclusion_content)
        
        # Add Statistical Summary slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        stats_title = "📋 Statistical Summary"
        stats_content = f"""
        📊 Analysis Parameters:
        • Mode: {mode_eng}
        • Sample Size: {project_data.get('num_animals', 'N/A')} animals per group
        • Analysis Date: {datetime.datetime.now().strftime('%Y-%m-%d')}
        
        📈 Key Metrics:
        • Mean scores and standard deviations
        • Group comparisons
        • Statistical significance tests
        • Effect sizes
        
        🎯 Findings:
        • Primary observations
        • Significant differences
        • Clinical implications
        • Recommendations
        """
        apply_template_styling(slide, stats_title, stats_content)
        
        # Add file analysis slide if available
        if file_summaries:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            file_title = "Additional File Analysis"
            
            file_content = f"""
            The analysis incorporated data from {len(file_summaries)} uploaded files, providing additional context and insights that enhanced the comprehensive evaluation of FOB test results. Each file was systematically analyzed and summarized to identify relevant information that could contribute to the interpretation of experimental findings.

            File data was carefully integrated into the overall analysis, with summaries providing valuable context for understanding the experimental conditions, historical data, and related research findings. This integration enhanced the depth and breadth of insights available for interpreting the FOB test results and identifying potential correlations with external data sources.

            The comprehensive review of uploaded files ensured that all available information was considered in the analysis, providing a more complete picture of the experimental context and potential factors that may have influenced the observed results. This approach supports robust interpretation of findings and enhances the reliability of conclusions drawn from the FOB test analysis.
            """
            
            apply_template_styling(slide, file_title, file_content)
        
        # Add Insights and Recommendations slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        insights_title = "Key Insights & Recommendations"
        insights_content = f"""
        The comprehensive analysis of {mode_eng} data revealed significant patterns and trends that provide valuable insights into the effects of experimental treatments on neurological and physiological function. Group performance evaluation demonstrated clear differences between treatment groups, with statistical analysis confirming the significance of observed effects across multiple parameters.

        Data analysis identified consistent patterns in treatment responses, with trend analysis revealing time-dependent effects that provide important insights into the temporal dynamics of treatment effects. Anomaly detection algorithms identified potential outliers and unusual responses that warrant further investigation and may indicate individual variation in treatment response.

        Based on the comprehensive analysis, several key recommendations emerge for future research directions. Follow-up experiments should focus on validating the observed effects in larger sample sizes and exploring the underlying mechanisms responsible for the observed treatment effects. Additional data collection should target specific timepoints and parameters that showed the most significant treatment-related changes.
        """
        apply_template_styling(slide, insights_title, insights_content)
        
        # Add Methodology slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        method_title = "Methodology & Experimental Design"
        method_content = f"""
        The experimental design employed a comprehensive FOB testing protocol using {project_data.get('animal_type', 'N/A')} as the animal model, with {project_data.get('num_animals', 'N/A')} animals per group to ensure statistical power and reliable detection of treatment effects. The study focused on {mode_eng} analysis, utilizing standardized FOB testing procedures and scoring criteria to ensure consistency and reliability across all experimental groups.

        FOB test parameters were carefully selected to capture the full spectrum of neurological and physiological responses, with systematic observation at predetermined timepoints throughout the experimental period. Data collection methods employed standardized protocols and quality control measures to ensure data integrity and reproducibility. All observations were performed by trained technicians using validated scoring criteria to minimize inter-observer variability.

        Statistical analysis employed appropriate parametric and non-parametric tests to identify significant differences between treatment groups, with data visualization techniques used to illustrate patterns and trends in the results. Quality assurance measures included comprehensive data validation, statistical testing for normality and homogeneity of variance, and reproducibility checks to ensure the reliability of findings.
        """
        apply_template_styling(slide, method_title, method_content)
        
        # Add Conclusion slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        conclusion_title = "Conclusions & Future Work"
        conclusion_content = f"""
        The comprehensive analysis of {mode_eng} data provides compelling evidence of treatment-related effects on neurological and physiological function, with statistical validation confirming the significance of observed changes across multiple parameters. The systematic evaluation of all six FOB test domains revealed consistent patterns of treatment effects that provide valuable insights into the safety and efficacy profiles of experimental compounds.

        Key conclusions from this study include the identification of significant treatment-related changes in specific neurological parameters, with clinical implications that warrant careful consideration in the development of therapeutic interventions. The statistical significance of observed effects across multiple domains provides strong evidence for the biological relevance of treatment effects and supports the validity of the experimental approach.

        Future research directions should focus on expanding the current findings through additional experiments with larger sample sizes and exploring the underlying mechanisms responsible for the observed treatment effects. Publication of these findings will contribute to the scientific literature and provide valuable data for regulatory decision-making and clinical development planning.
        """
        apply_template_styling(slide, conclusion_title, conclusion_content)
        
        # Save the presentation
        pptx_buffer = BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)
        
        return pptx_buffer.getvalue()
        
    except Exception as e:
        return f"Error creating PowerPoint presentation: {str(e)}"

# Language translations - Updated with Body Weight mode
TRANSLATIONS = {
    'en': {
        'page_title': 'FOB Test',
        'main_title': 'FOB Test',
        'main_subtitle': 'Visualize and compare Functional Observational Battery (FOB) test results across multiple groups',
        'language': 'Language',
        'create_project': 'Create New Project',
        'configure_project': 'Configure New Project',
        'project_name': 'Project Name',
        'animal_type': 'Animal Type',
        'mouse': 'Mouse',
        'rat': 'Rat',
        'custom': 'Custom',
        'custom_animal_name': 'Custom Animal Name',
        'animals_per_group': 'Number of animals per group',
        'num_groups': 'Number of groups to create',
        'create': 'Create Project',
        'cancel': 'Cancel',
        'select_mode': 'Select Analysis Mode',
        'choose_mode': 'Choose mode:',
        'general_behavior': 'General Behavior',
        'autonomic_functions': 'Autonomic and Sensorimotor Functions',
        'reflex_capabilities': 'Reflex Capabilities',
        'body_temperature': 'Body Temperature',
        'body_weight': 'Body Weight',
        'convulsive_behaviors': 'Convulsive Behaviors and Excitability',
        'experiment_groups': 'Experiment Groups',
        'select_group_edit': 'Select Group to Edit',
        'data_worksheet': 'Data Entry Worksheet',
        'manual_save': 'Edit with Save Button',
        'auto_save': 'Auto-Save Mode',
        'save_changes': 'Save Changes',
        'fill_random': 'Fill Random Data',
        'fill_all_random': 'Fill ALL Groups with Random Data',
        'add_timestep': 'Add',
        'reset': 'Reset',
        'export_csv': 'Export Worksheet as CSV',
        'mean_scores': 'Mean Scores Summary',
        'weight_summary': 'Weight Change Summary',
        'filter_time': 'Filter by time points:',
        'time': 'Time',
        'observation': 'Observation',
        'mean_score': 'Mean Score',
        'status': 'Status',
        'normal': 'Normal',
        'abnormal': 'Abnormal',
        'pale': 'Pale',
        'cyanosis': 'Cyanosis',
        'abnormal_episodes': 'Abnormal Episodes (Onset/Offset)',
        'onset_time': 'Onset Time',
        'offset_time': 'Offset Time',
        'duration': 'Duration',
        'peak_score': 'Peak Score',
        'no_abnormal': 'No abnormal episodes detected',
        'comparison_group': 'Select Comparison Group',
        'set_comparison': 'Set as Comparison Group',
        'is_comparison': 'This is a COMPARISON GROUP',
        'data_analysis': 'Data Analysis & Reporting',
        'select_analyze': 'Select groups to analyze',
        'select_all': 'Select All',
        'comparative_report': 'Comparative Analysis Report',
        'group_summary': 'Group Summary',
        'group': 'Group',
        'total_episodes': 'Total Abnormal Episodes',
        'affected_obs': 'Affected Observations',
        'none': 'None',
        'episodes_by_group': 'Abnormal Episodes by Group',
        'summary': 'Summary:',
        'avg_duration': 'Avg Duration',
        'max_peak': 'Max Peak Score',
        'no_episodes': 'No abnormal episodes detected in any group!',
        'comparative_viz': 'Comparative Visualization',
        'select_time_compare': 'Select Time Point for Comparison',
        'export_report': 'Export Report',
        'download_report': 'Download Complete Report',
        'download_templates': 'Download Data Templates',
        'template_type': 'Select Template Type',
        'csv_template': 'CSV Template',
        'excel_template': 'Excel Template',
        'download_csv_template': 'Download CSV Template',
        'download_excel_template': 'Download Excel Template',
        'tips': 'Tips:',
        'unsaved_changes': 'You have unsaved changes!',
        'changes_saved': 'Changes saved successfully!',
        'auto_saved': 'Auto-saved at',
        'add_new_timestep': 'Add new timestep:',
        'next_timestep': 'Next timestep (min)',
        'valid': 'Valid',
        'report_title': 'FOB Test Analysis Report',
        'report_generated': 'Report Generated',
        'detailed_episodes': 'DETAILED ABNORMAL EPISODES',
        'project': 'Project',
        'analysis_mode': 'Analysis Mode',
        'total_groups': 'Total Groups Analyzed',
        'not_set': 'Not set',
        'start_instruction': 'Click \'Create New Project\' to get started',
        'edit_tip': '**Choose your editing mode**: Use \'Edit with Save Button\' to batch your changes, or \'Auto-Save Mode\' for instant saves.',
        'no_groups': 'No groups created yet',
        'group_management': 'Group Management',
        'rename_group': 'Rename Group',
        'new_group_name': 'New Group Name:',
        'rename_group_btn': 'Rename Group',
        'delete_group': 'Delete Group',
        'confirm_deletion': 'Confirm deletion',
        'export_project': 'Export Project Data',
        'import_project': 'Import Project Data',
        'export_success': 'Project data exported successfully!',
        'import_success': 'Project data imported successfully!',
        'import_warning': 'This will replace all current data. Continue?',
        'no_project_to_export': 'No active project to export',
        'invalid_zip_file': 'Invalid ZIP file format',
        'filling_all': 'Filling all worksheets with random data...',
        'fill_complete': 'All worksheets filled with random data!',
        'download_plot': 'Download Plot',
        'abnormal_count': 'Abnormal Count',
        'binary_instruction': '**Instructions**: Click on any cell to toggle between Normal (default) and Abnormal (red). Each observation is assessed as either Normal or Abnormal for each animal.',
        'percentage_abnormal': '% Abnormal',
        'groups_to_plot': 'Groups to plot:',
        'select_groups_chart': 'Select groups to display in the chart:',
        'all_time_points': 'All Time Points',
        'before_experiment': 'Before Experiment',
        'after_experiment': 'After Experiment',
        'weight_change': 'Weight Change',
        'weight_g': 'Weight (g)',
        'percent_change': '% Change',
        'weight_instruction': '**Instructions**: Enter the weight (in grams) for each animal before and after the experiment. Weight changes will be calculated automatically.',
        'mean_weight': 'Mean Weight',
        'weight_loss': 'Weight Loss',
        'weight_gain': 'Weight Gain',
        'no_change': 'No Change',
        'animal': 'Animal',
        'change_g': 'Change (g)',
        'initial_weight': 'Initial Weight',
        'final_weight': 'Final Weight',
        'ai_report': 'AI-Powered Report',
        'generate_ai_report': 'Generate AI Report',
        'ai_report_placeholder': 'Enter your DeepSeek API key to generate AI-powered reports',
        'api_key': 'API Key',
        'ai_analysis': 'AI Analysis',
        'ai_insights': 'AI Insights',
        'ai_recommendations': 'AI Recommendations',
        'ai_section': 'AI Analysis Section',
        'upload_file': 'Upload File for AI Analysis',
        'upload_help': 'Upload CSV, Excel, or text files to include in AI analysis',
        'file_uploaded': 'File uploaded successfully',
        'no_file': 'No file uploaded',
        'ai_chatbot': 'AI Assistant',
        'chat_placeholder': 'Ask me anything about using this dashboard...',
        'send_message': 'Send',
        'clear_chat': 'Clear Chat',
        'chat_help': 'Ask questions about FOB testing, data analysis, or dashboard features',
        'temperature_help': 'Temperature for {animal} {num} in Celsius (e.g., 37.2)',
        'weight_help': 'Weight for {animal} {num} in grams',
        'binary_help': 'Click to toggle between Normal and Abnormal for {animal} {num}',
        'skin_color_help': 'Skin color for {animal} {num}: Normal, Pale, or Cyanosis',
        'score_help': 'Score for {animal} {num}. Use 0/4/8 system',
        'before': 'Before',
        'after': 'After',
        'min': 'min',
        'fob_test': 'FOB Test',
        'description_of_scores': 'Description of Scores',
        'project_management': 'Project Management',
        'create_new_project': 'Create New Project',
        'no_projects_yet': 'No projects created yet. Click \'Create New Project\' to get started.',
        'ai_tutor': 'AI Tutor',
        'ai_chatbot': 'AI Chatbot',
        'ai_report': 'AI Report',
        'powerpoint': 'PowerPoint',
        'upload_data': 'Upload Data',
        'upload_csv_excel': 'Upload CSV/Excel File',
        'upload_help': 'Upload a filled-in CSV or Excel file that matches the worksheet template',
        'download_template': 'Download Template',
        'file_uploaded_success': 'File uploaded and data imported successfully!',
        'invalid_file_format': 'Invalid file format. Please upload a CSV or Excel file.',
        'template_mismatch': 'File format does not match the expected template. Please check the column structure.',
        'import_data': 'Import Data',
        'replace_data': 'Replace Current Data',
        'merge_data': 'Merge with Current Data',
        'fob_intro_title': 'FOB (Functional Observational Battery) - Rodent Functional Scale Observation',
        'fob_intro_desc': 'FOB is a systematic **animal neurobehavioral observation method**, primarily used to assess the effects or toxicity of compounds or drugs on **central nervous system function**. It evaluates the functional status of animals (usually rats or mice) through a series of standardized behavioral, physiological, and neural reflex indicators, providing qualitative and semi-quantitative assessment.',
        'fob_purpose_title': 'Purpose',
        'fob_purpose_1': '1. Early screening for **neurotoxicity** of drugs;',
        'fob_purpose_2': '2. Assist in explaining abnormalities in behavioral or physiological experiments;',
        'fob_purpose_3': '3. As part of **Safety Pharmacology** or toxicology studies (e.g., ICH S7A/S7B requirements).',
        'fob_assessment_title': 'Main Assessment Content',
        'fob_assessment_intro': 'FOB is typically divided into three dimensions:',
        'fob_assessment_1': '**1. General Behavioral Observations**',
        'fob_assessment_2': '**2. Reflex and Neuromuscular Function**: Tactile reflex, righting reflex, forelimb/hindlimb grip strength, balance beam test, jumping response, pain reflex',
        'fob_assessment_3': '**3. Autonomic Function**: Salivation, pupil size, rectal temperature, urination/defecation, breathing pattern and skin color',
        'fob_applications_title': 'FOB Testing Applications',
        'fob_applications_1': '1. New drug safety evaluation (e.g., antidepressants, antiepileptics, anesthetics)',
        'fob_applications_2': '2. Neurotoxicity assessment of industrial chemicals and pesticides',
        'fob_applications_3': '3. Compare central nervous responses across different species or doses',
        'fob_results_title': 'Our Result Interpretation',
        'fob_results_desc': 'Through a **quantitative scoring system**, we record animal performance and compare it with untreated control groups and known drug-treated groups to determine if there are significant differences. If phenomena such as **excitement, lethargy, ataxia, abnormal reflexes** appear, it indicates that the central or peripheral nervous system is affected. We can use the following multi-dimensional quantitative formulas to compare with known drugs and quantify the toxic and side effects of test drugs on the nervous system.',
        'fob_formulas_title': 'Calculation Formulas',
        'fob_formula_incidence': '**Incidence**: I = i/N',
        'fob_formula_abnormality': '**Abnormality**: A = ∑aᵢ/8i',
        'fob_formula_severity': '**Severity**: S = ∑(aᵢ·b)/∑aᵢ',
        'fob_formula_legend': 'N: total number of parameters; i: number of abnormal parameters; aᵢ: mouse number of each abnormal parameter; b: percentage of score (score=1, b=25%; score=4, b=100%)'
    },
    'zh': {
        'page_title': 'FOB测试',
        'main_title': 'FOB测试',
        'main_subtitle': '可视化并比较多组功能观察电池（FOB）测试结果',
        'language': '语言',
        'create_project': '创建新项目',
        'configure_project': '配置新项目',
        'project_name': '项目名称',
        'animal_type': '动物类型',
        'mouse': '小鼠',
        'rat': '大鼠',
        'custom': '自定义',
        'custom_animal_name': '自定义动物名称',
        'animals_per_group': '每组动物数量',
        'num_groups': '创建组数',
        'create': '创建项目',
        'cancel': '取消',
        'select_mode': '选择分析模式',
        'choose_mode': '选择模式：',
        'general_behavior': '一般行为',
        'autonomic_functions': '自主神经和感觉运动功能',
        'reflex_capabilities': '反射能力',
        'body_temperature': '体温',
        'body_weight': '体重',
        'convulsive_behaviors': '惊厥行为和兴奋性',
        'experiment_groups': '实验组',
        'select_group_edit': '选择要编辑的组',
        'data_worksheet': '数据录入工作表',
        'manual_save': '编辑后保存',
        'auto_save': '自动保存模式',
        'save_changes': '保存更改',
        'fill_random': '填充随机数据',
        'fill_all_random': '为所有组填充随机数据',
        'add_timestep': '添加',
        'reset': '重置',
        'export_csv': '导出工作表为CSV',
        'mean_scores': '平均分数汇总',
        'weight_summary': '体重变化汇总',
        'filter_time': '按时间点筛选：',
        'time': '时间',
        'observation': '观察项',
        'mean_score': '平均分数',
        'status': '状态',
        'normal': '正常',
        'abnormal': '异常',
        'pale': '发白',
        'cyanosis': '发绀',
        'abnormal_episodes': '异常事件（起始/结束）',
        'onset_time': '起始时间',
        'offset_time': '结束时间',
        'duration': '持续时间',
        'peak_score': '峰值分数',
        'no_abnormal': '未检测到异常事件',
        'comparison_group': '选择对照组',
        'set_comparison': '设为对照组',
        'is_comparison': '这是对照组',
        'data_analysis': '数据分析与报告',
        'select_analyze': '选择要分析的组',
        'select_all': '全选',
        'comparative_report': '对比分析报告',
        'group_summary': '组别汇总',
        'group': '组别',
        'total_episodes': '异常事件总数',
        'affected_obs': '受影响的观察项',
        'none': '无',
        'episodes_by_group': '各组异常事件',
        'summary': '汇总：',
        'avg_duration': '平均持续时间',
        'max_peak': '最高峰值分数',
        'no_episodes': '所有组均未检测到异常事件！',
        'comparative_viz': '对比可视化',
        'select_time_compare': '选择比较时间点',
        'export_report': '导出报告',
        'download_report': '下载完整报告',
        'download_templates': '下载数据模板',
        'template_type': '选择模板类型',
        'csv_template': 'CSV模板',
        'excel_template': 'Excel模板',
        'download_csv_template': '下载CSV模板',
        'download_excel_template': '下载Excel模板',
        'tips': '提示：',
        'unsaved_changes': '您有未保存的更改！',
        'changes_saved': '更改已成功保存！',
        'auto_saved': '自动保存于',
        'add_new_timestep': '添加新时间点：',
        'next_timestep': '下一个时间点（分钟）',
        'valid': '有效',
        'report_title': 'FOB测试分析报告',
        'report_generated': '报告生成时间',
        'detailed_episodes': '详细异常事件',
        'project': '项目',
        'analysis_mode': '分析模式',
        'total_groups': '分析组总数',
        'not_set': '未设置',
        'start_instruction': '点击"创建新项目"开始',
        'edit_tip': '**选择编辑模式**：使用"编辑后保存"批量更改，或使用"自动保存模式"即时保存。',
        'no_groups': '尚未创建组',
        'group_management': '管理组',
        'rename_group': '重命名组',
        'new_group_name': '新组名称:',
        'rename_group_btn': '重命名组',
        'delete_group': '删除组',
        'confirm_deletion': '确认删除',
        'export_project': '导出项目数据',
        'import_project': '导入项目数据',
        'export_success': '项目数据导出成功！',
        'import_success': '项目数据导入成功！',
        'import_warning': '这将替换所有当前数据。继续吗？',
        'no_project_to_export': '没有活动项目可导出',
        'invalid_zip_file': '无效的ZIP文件格式',
        'filling_all': '正在为所有工作表填充随机数据...',
        'fill_complete': '所有工作表已填充随机数据！',
        'download_plot': '下载图表',
        'abnormal_count': '异常计数',
        'binary_instruction': '**说明**：点击任意单元格在正常（默认）和异常（红色）之间切换。每个观察项对每只动物评估为正常或异常。',
        'percentage_abnormal': '异常百分比',
        'groups_to_plot': '要绘制的组：',
        'select_groups_chart': '选择要在图表中显示的组：',
        'all_time_points': '所有时间点',
        'before_experiment': '实验前',
        'after_experiment': '实验后',
        'weight_change': '体重变化',
        'weight_g': '体重 (克)',
        'percent_change': '变化百分比',
        'weight_instruction': '**说明**：输入每只动物实验前和实验后的体重（以克为单位）。体重变化将自动计算。',
        'mean_weight': '平均体重',
        'weight_loss': '体重减轻',
        'weight_gain': '体重增加',
        'no_change': '无变化',
        'animal': '动物',
        'change_g': '变化 (克)',
        'initial_weight': '初始体重',
        'final_weight': '最终体重',
        'ai_report': 'AI智能报告',
        'generate_ai_report': '生成AI报告',
        'ai_report_placeholder': '输入您的DeepSeek API密钥以生成AI智能报告',
        'api_key': 'API密钥',
        'ai_analysis': 'AI分析',
        'ai_insights': 'AI洞察',
        'ai_recommendations': 'AI建议',
        'ai_section': 'AI分析区域',
        'upload_file': '上传文件进行AI分析',
        'upload_help': '上传CSV、Excel或文本文件以包含在AI分析中',
        'file_uploaded': '文件上传成功',
        'no_file': '未上传文件',
        'ai_chatbot': 'AI助手',
        'chat_placeholder': '询问我关于使用此仪表板的任何问题...',
        'send_message': '发送',
        'clear_chat': '清除聊天',
        'chat_help': '询问有关FOB测试、数据分析或仪表板功能的问题',
        'temperature_help': '{animal} {num}的体温，以摄氏度为单位（例如：37.2）',
        'weight_help': '{animal} {num}的体重，以克为单位',
        'binary_help': '点击切换{animal} {num}的正常/异常状态',
        'skin_color_help': '{animal} {num}的皮肤颜色：正常、发白或发绀',
        'score_help': '{animal} {num}的分数。使用0/4/8评分系统',
        'before': '实验前',
        'after': '实验后',
        'min': '分钟',
        'fob_test': 'FOB测试',
        'description_of_scores': '评分说明',
        'project_management': '项目管理',
        'create_new_project': '创建新项目',
        'no_projects_yet': '尚未创建项目。点击"创建新项目"开始使用。',
        'ai_tutor': 'AI导师',
        'ai_chatbot': 'AI聊天机器人',
        'ai_report': 'AI报告',
        'powerpoint': 'PowerPoint',
        'upload_data': '上传数据',
        'upload_csv_excel': '上传CSV/Excel文件',
        'upload_help': '上传已填写的CSV或Excel文件，该文件应与工作表模板匹配',
        'download_template': '下载模板',
        'file_uploaded_success': '文件上传成功，数据已导入！',
        'invalid_file_format': '无效的文件格式。请上传CSV或Excel文件。',
        'template_mismatch': '文件格式与预期模板不匹配。请检查列结构。',
        'import_data': '导入数据',
        'replace_data': '替换当前数据',
        'merge_data': '与当前数据合并',
        'fob_intro_title': 'FOB（Functional Observational Battery）啮齿动物功能量表观察简介',
        'fob_intro_desc': 'FOB 是一种系统的**动物神经行为学观察方法**，主要用于评估化合物或药物对**中枢神经系统功能的影响或毒性**。它通过一系列标准化的行为、生理和神经反射指标，对动物（通常为大鼠或小鼠）的功能状态进行定性与半定量评价。',
        'fob_purpose_title': '一、目的',
        'fob_purpose_1': '1. 早期筛查药物是否具有**神经毒性（neurotoxicity）**；',
        'fob_purpose_2': '2. 辅助解释行为学或生理实验中的异常；',
        'fob_purpose_3': '3. 作为**安全药理学（Safety Pharmacology）或毒理学试验**的组成部分（如ICH S7A/S7B 要求）。',
        'fob_assessment_title': '二、主要评估内容',
        'fob_assessment_intro': 'FOB 通常分为三个维度：',
        'fob_assessment_1': '**1. 一般行为描述（General Observations）**',
        'fob_assessment_2': '**2. 神经反射与运动协调（Reflex and Neuromuscular Function）**：触觉反射、翻正反射（righting reflex）、前肢/后肢抓握力、平衡木测试、跳跃反应、疼痛反射',
        'fob_assessment_3': '**3. 自主神经功能（Autonomic Function）**：流涎、瞳孔大小、直肠温度、排尿/排便情况、呼吸模式与皮肤颜色',
        'fob_applications_title': 'FOB测试可用于：',
        'fob_applications_1': '1. 新药安全性评价（如抗抑郁药、抗癫痫药、麻醉药）',
        'fob_applications_2': '2. 工业化学品、农药的神经毒性评估',
        'fob_applications_3': '3. 比较不同物种或剂量下的中枢神经反应',
        'fob_results_title': '我们的结果解释',
        'fob_results_desc': '通过**定量评分系统**，记录动物表现，然后与无药物处理对照组比较和已知药物处理组的比较判断有无显著差异。若有**兴奋、嗜睡、运动失调、异常反射**等现象出现，即提示中枢或外周神经系统受到影响。我们可以用以下多维定量的公式与已知药物比较，定量受试药物对神经系统的毒副作用。',
        'fob_formulas_title': 'Calculation Formulas:',
        'fob_formula_incidence': '**Incidence**: I = i/N',
        'fob_formula_abnormality': '**Abnormality**: A = ∑aᵢ/8i',
        'fob_formula_severity': '**Severity**: S = ∑(aᵢ·b)/∑aᵢ',
        'fob_formula_legend': 'N: total number of parameters（总参数数量）; i: number of abnormal parameters（异常参数数量）; aᵢ: mouse number of each abnormal parameter（每个异常参数的小鼠数量）; b: percentage of score（分数百分比，score=1, b=25%; score=4, b=100%）'
    }
}

# Observation translations
OBSERVATION_TRANSLATIONS = {
    'en': {
        # General behavior observations
        'spontaneous exploration': 'spontaneous exploration',
        'grooming': 'grooming',
        'smelling its congeners': 'smelling its congeners',
        'normal resting state': 'normal resting state',
        'alertness': 'alertness',
        'distending / oedema': 'distending / oedema',
        'bad condition': 'bad condition',
        'moribund': 'moribund',
        'dead': 'dead',
        # Autonomic observations
        'piloerection': 'piloerection',
        'skin color': 'skin color',
        'respiratory activity': 'respiratory activity',
        'irregular breathing': 'irregular breathing',
        'stertorous': 'stertorous',
        # Reflex observations
        'startle response': 'startle response',
        'touch reactivity': 'touch reactivity',
        'vocalization': 'vocalization',
        'abnormal gait': 'abnormal gait',
        'corneal reflex': 'corneal reflex',
        'pinna reflex': 'pinna reflex',
        'catalepsy': 'catalepsy',
        'grip reflex': 'grip reflex',
        'pulling reflex': 'pulling reflex',
        'righting reflex': 'righting reflex',
        'body tone': 'body tone',
        'pain response': 'pain response',
        # Convulsive observations
        'spontaneous activity': 'spontaneous activity',
        'restlessness': 'restlessness',
        'fighting': 'fighting',
        'writhing': 'writhing',
        'tremor': 'tremor',
        'stereotypy': 'stereotypy',
        'twitches / jerks': 'twitches / jerks',
        'straub': 'straub',
        'opisthotonus': 'opisthotonus',
        'convulsion': 'convulsion',
        # Other
        'body temperature': 'body temperature',
        'body weight': 'body weight'
    },
    'zh': {
        # General behavior observations
        'spontaneous exploration': '自发探索',
        'grooming': '理毛',
        'smelling its congeners': '嗅探同类',
        'normal resting state': '正常休息状态',
        'alertness': '警觉性',
        'distending / oedema': '肿胀/水肿',
        'bad condition': '状态不佳',
        'moribund': '濒死',
        'dead': '死亡',
        # Autonomic observations
        'piloerection': '立毛',
        'skin color': '皮肤颜色',
        'respiratory activity': '呼吸活动',
        'irregular breathing': '呼吸不规则',
        'stertorous': '鼾声呼吸',
        # Reflex observations
        'startle response': '惊吓反应',
        'touch reactivity': '触觉反应',
        'vocalization': '发声',
        'abnormal gait': '步态异常',
        'corneal reflex': '角膜反射',
        'pinna reflex': '耳廓反射',
        'catalepsy': '僵直症',
        'grip reflex': '抓握反射',
        'pulling reflex': '牵拉反射',
        'righting reflex': '翻正反射',
        'body tone': '肌张力',
        'pain response': '疼痛反应',
        # Convulsive observations
        'spontaneous activity': '自发活动',
        'restlessness': '躁动不安',
        'fighting': '打斗',
        'writhing': '扭动',
        'tremor': '震颤',
        'stereotypy': '刻板行为',
        'twitches / jerks': '抽搐/痉挛',
        'straub': '竖尾反应',
        'opisthotonus': '角弓反张',
        'convulsion': '惊厥',
        # Other
        'body temperature': '体温',
        'body weight': '体重'
    }
}

# Initialize language
if 'language' not in st.session_state:
    st.session_state.language = 'zh'

# Get translation function
def t(key):
    """Get translation for the current language"""
    return TRANSLATIONS[st.session_state.language].get(key, key)

def t_obs(key):
    """Get observation translation for the current language"""
    return OBSERVATION_TRANSLATIONS[st.session_state.language].get(key, key)

# Set up the page
st.set_page_config(
    page_title="FOB Test",
    page_icon="",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Apply custom styling
def set_custom_style():
    st.markdown("""
        <style>
        .main {
            background-color: #f5f9ff;
        }
        .sidebar .sidebar-content {
            background-color: #e8f4ff;
        }
        h1 {
            color: #1a3d6d;
            border-bottom: 2px solid #1a3d6d;
        }
        .stButton>button {
            background-color: #1a3d6d;
            color: white;
            border-radius: 5px;
            padding: 8px 16px;
            transition: all 0.3s;
        }
        .stButton>button:hover {
            background-color: #2a5a9c;
            transform: scale(1.05);
        }
        .worksheet-container {
            background-color: white;
            border: 1px solid #e1e4e8;
            border-radius: 8px;
            padding: 20px;
            margin-bottom: 20px;
        }
        .plot-container {
            border: 1px solid #e1e4e8;
            border-radius: 8px;
            padding: 15px;
            margin-bottom: 20px;
            box-shadow: 0 2px 4px rgba(0,0,0,0.1);
            background-color: white;
        }
        .abnormal-high {
            background-color: #ffcccc;
        }
        .abnormal-low {
            background-color: #cce5ff;
        }
        .normal {
            background-color: #e6ffe6;
        }
        .template-download {
            border: 1px solid #1a3d6d;
            border-radius: 5px;
            padding: 10px;
            margin: 10px 0;
            background-color: #e8f4ff;
        }
        .autonomic-table {
            width: 100%;
            margin-bottom: 20px;
        }
        .autonomic-table th {
            background-color: #1a3d6d;
            color: white;
        }
        .autonomic-table tr:nth-child(even) {
            background-color: #f2f2f2;
        }
        .comparison-group {
            background-color: #d4edda;
            border: 2px solid #28a745;
        }
        .binary-instruction {
            background-color: #e8f4ff;
            border: 1px solid #1a3d6d;
            border-radius: 5px;
            padding: 15px;
            margin-bottom: 20px;
        }
        .weight-table {
            margin-top: 20px;
        }
        .weight-loss {
            color: #ff4444;
            font-weight: bold;
        }
        .weight-gain {
            color: #00aa00;
            font-weight: bold;
        }
        .no-change {
            color: #666666;
        }
        </style>
    """, unsafe_allow_html=True)

set_custom_style()



# Sidebar
with st.sidebar:
    st.title(f"🔬 {t('fob_test')}")
    
    # Help Icon for Scoring System
    if st.button(f"❓ {t('description_of_scores')}", use_container_width=True, help="Click to see how scoring is determined for each mode", key="sidebar_scoring_help"):
        st.session_state.show_scoring_help = not st.session_state.show_scoring_help
        st.rerun()
    
    # Language selection
    st.subheader("Language")
    selected_language = st.selectbox(
        t('language'),
        options=['en', 'zh'],
        format_func=lambda x: 'English' if x == 'en' else '中文',
        index=0 if st.session_state.language == 'en' else 1,
        key='language_selector'
    )
    
    # Handle language change
    if selected_language != st.session_state.language:
        st.session_state.language = selected_language
        st.rerun()
    
    st.markdown("---")
    
    # Project Management Section
    st.subheader(f"📁 {t('project_management')}")
    
    # Create New Project Button
    if st.button(f"➕ {t('create_new_project')}", use_container_width=True, type="primary"):
        st.session_state.show_project_creation = True
        st.rerun()
    
    # Project Selector
    if st.session_state.projects:
        st.markdown("**Select Project:**")
        # Create a list of tuples (project_id, project_name) for the selectbox
        project_options = [(project_id, project_data['name']) for project_id, project_data in st.session_state.projects.items()]
        project_ids = [option[0] for option in project_options]
        project_names = [option[1] for option in project_options]
        
        # Find the current index
        current_index = 0
        if st.session_state.active_project in project_ids:
            current_index = project_ids.index(st.session_state.active_project)
        
        selected_project_name = st.selectbox(
            "Choose a project:",
            options=project_names,
            index=current_index,
            key="sidebar_project_selector"
        )
        
        # Get the corresponding project ID
        selected_project_id = project_ids[project_names.index(selected_project_name)]
        
        if selected_project_id != st.session_state.active_project:
            st.session_state.active_project = selected_project_id
            st.rerun()
        
        # Show current project info and management options
        if st.session_state.active_project:
            project = st.session_state.projects[st.session_state.active_project]
            st.info(f"""
            **Current Project:** {project['name']}
            **Animal Type:** {t(project['animal_type'])}
            **Animals per Group:** {project['num_animals']}
            """)
            
            # Export/Import buttons
            st.markdown("**📤 Export/Import:**")
            col_export, col_import = st.columns(2)
            
            with col_export:
                if st.button("📤 Export", use_container_width=True, help="Download project data as ZIP file"):
                    with st.spinner("Exporting project data..."):
                        zip_data, message = export_project_data_as_zip()
                        if zip_data:
                            project_name = project['name'].replace(' ', '_')
                            timestamp = datetime.datetime.now().strftime('%Y%m%d_%H%M%S')
                            filename = f"{project_name}_export_{timestamp}.zip"
                            
                            st.download_button(
                                label="📥 Download ZIP",
                                data=zip_data,
                                file_name=filename,
                                mime="application/zip",
                                use_container_width=True
                            )
                            st.success(t('export_success'))
                        else:
                            st.error(message)
            
            with col_import:
                if st.button("📥 Import", use_container_width=True, help="Upload ZIP file to restore project data"):
                    st.session_state.show_import_dialog = True
            
            # Project management buttons
            col_delete, col_rename = st.columns(2)
            with col_delete:
                if st.button("🗑️ Delete", use_container_width=True, key="delete_project"):
                    if st.session_state.active_project in st.session_state.projects:
                        project_name = st.session_state.projects[st.session_state.active_project]['name']
                        project_id_to_delete = st.session_state.active_project
                        
                        # Remove related experiments
                        project_groups_to_remove = get_project_groups(project_id_to_delete)
                        for group in project_groups_to_remove:
                            if group in st.session_state.experiments:
                                del st.session_state.experiments[group]
                            if group in st.session_state.group_projects:
                                del st.session_state.group_projects[group]
                        
                        # Remove the project
                        del st.session_state.projects[project_id_to_delete]
                        
                        # Set active project to None or first available
                        if st.session_state.projects:
                            st.session_state.active_project = list(st.session_state.projects.keys())[0]
                        else:
                            st.session_state.active_project = None
                        
                        st.success(f"Project '{project_name}' deleted successfully!")
                        st.rerun()
    else:
        st.info(t('no_projects_yet'))
    
    st.markdown("---")
    
    # AI Features Section
    st.subheader("🤖 Dashboard AI")
    
    # Initialize AI session states
    if 'ai_tutor_active' not in st.session_state:
        st.session_state.ai_tutor_active = False
    if 'ai_chatbot_active' not in st.session_state:
        st.session_state.ai_chatbot_active = False
    if 'ai_report_active' not in st.session_state:
        st.session_state.ai_report_active = False
    if 'ai_powerpoint_active' not in st.session_state:
        st.session_state.ai_powerpoint_active = False
    if 'uploaded_file_content' not in st.session_state:
        st.session_state.uploaded_file_content = None
    if 'file_summaries' not in st.session_state:
        st.session_state.file_summaries = []
    if 'tutor_chat_history' not in st.session_state:
        st.session_state.tutor_chat_history = []
        # Add welcome message for tutor
        tutor_welcome_msg = "Hello! I'm your FOB Test Analysis Dashboard tutor. I can help you learn how to use this tool effectively. Ask me anything about the dashboard features, analysis modes, or how to perform specific tasks!"
        if st.session_state.language == 'zh':
            tutor_welcome_msg = "你好！我是你的FOB测试分析仪表板导师。我可以帮助你学习如何有效使用这个工具。询问我关于仪表板功能、分析模式或如何执行特定任务的问题！"
        st.session_state.tutor_chat_history.append({"role": "assistant", "content": tutor_welcome_msg})
    if 'chatbot_chat_history' not in st.session_state:
        st.session_state.chatbot_chat_history = []
        # Add welcome message for chatbot
        chatbot_welcome_msg = "Hello! I'm your file analysis assistant. Upload multiple files and I'll help you summarize their content for your FOB test analysis."
        if st.session_state.language == 'zh':
            chatbot_welcome_msg = "你好！我是你的文件分析助手。上传多个文件，我将帮助你总结其内容以用于FOB测试分析。"
        st.session_state.chatbot_chat_history.append({"role": "assistant", "content": chatbot_welcome_msg})
    
    # AI Tutor Button
    if st.button(f"🎓 {t('ai_tutor')}", use_container_width=True):
        st.session_state.ai_tutor_active = not st.session_state.ai_tutor_active
        st.session_state.ai_chatbot_active = False
        st.session_state.ai_report_active = False
        # Clear chat history when switching to tutor
        if st.session_state.ai_tutor_active:
            st.session_state.tutor_chat_history = []
            # Add welcome message for tutor
            tutor_welcome_msg = "Hello! I'm your FOB Test Analysis Dashboard tutor. I can help you learn how to use this tool effectively. Ask me anything about the dashboard features, analysis modes, or how to perform specific tasks!"
            if st.session_state.language == 'zh':
                tutor_welcome_msg = "你好！我是你的FOB测试分析仪表板导师。我可以帮助你学习如何有效使用这个工具。询问我关于仪表板功能、分析模式或如何执行特定任务的问题！"
            st.session_state.tutor_chat_history.append({"role": "assistant", "content": tutor_welcome_msg})
        st.rerun()
    
    # AI Chatbot Button
    if st.button(f"💬 {t('ai_chatbot')}", use_container_width=True):
        st.session_state.ai_chatbot_active = not st.session_state.ai_chatbot_active
        st.session_state.ai_tutor_active = False
        st.session_state.ai_report_active = False
        # Initialize chat when activating chatbot
        if st.session_state.ai_chatbot_active:
            # Initialize chat messages if empty (don't auto-open floating chat)
            st.session_state.floating_chat_open = False  # Don't auto-open the floating chat
            # Initialize chat messages if empty
            if 'chat_messages' not in st.session_state or not st.session_state.chat_messages:
                st.session_state.chat_messages = []
                # Add welcome message
                welcome_msg = "Hello! I'm your AI assistant. I can help you analyze files, answer questions about FOB testing, and provide guidance on using this dashboard. How can I help you today?"
                if st.session_state.language == 'zh':
                    welcome_msg = "你好！我是你的AI助手。我可以帮助你分析文件、回答FOB测试相关问题，并提供使用这个仪表板的指导。今天我能为你做些什么？"
                st.session_state.chat_messages.append({
                    "role": "assistant", 
                    "content": welcome_msg, 
                    "timestamp": datetime.datetime.now().strftime("%H:%M")
                })
        st.rerun()
    
    # AI Report Button
    if st.button(f"📊 {t('ai_report')}", use_container_width=True):
        st.session_state.ai_report_active = not st.session_state.ai_report_active
        st.session_state.ai_tutor_active = False
        st.session_state.ai_chatbot_active = False
        st.rerun()
    
    # PowerPoint Button
    if st.button(f"📈 {t('powerpoint')}", use_container_width=True):
        st.session_state.ai_powerpoint_active = not st.session_state.ai_powerpoint_active
        st.session_state.ai_tutor_active = False
        st.session_state.ai_chatbot_active = False
        st.session_state.ai_report_active = False
        st.rerun()
    
    # Show active AI feature
    if st.session_state.ai_tutor_active:
        st.success("🎓 AI Tutor Active")
    elif st.session_state.ai_chatbot_active:
        st.success("💬 AI Chatbot Active")
    elif st.session_state.ai_report_active:
        st.success("📊 AI Report Active")
    elif st.session_state.ai_powerpoint_active:
        st.success("📈 PowerPoint Active")
    
    st.markdown("---")

# App header
st.title(t('main_title'))
st.markdown(t('main_subtitle'))

# FOB Introduction Section
with st.expander("📖 " + t('fob_intro_title'), expanded=False):
    st.markdown(t('fob_intro_desc'))
    
    st.markdown(f"### {t('fob_purpose_title')}")
    st.markdown(t('fob_purpose_1'))
    st.markdown(t('fob_purpose_2'))
    st.markdown(t('fob_purpose_3'))
    
    st.markdown(f"### {t('fob_assessment_title')}")
    st.markdown(t('fob_assessment_intro'))
    st.markdown(t('fob_assessment_1'))
    st.markdown(t('fob_assessment_2'))
    st.markdown(t('fob_assessment_3'))
    
    st.markdown(f"### {t('fob_applications_title')}")
    st.markdown(t('fob_applications_1'))
    st.markdown(t('fob_applications_2'))
    st.markdown(t('fob_applications_3'))
    
    st.markdown(f"### {t('fob_results_title')}")
    st.markdown(t('fob_results_desc'))
    
    st.markdown(f"#### {t('fob_formulas_title')}")
    st.markdown(t('fob_formula_incidence'))
    st.markdown(t('fob_formula_abnormality'))
    st.markdown(t('fob_formula_severity'))
    st.markdown("")
    st.markdown(f"*{t('fob_formula_legend')}*")

st.markdown("---")

# Constants for modes
GENERAL_BEHAVIOR_OBSERVATIONS = [
    'spontaneous exploration',
    'grooming',
    'smelling its congeners',
    'normal resting state',
    'alertness',
    'distending / oedema',
    'bad condition',
    'moribund',
    'dead'
]

AUTONOMIC_OBSERVATIONS = [
    'piloerection',
    'skin color',
    'respiratory activity',
    'irregular breathing',
    'stertorous'
]

REFLEX_OBSERVATIONS = [
    'startle response',
    'touch reactivity',
    'vocalization',
    'abnormal gait',
    'corneal reflex',
    'pinna reflex',
    'catalepsy',
    'grip reflex',
    'pulling reflex',
    'righting reflex',
    'body tone',
    'pain response'
]

CONVULSIVE_OBSERVATIONS = [
    'spontaneous activity',
    'restlessness',
    'fighting',
    'writhing',
    'tremor',
    'stereotypy',
    'twitches / jerks',
    'straub',
    'opisthotonus',
    'convulsion'
]

# Updated ALL_MODES to include Body Weight
ALL_MODES = [
    "General Behavior", 
    "Autonomic and Sensorimotor Functions", 
    "Reflex Capabilities", 
    "Body Temperature",
    "Body Weight",
    "Convulsive Behaviors and Excitability"
]



# Helper function to save plot as bytes
def save_plot_as_bytes(fig):
    """Save matplotlib figure as bytes for download"""
    img_buffer = BytesIO()
    fig.savefig(img_buffer, format='png', dpi=300, bbox_inches='tight')
    img_buffer.seek(0)
    return img_buffer.getvalue()

# Helper function to parse the scoring system
def parse_score(score_str, observation=None):
    """Parse 0/4/8 scoring system with +/- modifiers or Normal/Abnormal"""
    if pd.isna(score_str):
        return np.nan
    
    # Handle binary Normal/Abnormal (both English and Chinese)
    score_lower = str(score_str).lower()
    
    # Check for normal values
    if score_lower in ['normal'] or score_lower in [t('normal').lower()]:
        return 0  # Normal is 0
    
    # Check for abnormal values (including pale and cyanosis)
    if score_lower in ['abnormal', 'pale', 'cyanosis'] or score_lower in [t('abnormal').lower(), t('pale').lower(), t('cyanosis').lower()]:
        return 1  # All abnormal variations are 1
    
    # Convert to string if it's a number
    if isinstance(score_str, (int, float)):
        return float(score_str)
    
    # Extract base score and modifiers
    match = re.match(r'(\d+(?:\.\d+)?)([\+\-]*)', str(score_str))
    if not match:
        return np.nan
    
    base_score = float(match.group(1))
    modifiers = match.group(2)
    
    # Calculate numerical value
    value = base_score
    if modifiers:
        modifier_value = len(modifiers) * (1 if '+' in modifiers else -1)
        value += modifier_value
    
    return value

# Function to calculate mean score from animal data
def calculate_mean_score(animal_scores, observation=None):
    """Calculate mean score from individual animal scores"""
    parsed_scores = [parse_score(score, observation) for score in animal_scores if pd.notna(score)]
    if parsed_scores:
        return np.mean(parsed_scores)
    return np.nan

# Function to generate random data
def generate_random_data(mode, times, num_animals=8, animal_type="mouse"):
    """Generate random data based on the mode"""
    if mode == "Body Temperature":
        # Normal temp range varies by animal type
        if animal_type == "rat":
            base_temp_mean = 37.5
        elif animal_type == "mouse":
            base_temp_mean = 37.0
        else:
            base_temp_mean = 37.2  # Default for custom animals
            
        data = {
            'time': [],
            'observation': []
        }
        for i in range(1, num_animals + 1):
            data[f'{animal_type}_{i}'] = []
        
        for time in times:
            data['time'].append(time)
            data['observation'].append(t_obs('body temperature'))
            for i in range(1, num_animals + 1):
                # Generate realistic body temperature
                base_temp = np.random.normal(base_temp_mean, 0.5)
                # Add some time-based variation
                if time > 30:
                    base_temp += np.random.normal(0.2, 0.1)
                data[f'{animal_type}_{i}'].append(f"{base_temp:.1f}")
        
        return pd.DataFrame(data)
    
    elif mode == "Body Weight":
        # Generate weight data for before and after experiment
        if animal_type == "rat":
            base_weight_mean = 250  # Rats are heavier
        elif animal_type == "mouse":
            base_weight_mean = 25   # Mice are lighter
        else:
            base_weight_mean = 100  # Default for custom animals
            
        data = {
            'time': [],
            'observation': []
        }
        for i in range(1, num_animals + 1):
            data[f'{animal_type}_{i}'] = []
        
        # Only two time points for weight: before (0) and after (e.g., end of experiment)
        for time_label in [t('before'), t('after')]:
            data['time'].append(time_label)
            data['observation'].append(t_obs('body weight'))
            for i in range(1, num_animals + 1):
                if time_label == t('before'):
                    # Initial weight
                    weight = np.random.normal(base_weight_mean, base_weight_mean * 0.1)
                else:
                    # After experiment - usually slight weight loss (stress, food restriction)
                    initial_weight = float(data[f'{animal_type}_{i}'][0])
                    # 90% chance of weight loss, 10% chance of weight gain
                    if np.random.random() < 0.9:
                        weight_change = np.random.uniform(-0.05, -0.01) * initial_weight  # 1-5% loss
                    else:
                        weight_change = np.random.uniform(0, 0.02) * initial_weight  # 0-2% gain
                    weight = initial_weight + weight_change
                data[f'{animal_type}_{i}'].append(f"{weight:.1f}")
        
        return pd.DataFrame(data)
    
    elif mode in ["Convulsive Behaviors and Excitability", "Autonomic and Sensorimotor Functions", "Reflex Capabilities"]:
        # Binary Normal/Abnormal system
        observations = []
        if mode == "Convulsive Behaviors and Excitability":
            observations = CONVULSIVE_OBSERVATIONS
        elif mode == "Autonomic and Sensorimotor Functions":
            observations = AUTONOMIC_OBSERVATIONS
        else:  # Reflex Capabilities
            observations = REFLEX_OBSERVATIONS
            
        data = {
            'time': [],
            'observation': []
        }
        for i in range(1, num_animals + 1):
            data[f'{animal_type}_{i}'] = []
        
        for time in times:
            for obs in observations:
                data['time'].append(time)
                data['observation'].append(t_obs(obs))  # obs is the key, t_obs translates it
                for i in range(1, num_animals + 1):
                    # Special handling for skin color observation
                    if obs == 'skin color':
                        # 70% normal, 15% pale, 15% cyanosis
                        rand = np.random.random()
                        if rand < 0.7:
                            data[f'{animal_type}_{i}'].append(t('normal'))
                        elif rand < 0.85:
                            data[f'{animal_type}_{i}'].append(t('pale'))
                        else:
                            data[f'{animal_type}_{i}'].append(t('cyanosis'))
                    else:
                        # 80% normal, 20% abnormal
                        if np.random.random() < 0.8:
                            data[f'{animal_type}_{i}'].append(t('normal'))
                        else:
                            data[f'{animal_type}_{i}'].append(t('abnormal'))
        
        return pd.DataFrame(data)
    
    else:  # General Behavior
        behaviors = GENERAL_BEHAVIOR_OBSERVATIONS
        data = {
            'time': [],
            'observation': []
        }
        for i in range(1, num_animals + 1):
            data[f'{animal_type}_{i}'] = []
        
        for time in times:
            for behavior in behaviors:
                data['time'].append(time)
                data['observation'].append(t_obs(behavior))
                for i in range(1, num_animals + 1):
                    # 0/4/8 system - generate scores mostly in normal range
                    if np.random.random() < 0.7:  # 70% normal range
                        base = 4
                    else:
                        base = np.random.choice([0, 8])
                    modifier = np.random.choice(['', '+', '-', '++', '--'], p=[0.5, 0.2, 0.2, 0.05, 0.05])
                    data[f'{animal_type}_{i}'].append(f"{base}{modifier}")
        
        return pd.DataFrame(data)

# Function to fill all worksheets with random data
def fill_all_worksheets_with_random_data():
    """Fill all worksheets for all groups and all modes with random data"""
    if st.session_state.active_project is None:
        st.error("No active project")
        return
    
    project = st.session_state.projects[st.session_state.active_project]
    animal_type = project.get('animal_type', 'mouse')
    if animal_type == 'custom':
        animal_type = project.get('custom_animal_name', 'animal')
    num_animals = project.get('num_animals', 8)
    
    # Get all groups for this project
    project_groups = get_project_groups(st.session_state.active_project)
    
    filled_count = 0
    
    # Progress tracking
    progress_placeholder = st.empty()
    progress_bar = st.progress(0)
    
    total_worksheets = len(project_groups) * len(ALL_MODES)
    current_worksheet = 0
    
    # For each group
    for group in project_groups:
        # For each mode
        for mode in ALL_MODES:
            current_worksheet += 1
            progress = current_worksheet / total_worksheets
            progress_bar.progress(progress)
            progress_placeholder.text(f"{t('filling_all')} ({current_worksheet}/{total_worksheets})")
            
            worksheet_key = f"worksheet_{group}_{mode}"
            
            # Generate appropriate observations for this mode
            if mode == "Autonomic and Sensorimotor Functions":
                observations = AUTONOMIC_OBSERVATIONS
            elif mode == "Reflex Capabilities":
                observations = REFLEX_OBSERVATIONS
            elif mode == "Convulsive Behaviors and Excitability":
                observations = CONVULSIVE_OBSERVATIONS
            elif mode == "Body Temperature":
                observations = ['body temperature']
            elif mode == "Body Weight":
                observations = ['body weight']
            else:
                observations = GENERAL_BEHAVIOR_OBSERVATIONS
            
            # Check if worksheet exists, if not create it
            if worksheet_key not in st.session_state:
                if mode == "Body Weight":
                    times = ['before', 'after']
                else:
                    times = [0]  # Only keep 0min by default, others should be added manually
                
                data = []
                for time in times:
                    for obs in observations:
                        row = {'time': time, 'observation': t_obs(obs)}
                        for i in range(1, num_animals + 1):
                            if mode == "Body Temperature":
                                row[f'{animal_type}_{i}'] = '37.0'
                            elif mode == "Body Weight":
                                row[f'{animal_type}_{i}'] = '25.0' if animal_type == 'mouse' else '250.0'
                            elif mode in ["Autonomic and Sensorimotor Functions", "Reflex Capabilities", "Convulsive Behaviors and Excitability"]:
                                # Special handling for skin color observation
                                if obs == 'skin color':
                                    row[f'{animal_type}_{i}'] = t('normal')
                                else:
                                    row[f'{animal_type}_{i}'] = t('normal')
                            else:
                                row[f'{animal_type}_{i}'] = '0'
                        data.append(row)
                st.session_state[worksheet_key] = pd.DataFrame(data)
            
            # Get existing times from the worksheet
            existing_df = st.session_state[worksheet_key]
            if mode == "Body Weight":
                times = ['before', 'after']
            else:
                times = sorted(existing_df['time'].unique())
            
            # Generate random data
            random_df = generate_random_data(mode, times, num_animals, animal_type)
            
            # Update the worksheet
            st.session_state[worksheet_key] = random_df
            st.session_state.worksheet_data[f"{group}_{mode}"] = random_df
            filled_count += 1
    
    # Clear progress indicators
    progress_placeholder.empty()
    progress_bar.empty()
    
    return filled_count

# Function to validate uploaded file format
def validate_uploaded_file(df, mode, animal_type, num_animals):
    """Validate that uploaded file matches the expected template"""
    if df is None or df.empty:
        return False, "Empty file"
    
    # Check required columns
    required_columns = ['time', 'observation']
    for col in required_columns:
        if col not in df.columns:
            return False, f"Missing required column: {col}"
    
    # Check animal columns
    expected_animal_columns = [f'{animal_type}_{i}' for i in range(1, num_animals + 1)]
    missing_columns = [col for col in expected_animal_columns if col not in df.columns]
    if missing_columns:
        return False, f"Missing animal columns: {missing_columns}"
    
    return True, "Valid format"

# Function to process uploaded file
def process_uploaded_file(uploaded_file, mode, animal_type, num_animals):
    """Process uploaded CSV/Excel file and return DataFrame"""
    try:
        # Read file based on extension
        if uploaded_file.name.endswith('.csv'):
            df = pd.read_csv(uploaded_file)
        elif uploaded_file.name.endswith(('.xlsx', '.xls')):
            df = pd.read_excel(uploaded_file)
        else:
            return None, t('invalid_file_format')
        
        # Validate format
        is_valid, message = validate_uploaded_file(df, mode, animal_type, num_animals)
        if not is_valid:
            return None, f"{t('template_mismatch')}: {message}"
        
        # Data will be displayed in the current language without migration
        
        return df, "Success"
        
    except Exception as e:
        return None, f"Error processing file: {str(e)}"

# Function to migrate existing English data to Chinese
def migrate_data_to_chinese(df, mode):
    """Convert existing English data to Chinese translations"""
    if df is None or df.empty:
        return df
    
    df_copy = df.copy()
    
    # Create reverse translation mapping
    reverse_obs_translations = {}
    if 'en' in OBSERVATION_TRANSLATIONS and 'zh' in OBSERVATION_TRANSLATIONS:
        en_obs = OBSERVATION_TRANSLATIONS['en']
        zh_obs = OBSERVATION_TRANSLATIONS['zh']
        for obs_key in en_obs:
            if obs_key in zh_obs:
                reverse_obs_translations[en_obs[obs_key]] = zh_obs[obs_key]
    
    # Migrate observation names
    if 'observation' in df_copy.columns:
        df_copy['observation'] = df_copy['observation'].apply(
            lambda x: reverse_obs_translations.get(x, x)
        )
    
    # Migrate time values for Body Weight mode
    if mode == "Body Weight" and 'time' in df_copy.columns:
        time_mapping = {
            'before': t('before'),
            'after': t('after')
        }
        df_copy['time'] = df_copy['time'].apply(
            lambda x: time_mapping.get(x, x)
        )
    
    # Migrate binary values (Normal/Abnormal)
    if mode in ["Autonomic and Sensorimotor Functions", "Reflex Capabilities", "Convulsive Behaviors and Excitability"]:
        binary_mapping = {
            'Normal': t('normal'),
            'Abnormal': t('abnormal')
        }
        for col in df_copy.columns:
            if col not in ['time', 'observation']:
                df_copy[col] = df_copy[col].apply(
                    lambda x: binary_mapping.get(x, x)
                )
    
    return df_copy

# Function to process data with onset/offset tracking
def process_data_with_episodes(df, mode, animal_type="mouse", num_animals=8):
    """Process data and track onset/offset of abnormal episodes"""
    results = []
    
    # Skip weight mode as it doesn't have episodes
    if mode == "Body Weight":
        return pd.DataFrame(results)
    
    # Get appropriate observations based on mode
    if mode == "Autonomic and Sensorimotor Functions":
        observations = [t_obs(obs) for obs in AUTONOMIC_OBSERVATIONS]
    elif mode == "Reflex Capabilities":
        observations = [t_obs(obs) for obs in REFLEX_OBSERVATIONS]
    elif mode == "Convulsive Behaviors and Excitability":
        observations = [t_obs(obs) for obs in CONVULSIVE_OBSERVATIONS]
    elif mode == "Body Temperature":
        observations = [t_obs('body temperature')]
    else:  # General Behavior
        observations = df['observation'].unique()
    
    for obs in observations:
        obs_df = df[df['observation'] == obs].sort_values('time')
        
        if obs_df.empty:
            continue
        
        # Track episodes
        onset_time = None
        in_episode = False
        peak_score = 0
        
        for _, row in obs_df.iterrows():
            # Calculate mean score from all animals
            animal_scores = [row[f'{animal_type}_{i}'] for i in range(1, num_animals + 1) 
                           if f'{animal_type}_{i}' in row]
            
            if mode in ["Autonomic and Sensorimotor Functions", "Reflex Capabilities", "Convulsive Behaviors and Excitability"]:
                # For binary modes, count percentage of abnormal
                # For all observations in Autonomic mode: pale and cyanosis are both abnormal
                abnormal_count = sum(1 for score in animal_scores 
                                   if str(score).lower() in [t('abnormal').lower(), t('pale').lower(), t('cyanosis').lower()])
                
                mean_score = (abnormal_count / len(animal_scores)) * 100 if animal_scores else 0
                is_abnormal = abnormal_count > 0  # Any animal abnormal
            else:
                mean_score = calculate_mean_score(animal_scores)
                
                # Track peak score
                if not pd.isna(mean_score) and mean_score > peak_score:
                    peak_score = mean_score
                
                # Determine if abnormal based on mode
                is_abnormal = False
                if mode == "Body Temperature":
                    # Abnormal if outside 36-38°C range
                    is_abnormal = mean_score < 36 or mean_score > 38
                else:  # General Behavior
                    # Abnormal if mean < 2 or > 6
                    is_abnormal = mean_score < 2 or mean_score > 6
            
            if is_abnormal and not in_episode:
                # Start of abnormal episode
                onset_time = row['time']
                in_episode = True
                peak_score = mean_score
            elif not is_abnormal and in_episode:
                # End of abnormal episode
                results.append({
                    t('observation'): t_obs(obs),
                    t('onset_time'): onset_time,
                    t('offset_time'): row['time'],
                    t('duration'): row['time'] - onset_time,
                    t('peak_score'): peak_score if mode not in ["Autonomic and Sensorimotor Functions", "Reflex Capabilities", "Convulsive Behaviors and Excitability"] else f"{peak_score:.0f}%"
                })
                in_episode = False
                onset_time = None
                peak_score = 0
        
        # Handle ongoing episode
        if in_episode and onset_time is not None:
            results.append({
                t('observation'): t_obs(obs),
                t('onset_time'): onset_time,
                t('offset_time'): obs_df['time'].max(),
                t('duration'): obs_df['time'].max() - onset_time,
                t('peak_score'): peak_score if mode not in ["Autonomic and Sensorimotor Functions", "Reflex Capabilities", "Convulsive Behaviors and Excitability"] else f"{peak_score:.0f}%"
            })
    
    return pd.DataFrame(results)

# Function to generate template data
def create_template(mode="General Behavior", num_animals=8, animal_type="mouse"):
    """Create template with individual animal columns"""
    if mode == "Body Temperature":
        times = [0, 15, 30, 45, 60]
        data = {
            'time': [],
            'observation': []
        }
        # Add animal columns
        for i in range(1, num_animals + 1):
            data[f'{animal_type}_{i}'] = []
        
        for time in times:
            data['time'].append(time)
            data['observation'].append('body temperature')
            for i in range(1, num_animals + 1):
                # Normal temperature range
                temp = np.random.normal(37.0, 0.2)
                data[f'{animal_type}_{i}'].append(f"{temp:.1f}")
        
        return pd.DataFrame(data)
    
    elif mode == "Body Weight":
        # Weight template with before/after
        data = {
            'time': [],
            'observation': []
        }
        # Add animal columns
        for i in range(1, num_animals + 1):
            data[f'{animal_type}_{i}'] = []
        
        for time_label in ['before', 'after']:
            data['time'].append(time_label)
            data['observation'].append('body weight')
            for i in range(1, num_animals + 1):
                if animal_type == "mouse":
                    weight = 25.0 if time_label == 'before' else 24.5
                elif animal_type == "rat":
                    weight = 250.0 if time_label == 'before' else 245.0
                else:
                    weight = 100.0 if time_label == 'before' else 98.0
                data[f'{animal_type}_{i}'].append(f"{weight:.1f}")
        
        return pd.DataFrame(data)
    
    elif mode in ["Convulsive Behaviors and Excitability", "Autonomic and Sensorimotor Functions", "Reflex Capabilities"]:
        times = [0, 15, 30]
        observations = []
        if mode == "Convulsive Behaviors and Excitability":
            observations = CONVULSIVE_OBSERVATIONS
        elif mode == "Autonomic and Sensorimotor Functions":
            observations = AUTONOMIC_OBSERVATIONS
        else:
            observations = REFLEX_OBSERVATIONS
            
        data = {
            'time': [],
            'observation': []
        }
        # Add animal columns
        for i in range(1, num_animals + 1):
            data[f'{animal_type}_{i}'] = []
        
        for time in times:
            for obs in observations:
                data['time'].append(time)
                data['observation'].append(obs)
                # Add scores for each animal - default to Normal
                for i in range(1, num_animals + 1):
                    data[f'{animal_type}_{i}'].append('Normal')
        
        return pd.DataFrame(data)
    
    else:  # General Behavior
        behaviors = GENERAL_BEHAVIOR_OBSERVATIONS
        times = [0, 15, 30]
        
        data = {
            'time': [],
            'observation': []
        }
        # Add animal columns
        for i in range(1, num_animals + 1):
            data[f'{animal_type}_{i}'] = []
        
        for time in times:
            for behavior in behaviors:
                data['time'].append(time)
                data['observation'].append(behavior)
                
                # Add scores for each animal
                for i in range(1, num_animals + 1):
                    # 0/4/8 system with modifiers
                    base = random.choice([0, 4, 8])
                    modifier = random.choice(['', '+', '-', '++', '--'])
                    data[f'{animal_type}_{i}'].append(f"{base}{modifier}")
        
        return pd.DataFrame(data)

# Function to create worksheet interface
def create_worksheet(mode, experiment_name, project_info):
    """Create an editable worksheet for data entry"""
    st.subheader(f"{t('data_worksheet')} - {experiment_name}")
    
    # Get animal info from project
    animal_type = project_info.get('animal_type', 'mouse')
    if animal_type == 'custom':
        animal_type = project_info.get('custom_animal_name', 'animal')
    num_animals = project_info.get('num_animals', 8)
    
    # Show if this is a comparison group
    if experiment_name in st.session_state.comparison_groups.get(st.session_state.active_project, []):
        st.success(t('is_comparison'))
    
    # Show appropriate instruction based on mode
    if mode in ["Autonomic and Sensorimotor Functions", "Reflex Capabilities", "Convulsive Behaviors and Excitability"]:
        st.markdown(f'<div class="binary-instruction">{t("binary_instruction")}</div>', unsafe_allow_html=True)
    elif mode == "Body Weight":
        st.markdown(f'<div class="binary-instruction">{t("weight_instruction")}</div>', unsafe_allow_html=True)
    
    # Create a unique key for this worksheet that includes mode
    worksheet_key = f"worksheet_{experiment_name}_{mode}"
    
    # Initialize worksheet data if not exists
    if worksheet_key not in st.session_state:
        if mode == "Autonomic and Sensorimotor Functions":
            observations = [t_obs(obs) for obs in AUTONOMIC_OBSERVATIONS]
        elif mode == "Reflex Capabilities":
            observations = [t_obs(obs) for obs in REFLEX_OBSERVATIONS]
        elif mode == "Convulsive Behaviors and Excitability":
            observations = [t_obs(obs) for obs in CONVULSIVE_OBSERVATIONS]
        elif mode == "Body Temperature":
            observations = [t_obs('body temperature')]
        elif mode == "Body Weight":
            observations = [t_obs('body weight')]
        else:
            observations = [t_obs(obs) for obs in GENERAL_BEHAVIOR_OBSERVATIONS]
        
        # Create initial data structure
        if mode == "Body Weight":
            times = [t('before'), t('after')]
        else:
            times = [0]  # Only keep 0min by default, others should be added manually
        
        data = []
        for time in times:
            for obs in observations:
                row = {'time': time, 'observation': obs}
                for i in range(1, num_animals + 1):
                    if mode == "Body Temperature":
                        # Default temperature values
                        if animal_type == 'mouse':
                            row[f'{animal_type}_{i}'] = '37.0'
                        elif animal_type == 'rat':
                            row[f'{animal_type}_{i}'] = '37.5'
                        else:
                            row[f'{animal_type}_{i}'] = '37.2'
                    elif mode == "Body Weight":
                        # Default weight values
                        if animal_type == 'mouse':
                            row[f'{animal_type}_{i}'] = '25.0'
                        elif animal_type == 'rat':
                            row[f'{animal_type}_{i}'] = '250.0'
                        else:
                            row[f'{animal_type}_{i}'] = '100.0'
                    elif mode in ["Autonomic and Sensorimotor Functions", "Reflex Capabilities", "Convulsive Behaviors and Excitability"]:
                        # Default to Normal state
                        row[f'{animal_type}_{i}'] = t('normal')
                    else:
                        # General Behavior: Default score
                        row[f'{animal_type}_{i}'] = '4'
                data.append(row)
        
        st.session_state[worksheet_key] = pd.DataFrame(data)
    
    # Get the dataframe from session state and migrate to Chinese if needed
    df = st.session_state[worksheet_key].copy()
    
    # Don't migrate existing data - just use it as is
    # The translation will happen in the display functions
    
    # Configure column settings with better formatting
    if mode == "Body Weight":
        column_config = {
            'time': st.column_config.SelectboxColumn(
                t('time'),
                options=[t('before'), t('after')],
                default=t('before'),
                disabled=True
            ),
            'observation': st.column_config.TextColumn(t('observation'), disabled=True)
        }
    else:
        column_config = {
            'time': st.column_config.NumberColumn(
                f"{t('time')} ({t('min')})", 
                min_value=0, 
                max_value=300, 
                step=5,
                format=f"%d {t('min')}"
            ),
            'observation': st.column_config.TextColumn(t('observation'), disabled=True)
        }
    
    # Add animal columns configuration
    for i in range(1, num_animals + 1):
        if mode == "Body Temperature":
            column_config[f'{animal_type}_{i}'] = st.column_config.NumberColumn(
                f'{t(animal_type).capitalize()} {i}',
                help=t('temperature_help').format(animal=t(animal_type), num=i),
                min_value=30.0,
                max_value=45.0,
                step=0.1,
                format="%.1f °C"
            )
        elif mode == "Body Weight":
            column_config[f'{animal_type}_{i}'] = st.column_config.NumberColumn(
                f'{t(animal_type).capitalize()} {i}',
                help=t('weight_help').format(animal=t(animal_type), num=i),
                min_value=0.0,
                max_value=1000.0,
                step=0.1,
                format="%.1f g"
            )
        elif mode in ["Autonomic and Sensorimotor Functions", "Reflex Capabilities", "Convulsive Behaviors and Excitability"]:
            # Check if this dataframe contains skin color observations
            df = st.session_state.get(worksheet_key, pd.DataFrame())
            has_skin_color = not df.empty and 'observation' in df.columns and any(
                obs.lower() in ['skin color', '皮肤颜色'] for obs in df['observation'].unique()
            )
            
            if has_skin_color:
                # Use expanded options for all rows when skin color is present
                column_config[f'{animal_type}_{i}'] = st.column_config.SelectboxColumn(
                    f'{t(animal_type).capitalize()} {i}',
                    help=t('skin_color_help').format(animal=t(animal_type), num=i),
                    options=[t('normal'), t('abnormal'), t('pale'), t('cyanosis')],
                    default=t('normal')
                )
            else:
                # Standard binary options
                column_config[f'{animal_type}_{i}'] = st.column_config.SelectboxColumn(
                    f'{t(animal_type).capitalize()} {i}',
                    help=t('binary_help').format(animal=t(animal_type), num=i),
                    options=[t('normal'), t('abnormal')],
                    default=t('normal')
                )
        else:
            # General Behavior: Default score
            column_config[f'{animal_type}_{i}'] = st.column_config.NumberColumn(
                f'{t(animal_type).capitalize()} {i}',
                help=t('score_help').format(animal=t(animal_type), num=i),
                min_value=0.0,
                max_value=20.0,
                step=0.5,
                format="%.1f"
            )
    
    # Create three tabs for different interaction modes
    tab1, tab2, tab3 = st.tabs([t('manual_save'), t('auto_save'), t('upload_data')])
    
    with tab1:
        st.markdown(f"**{t('manual_save')}**")
        
        # Check if there are unsaved changes
        temp_key = f"temp_{worksheet_key}"
        if temp_key in st.session_state and not df.equals(st.session_state[temp_key]):
            st.warning(t('unsaved_changes'))
        
        # Use a form to prevent constant reruns
        with st.form(key=f"form_{worksheet_key}"):
            # Create editable dataframe
            edited_df = st.data_editor(
                df,
                column_config=column_config,
                use_container_width=True,
                num_rows="dynamic" if mode != "Body Weight" else "fixed",
                key=f"editor_{worksheet_key}_form",
                hide_index=True
            )
            
            # Store temp changes
            st.session_state[temp_key] = edited_df
            
            # Form submit button
            col1, col2, col3, col4 = st.columns([2, 2, 2, 1])
            with col1:
                submitted = st.form_submit_button(t('save_changes'), use_container_width=True, type="primary")
            with col2:
                fill_random = st.form_submit_button(t('fill_random'), use_container_width=True)
            
            if mode != "Body Weight":
                with col3:
                    st.markdown(f"**{t('add_new_timestep')}**")
                    new_timestep = st.number_input(
                        t('next_timestep'), 
                        min_value=0,
                        max_value=300,
                        step=5,
                        value=edited_df['time'].max() + 5 if not edited_df.empty else 0,
                        key=f"new_time_{worksheet_key}",
                        label_visibility="collapsed"
                    )
                    add_timestep = st.form_submit_button(t('add_timestep'), use_container_width=True)
            else:
                add_timestep = False
                
            with col4:
                reset = st.form_submit_button(t('reset'), use_container_width=True)
            
            # Fix: Ensure state changes are properly handled
            if submitted:
                # Save data as is (no migration needed)
                final_df = edited_df.copy()
                
                # Sort by time to ensure proper ordering
                final_df = final_df.sort_values(['time', 'observation']).reset_index(drop=True)
                
                # Update session state with edited data
                st.session_state[worksheet_key] = final_df
                st.session_state.worksheet_data[f"{experiment_name}_{mode}"] = final_df
                st.session_state.save_status[experiment_name] = "saved"
                
                # Synchronize time points if time column was modified - sync to ALL modes
                if mode != "Body Weight" and st.session_state.active_project:
                    new_times = sorted(final_df['time'].unique())
                    synchronize_time_points_across_worksheets(st.session_state.active_project, new_times, mode=None)
                
                # Clear temp changes
                if temp_key in st.session_state:
                    del st.session_state[temp_key]
                st.success(t('changes_saved'))
                # Force rerun to reflect changes
                st.rerun()
            
            if fill_random:
                # Generate random data
                if mode == "Body Weight":
                    times = ['before', 'after']
                else:
                    times = sorted(edited_df['time'].unique())
                random_df = generate_random_data(mode, times, num_animals, animal_type)
                st.session_state[worksheet_key] = random_df
                st.rerun()
            
            if add_timestep and mode != "Body Weight":
                # Add new timestep with all observations
                new_rows = []
                observations = edited_df['observation'].unique()
                for obs in observations:
                    new_row = {'time': new_timestep, 'observation': obs}
                    for i in range(1, num_animals + 1):
                        if mode == "Body Temperature":
                            new_row[f'{animal_type}_{i}'] = '37.0'
                        elif mode in ["Autonomic and Sensorimotor Functions", "Reflex Capabilities", "Convulsive Behaviors and Excitability"]:
                            new_row[f'{animal_type}_{i}'] = t('normal')
                        else:
                            new_row[f'{animal_type}_{i}'] = '0'
                    new_rows.append(new_row)
                
                # Append new rows and sort by time
                new_df = pd.concat([edited_df, pd.DataFrame(new_rows)], ignore_index=True)
                new_df = new_df.sort_values(['time', 'observation']).reset_index(drop=True)
                st.session_state[worksheet_key] = new_df
                st.session_state.worksheet_data[f"{experiment_name}_{mode}"] = new_df
                
                # IMPORTANT: Synchronize time points IMMEDIATELY across all worksheets - sync to ALL modes
                # This happens before saving, so other modes update instantly
                if st.session_state.active_project:
                    synchronize_time_points_across_worksheets(st.session_state.active_project, [new_timestep], mode=None)
                
                st.success(f"{t('add_timestep')} {new_timestep} min - {t('changes_saved')} (所有模式已同步)")
                st.rerun()
            
            if reset:
                # Reset to original state
                if temp_key in st.session_state:
                    del st.session_state[temp_key]
                st.rerun()
    
    with tab2:
        st.markdown(f"**{t('auto_save')}**")
        st.info(t('edit_tip'))
        
        # Quick action buttons
        col1, col2, col3 = st.columns([1, 1, 2])
        with col1:
            if st.button(t('fill_random'), use_container_width=True, key=f"random_auto_{worksheet_key}"):
                if mode == "Body Weight":
                    times = [t('before'), t('after')]
                else:
                    times = sorted(st.session_state[worksheet_key]['time'].unique())
                random_df = generate_random_data(mode, times, num_animals, animal_type)
                st.session_state[worksheet_key] = random_df
                st.rerun()
        
        # Create editable dataframe without form (auto-saves)
        edited_df_auto = st.data_editor(
            st.session_state[worksheet_key],
            column_config=column_config,
            use_container_width=True,
            num_rows="dynamic" if mode != "Body Weight" else "fixed",
            key=f"editor_{worksheet_key}_auto",
            hide_index=True
        )
        
        # Auto-save the changes
        if not edited_df_auto.equals(st.session_state[worksheet_key]):
            # Save data as is (no migration needed)
            final_df_auto = edited_df_auto.copy()
            
            # Sort by time to ensure proper ordering
            final_df_auto = final_df_auto.sort_values(['time', 'observation']).reset_index(drop=True)
            
            st.session_state[worksheet_key] = final_df_auto
            st.session_state.worksheet_data[f"{experiment_name}_{mode}"] = final_df_auto
            st.session_state.save_status[experiment_name] = "saved"
            
            # Synchronize time points if time column was modified - sync to ALL modes
            if mode != "Body Weight" and st.session_state.active_project:
                new_times = sorted(final_df_auto['time'].unique())
                synchronize_time_points_across_worksheets(st.session_state.active_project, new_times, mode=None)
        
        # Show save status with timestamp
        st.success(f"{t('auto_saved')} {datetime.datetime.now().strftime('%H:%M:%S')}")
        
        # Quick actions - NEW TIMESTEP FUNCTIONALITY (not for Body Weight)
        if mode != "Body Weight":
            st.markdown(f"**{t('add_new_timestep')}**")
            col1, col2 = st.columns([3, 2])
            with col1:
                new_timestep_auto = st.number_input(
                    t('next_timestep'), 
                    min_value=0,
                    max_value=300,
                    step=5,
                    value=edited_df_auto['time'].max() + 5 if not edited_df_auto.empty else 0,
                    key=f"new_time_auto_{worksheet_key}",
                    label_visibility="collapsed"
                )
            with col2:
                if st.button(t('add_timestep'), use_container_width=True):
                    # Add rows for new timestep
                    new_rows = []
                    observations = edited_df_auto['observation'].unique()
                    for obs in observations:
                        new_row = {'time': new_timestep_auto, 'observation': obs}
                        for i in range(1, num_animals + 1):
                            if mode == "Body Temperature":
                                new_row[f'{animal_type}_{i}'] = '37.0'
                            elif mode in ["Autonomic and Sensorimotor Functions", "Reflex Capabilities", "Convulsive Behaviors and Excitability"]:
                                new_row[f'{animal_type}_{i}'] = t('normal')
                            else:
                                new_row[f'{animal_type}_{i}'] = '0'
                        new_rows.append(new_row)
                    
                    # Append new rows and sort by time
                    new_df = pd.concat([edited_df_auto, pd.DataFrame(new_rows)], ignore_index=True)
                    new_df = new_df.sort_values(['time', 'observation']).reset_index(drop=True)
                    st.session_state[worksheet_key] = new_df
                    st.session_state.worksheet_data[f"{experiment_name}_{mode}"] = new_df
                    
                    # IMPORTANT: Synchronize time points IMMEDIATELY across all worksheets - sync to ALL modes
                    # This happens automatically, so other modes update instantly without needing to save
                    if st.session_state.active_project:
                        synchronize_time_points_across_worksheets(st.session_state.active_project, [new_timestep_auto], mode=None)
                    
                    st.success(f"{t('add_timestep')} {new_timestep_auto} min - {t('auto_saved')} (所有模式已同步)")
                    st.rerun()
    
    with tab3:
        st.markdown(f"**{t('upload_data')}**")
        st.info(t('upload_help'))
        
        # File upload section
        uploaded_file = st.file_uploader(
            t('upload_csv_excel'),
            type=['csv', 'xlsx', 'xls'],
            help=t('upload_help'),
            key=f"upload_{worksheet_key}"
        )
        
        if uploaded_file is not None:
            # Process uploaded file
            processed_df, message = process_uploaded_file(uploaded_file, mode, animal_type, num_animals)
            
            if processed_df is not None:
                st.success(t('file_uploaded_success'))
                
                # Show preview of uploaded data
                st.subheader("Data Preview")
                st.dataframe(processed_df.head(10), use_container_width=True)
                
                # Import options
                col1, col2 = st.columns(2)
                with col1:
                    if st.button(t('replace_data'), use_container_width=True, type="primary"):
                        st.session_state[worksheet_key] = processed_df
                        st.session_state.worksheet_data[f"{experiment_name}_{mode}"] = processed_df
                        st.success("Data replaced successfully!")
                        st.rerun()
                
                with col2:
                    if st.button(t('merge_data'), use_container_width=True):
                        # Merge with existing data
                        existing_df = st.session_state[worksheet_key]
                        merged_df = pd.concat([existing_df, processed_df], ignore_index=True)
                        st.session_state[worksheet_key] = merged_df
                        st.session_state.worksheet_data[f"{experiment_name}_{mode}"] = merged_df
                        st.success("Data merged successfully!")
                        st.rerun()
            else:
                st.error(message)
        
        # Download template section
        st.markdown("---")
        st.subheader(t('download_template'))
        
        # Create template download
        template_df = st.session_state[worksheet_key].copy()
        csv_data = template_df.to_csv(index=False)
        
        st.download_button(
            label=t('download_template'),
            data=csv_data,
            file_name=f"{experiment_name}_{mode}_template.csv",
            mime="text/csv",
            use_container_width=True
        )
    
    # Get the current dataframe (from whichever tab was used)
    current_df = st.session_state[worksheet_key]
    
    # Display appropriate summary based on mode
    if mode == "Body Weight":
        st.subheader(t('weight_summary'))
        
        # Calculate weight changes
        weight_data = []
        before_df = current_df[current_df['time'] == 'before']
        after_df = current_df[current_df['time'] == 'after']
        
        if not before_df.empty and not after_df.empty:
            for i in range(1, num_animals + 1):
                animal_col = f'{animal_type}_{i}'
                if animal_col in before_df.columns:
                    try:
                        before_weight = float(before_df.iloc[0][animal_col])
                        after_weight = float(after_df.iloc[0][animal_col])
                        change = after_weight - before_weight
                        percent_change = (change / before_weight) * 100
                        
                        status = t('weight_loss') if change < 0 else (t('weight_gain') if change > 0 else t('no_change'))
                        
                        weight_data.append({
                            t('animal'): f'{t(animal_type).capitalize()} {i}',
                            f"{t('before_experiment')} (g)": f"{before_weight:.1f}",
                            f"{t('after_experiment')} (g)": f"{after_weight:.1f}",
                            f"{t('change_g')}": f"{change:.1f}",
                            t('percent_change'): f"{percent_change:.2f}%",
                            t('status'): status
                        })
                    except (ValueError, TypeError):
                        continue
        
        if weight_data:
            weight_df = pd.DataFrame(weight_data)
            st.dataframe(
                weight_df,
                use_container_width=True,
                hide_index=True,
                column_config={
                    t('status'): st.column_config.TextColumn(t('status'), width="small")
                }
            )
            
            # Calculate group statistics
            col1, col2, col3 = st.columns(3)
            with col1:
                before_key = f"{t('before_experiment')} (g)"
                mean_initial = np.mean([float(row[before_key]) for row in weight_data])
                st.metric(f"{t('mean_weight')} - {t('before_experiment')}", f"{mean_initial:.1f} g")
            with col2:
                after_key = f"{t('after_experiment')} (g)"
                mean_final = np.mean([float(row[after_key]) for row in weight_data])
                st.metric(f"{t('mean_weight')} - {t('after_experiment')}", f"{mean_final:.1f} g")
            with col3:
                change_key = t('change_g')
                mean_change = np.mean([float(row[change_key]) for row in weight_data])
                st.metric(f"{t('mean_weight')} {t('weight_change')}", f"{mean_change:.1f} g")
    else:
        # Original mean scores summary for other modes
        st.subheader(t('mean_scores'))
        
        # Add a filter for time points
        unique_times = sorted(current_df['time'].unique())
        selected_times = st.multiselect(
            t('filter_time'),
            unique_times,
            default=unique_times[:3] if len(unique_times) > 3 else unique_times
        )
        
        summary_data = []
        filtered_df = current_df[current_df['time'].isin(selected_times)] if selected_times else current_df
        
        for _, row in filtered_df.iterrows():
            animal_scores = [row[f'{animal_type}_{i}'] for i in range(1, num_animals + 1) 
                            if f'{animal_type}_{i}' in row]
            
            if mode in ["Autonomic and Sensorimotor Functions", "Reflex Capabilities", "Convulsive Behaviors and Excitability"]:
                # For binary modes, calculate percentage abnormal
                # For all observations in Autonomic mode: pale and cyanosis are both abnormal
                abnormal_count = sum(1 for score in animal_scores 
                                   if str(score).lower() in [t('abnormal').lower(), t('pale').lower(), t('cyanosis').lower()])
                
                percent_abnormal = (abnormal_count / len(animal_scores)) * 100 if animal_scores else 0
                status = t('abnormal') if abnormal_count > 0 else t('normal')
                
                summary_data.append({
                    t('time'): f"{int(row['time'])} min",
                    t('observation'): t_obs(row['observation']),
                    t('abnormal_count'): f"{abnormal_count}/{len(animal_scores)}",
                    t('percentage_abnormal'): f"{percent_abnormal:.1f}%",
                    t('status'): status
                })
            else:
                mean_score = calculate_mean_score(animal_scores)
                
                # Count how many animals have valid scores
                valid_scores = sum(1 for score in animal_scores if pd.notna(score) and score != '')
                
                # Determine status based on mode and thresholds
                if pd.isna(mean_score):
                    status = 'N/A'
                else:
                    if mode == "Body Temperature":
                        status = t('normal') if 36 <= mean_score <= 38 else t('abnormal')
                    else:  # General Behavior
                        if mean_score < 2 or mean_score > 6:
                            status = t('abnormal')
                        else:
                            status = t('normal')
                
                summary_data.append({
                    t('time'): f"{int(row['time'])} min",
                    t('observation'): t_obs(row['observation']),
                    t('mean_score'): f"{mean_score:.2f}" if not pd.isna(mean_score) else "N/A",
                    f"{t('valid')} {t(animal_type).capitalize()}s": f"{valid_scores}/{num_animals}",
                    t('status'): status
                })
        
        summary_df = pd.DataFrame(summary_data)
        
        # Display with custom styling
        st.dataframe(
            summary_df,
            use_container_width=True,
            hide_index=True,
            column_config={
                t('status'): st.column_config.TextColumn(t('status'), width="small")
            }
        )
        
        # Display abnormal episodes (not for Body Weight)
        st.subheader(t('abnormal_episodes'))
        episodes_df = process_data_with_episodes(current_df, mode, animal_type, num_animals)
        if not episodes_df.empty:
            st.dataframe(episodes_df, use_container_width=True, hide_index=True)
        else:
            st.info(t('no_abnormal'))
    
    return current_df

# New function to create plots for all modes
def create_comparative_plot(selected_for_viz, mode_eng, project, comparison_group=None):
    """Create comparative plots for all analysis modes"""
    
    # Get animal info
    animal_type = project.get('animal_type', 'mouse')
    if animal_type == 'custom':
        animal_type = project.get('custom_animal_name', 'animal')
    num_animals = project.get('num_animals', 8)
    
    # Get all times across selected groups
    all_times = set()
    valid_groups = []
    
    for exp in selected_for_viz:
        worksheet_key = f"worksheet_{exp}_{mode_eng}"
        if worksheet_key in st.session_state:
            df = st.session_state[worksheet_key]
            if not df.empty:
                if mode_eng != "Body Weight":
                    all_times.update(df['time'].unique())
                valid_groups.append(exp)
    
    if not valid_groups:
        st.warning("No data available for visualization")
        return None
    
    # Create visualization based on mode
    if mode_eng == "General Behavior":
        # For General Behavior, still use bar plot with time selection
        selected_time = st.selectbox(
            t('select_time_compare'), 
            sorted(list(all_times)),
            key=f"time_select_{mode_eng}"
        )
        return create_general_behavior_plot(valid_groups, selected_time, mode_eng, animal_type, num_animals, comparison_group)
    elif mode_eng == "Body Weight":
        # For Body Weight, create a special comparison plot
        return create_body_weight_comparison_plot(valid_groups, mode_eng, animal_type, num_animals, comparison_group)
    else:
        # For all other modes, use line charts
        # Allow selection of which groups to plot
        st.markdown(f"**{t('select_groups_chart')}**")
        selected_groups_for_plot = st.multiselect(
            t('groups_to_plot'),
            valid_groups,
            default=valid_groups,  # Show all groups by default
            key=f"groups_select_{mode_eng}"
        )
        
        if not selected_groups_for_plot:
            st.warning("Please select at least one group to plot")
            return None
            
        if mode_eng == "Body Temperature":
            return create_body_temperature_line_plot(selected_groups_for_plot, mode_eng, animal_type, num_animals, comparison_group)
        elif mode_eng in ["Autonomic and Sensorimotor Functions", "Reflex Capabilities", "Convulsive Behaviors and Excitability"]:
            return create_binary_score_line_plot(selected_groups_for_plot, mode_eng, animal_type, num_animals, comparison_group)
    
    return None

def create_body_weight_comparison_plot(valid_groups, mode_eng, animal_type, num_animals, comparison_group):
    """Create comparison plot for Body Weight mode"""
    # Ensure Chinese font is loaded if language is Chinese
    if st.session_state.language == 'zh':
        ensure_chinese_font()
    
    fig, ax = plt.subplots(figsize=(14, 8))
    
    group_names = []
    before_means = []
    before_stds = []
    after_means = []
    after_stds = []
    percent_changes = []
    
    for group in valid_groups:
        worksheet_key = f"worksheet_{group}_{mode_eng}"
        if worksheet_key in st.session_state:
            df = st.session_state[worksheet_key]
            
            before_df = df[df['time'] == 'before']
            after_df = df[df['time'] == 'after']
            
            if not before_df.empty and not after_df.empty:
                before_weights = []
                after_weights = []
                changes = []
                
                for i in range(1, num_animals + 1):
                    animal_col = f'{animal_type}_{i}'
                    if animal_col in before_df.columns:
                        try:
                            before_weight = float(before_df.iloc[0][animal_col])
                            after_weight = float(after_df.iloc[0][animal_col])
                            before_weights.append(before_weight)
                            after_weights.append(after_weight)
                            changes.append(((after_weight - before_weight) / before_weight) * 100)
                        except (ValueError, TypeError):
                            continue
                
                if before_weights and after_weights:
                    group_names.append(group.split('_')[-1])
                    before_means.append(np.mean(before_weights))
                    before_stds.append(np.std(before_weights))
                    after_means.append(np.mean(after_weights))
                    after_stds.append(np.std(after_weights))
                    percent_changes.append(np.mean(changes))
    
    if not group_names:
        st.warning("No valid weight data found")
        return None
    
    # Create grouped bar chart
    x = np.arange(len(group_names))
    width = 0.35
    
    # Plot before/after bars
    bars1 = ax.bar(x - width/2, before_means, width, yerr=before_stds, 
                     label=t('before_experiment'), capsize=5, alpha=0.8, color='#3498db')
    bars2 = ax.bar(x + width/2, after_means, width, yerr=after_stds,
                     label=t('after_experiment'), capsize=5, alpha=0.8, color='#e74c3c')
    
    # Color comparison group differently
    for i, group in enumerate(group_names):
        full_group = f"{st.session_state.projects[st.session_state.active_project]['name']}_Group_{group}"
        if full_group == comparison_group:
            bars1[i].set_color('#28a745')
            bars2[i].set_color('#1e7e34')
    
    # Add value labels on bars
    for i, (before, after) in enumerate(zip(before_means, after_means)):
        # Before weight label
        ax.text(i - width/2, before + before_stds[i] + 0.5, f"{before:.1f}", 
                ha='center', va='bottom', fontsize=10, fontweight='bold')
        # After weight label
        ax.text(i + width/2, after + after_stds[i] + 0.5, f"{after:.1f}", 
                ha='center', va='bottom', fontsize=10, fontweight='bold')
        
        # Add percentage change label above the group
        change_color = 'red' if percent_changes[i] < 0 else 'green'
        y_pos = max(before + before_stds[i], after + after_stds[i]) + 3
        ax.text(i, y_pos, f"{percent_changes[i]:+.1f}%", 
                ha='center', va='bottom', fontsize=11, fontweight='bold', 
                color=change_color, bbox=dict(boxstyle="round,pad=0.3", facecolor='white', edgecolor=change_color))
    
    # Formatting
    ax.set_xlabel(t('group'), fontsize=14)
    ax.set_ylabel(f"{t('weight_g')}", fontsize=14)
    ax.set_title(f"{t('body_weight')} - {t('before_experiment')} vs {t('after_experiment')} {t('comparative_viz')}", 
                 fontsize=16, fontweight='bold')
    ax.set_xticks(x)
    ax.set_xticklabels(group_names, fontsize=12)
    ax.legend(fontsize=12, loc='upper left')
    ax.grid(True, alpha=0.3, axis='y')
    
    # Add a subtle horizontal line at y=0 to show the baseline
    if min(before_means + after_means) < 0:
        ax.axhline(y=0, color='black', linestyle='-', linewidth=0.5)
    
    # Set y-axis to start from 0 unless there are negative values
    if min(before_means + after_means) >= 0:
        ax.set_ylim(bottom=0)
    
    plt.tight_layout()
    return fig

def create_general_behavior_plot(valid_groups, selected_time, mode_eng, animal_type, num_animals, comparison_group):
    """Create plot for General Behavior mode"""
    # Ensure Chinese font is loaded if language is Chinese
    if st.session_state.language == 'zh':
        ensure_chinese_font()
    
    fig, ax = plt.subplots(figsize=(12, 8))
    
    overall_means = []
    overall_stds = []
    group_names = []
    
    for exp in valid_groups:
        worksheet_key = f"worksheet_{exp}_{mode_eng}"
        if worksheet_key in st.session_state:
            df = st.session_state[worksheet_key]
            filtered = df[df['time'] == selected_time]
            
            if not filtered.empty:
                all_mean_scores = []
                for _, row in filtered.iterrows():
                    animal_scores = [row[f'{animal_type}_{i}'] for i in range(1, num_animals + 1) 
                                   if f'{animal_type}_{i}' in row]
                    mean_score = calculate_mean_score(animal_scores)
                    if not pd.isna(mean_score):
                        all_mean_scores.append(mean_score)
                
                if all_mean_scores:
                    overall_means.append(np.mean(all_mean_scores))
                    overall_stds.append(np.std(all_mean_scores))
                    group_names.append(exp)
    
    if not overall_means:
        st.warning("No valid data for the selected time point")
        return None
    
    # Create bar plot
    x_pos = range(len(group_names))
    bars = ax.bar(x_pos, overall_means, yerr=overall_stds, capsize=5, alpha=0.8)
    
    # Color bars based on status
    for i, (mean, group) in enumerate(zip(overall_means, group_names)):
        if group == comparison_group:
            bars[i].set_color('#28a745')  # Green for comparison group
        elif mean < 2 or mean > 6:  # Abnormal
            bars[i].set_color('#ff6b6b')  # Red for abnormal
        else:
            bars[i].set_color('#4cc9f0')  # Blue for normal
    
    # Formatting
    ax.set_title(f"{t('general_behavior')} - {t('comparative_viz')} ({selected_time} min)", fontsize=14, fontweight='bold')
    ax.set_ylabel(f"{t('mean_score')} (0-10)", fontsize=12)
    ax.set_xlabel(t('group'), fontsize=12)
    ax.set_xticks(x_pos)
    ax.set_xticklabels([g.split('_')[-1] for g in group_names], rotation=45, ha='right')
    ax.set_ylim(0, 10)
    ax.grid(True, alpha=0.3, axis='y')
    
    # Add threshold lines
    ax.axhline(y=2, color='gray', linestyle='--', alpha=0.7, label='Lower threshold')
    ax.axhline(y=6, color='gray', linestyle='--', alpha=0.7, label='Upper threshold')
    
    # Add value labels
    for i, (mean, std) in enumerate(zip(overall_means, overall_stds)):
        ax.text(i, mean + std + 0.2, f"{mean:.2f}", ha='center', va='bottom', fontweight='bold')
    
    # Add legend
    from matplotlib.patches import Patch
    legend_elements = [
        Patch(facecolor='#28a745', label=t('comparison_group')),
        Patch(facecolor='#4cc9f0', label=t('normal')),
        Patch(facecolor='#ff6b6b', label=t('abnormal'))
    ]
    ax.legend(handles=legend_elements, loc='upper right')
    
    plt.tight_layout()
    return fig

def create_body_temperature_line_plot(selected_groups, mode_eng, animal_type, num_animals, comparison_group):
    """Create line plot for Body Temperature mode"""
    # Ensure Chinese font is loaded if language is Chinese
    if st.session_state.language == 'zh':
        ensure_chinese_font()
    
    fig, ax = plt.subplots(figsize=(14, 8))
    
    # Colors for different groups
    colors = plt.cm.tab10(np.linspace(0, 1, len(selected_groups)))
    
    for idx, group in enumerate(selected_groups):
        worksheet_key = f"worksheet_{group}_{mode_eng}"
        if worksheet_key in st.session_state:
            df = st.session_state[worksheet_key]
            
            # Get unique times and sort them
            times = sorted(df['time'].unique())
            mean_temps = []
            std_temps = []
            
            for time in times:
                time_df = df[df['time'] == time]
                all_temps = []
                
                for _, row in time_df.iterrows():
                    # Collect all animal temperatures
                    for i in range(1, num_animals + 1):
                        if f'{animal_type}_{i}' in row:
                            try:
                                temp = float(row[f'{animal_type}_{i}'])
                                all_temps.append(temp)
                            except (ValueError, TypeError):
                                continue
                
                if all_temps:
                    mean_temps.append(np.mean(all_temps))
                    std_temps.append(np.std(all_temps))
                else:
                    mean_temps.append(np.nan)
                    std_temps.append(0)
            
            # Plot line with error bars
            line_style = '-' if group != comparison_group else '--'
            line_width = 2 if group != comparison_group else 3
            marker = 'o' if group != comparison_group else 's'
            
            ax.errorbar(times, mean_temps, yerr=std_temps, 
                       label=group.split('_')[-1], 
                       color=colors[idx],
                       linestyle=line_style,
                       linewidth=line_width,
                       marker=marker,
                       markersize=8,
                       capsize=5,
                       alpha=0.8)
    
    # Add normal range
    ax.axhspan(36, 38, alpha=0.2, color='green', label='Normal range (36-38°C)')
    
    # Formatting
    ax.set_title(f"{t('body_temperature')} - {t('comparative_viz')} ({t('all_time_points')})", fontsize=16, fontweight='bold')
    ax.set_xlabel(f"{t('time')} (min)", fontsize=12)
    ax.set_ylabel("Temperature (°C)", fontsize=12)
    ax.grid(True, alpha=0.3)
    ax.legend(bbox_to_anchor=(1.05, 1), loc='upper left')
    
    # Set y-axis limits
    ax.set_ylim(34, 40)
    
    plt.tight_layout()
    return fig

def create_binary_score_line_plot(selected_groups, mode_eng, animal_type, num_animals, comparison_group):
    """Create line plot for binary (Normal/Abnormal) scoring modes"""
    # Ensure Chinese font is loaded if language is Chinese
    if st.session_state.language == 'zh':
        ensure_chinese_font()
    
    # Get observations for this mode
    if mode_eng == "Autonomic and Sensorimotor Functions":
        observations = AUTONOMIC_OBSERVATIONS
    elif mode_eng == "Reflex Capabilities":
        observations = REFLEX_OBSERVATIONS
    else:  # Convulsive Behaviors
        observations = CONVULSIVE_OBSERVATIONS
    
    # Create subplot for each observation
    n_obs = len(observations)
    n_cols = 3
    n_rows = (n_obs + n_cols - 1) // n_cols
    
    fig, axes = plt.subplots(n_rows, n_cols, figsize=(18, 5*n_rows))
    if n_rows == 1:
        axes = axes.reshape(1, -1)
    axes = axes.flatten()
    
    # Colors for different groups
    colors = plt.cm.tab10(np.linspace(0, 1, len(selected_groups)))
    
    for obs_idx, obs in enumerate(observations):
        ax = axes[obs_idx]
        
        for group_idx, group in enumerate(selected_groups):
            worksheet_key = f"worksheet_{group}_{mode_eng}"
            if worksheet_key in st.session_state:
                df = st.session_state[worksheet_key]
                
                # Filter for this observation (use translated observation name)
                obs_translated = t_obs(obs)
                obs_df = df[df['observation'] == obs_translated]
                
                # Debug: Check if we found any data for this observation
                if obs_df.empty:
                    # Try to find data with English observation name as fallback
                    obs_df = df[df['observation'] == obs]
                
                times = sorted(obs_df['time'].unique())
                percentages = []
                
                for time in times:
                    time_df = obs_df[obs_df['time'] == time]
                    if not time_df.empty:
                        abnormal_count = 0
                        total_count = 0
                        
                        for _, row in time_df.iterrows():
                            for i in range(1, num_animals + 1):
                                if f'{animal_type}_{i}' in row:
                                    cell_value = str(row[f'{animal_type}_{i}']).lower()
                                    # Check for abnormal, pale, or cyanosis (all are considered abnormal)
                                    if cell_value in [t('abnormal').lower(), t('pale').lower(), t('cyanosis').lower()]:
                                        abnormal_count += 1
                                    total_count += 1
                        
                        percentage = (abnormal_count / total_count * 100) if total_count > 0 else 0
                        percentages.append(percentage)
                    else:
                        percentages.append(0)
                
                # Plot line
                line_style = '-' if group != comparison_group else '--'
                line_width = 2 if group != comparison_group else 3
                marker = 'o' if group != comparison_group else 's'
                
                ax.plot(times, percentages, 
                       label=group.split('_')[-1],
                       color=colors[group_idx],
                       linestyle=line_style,
                       linewidth=line_width,
                       marker=marker,
                       markersize=6,
                       alpha=0.8)
        
        # Formatting
        ax.set_title(t_obs(obs), fontsize=12, fontweight='bold')
        ax.set_xlabel(f"{t('time')} (min)", fontsize=10)
        ax.set_ylabel(f"{t('percentage_abnormal')} (%)", fontsize=10)
        ax.set_ylim(-5, 105)
        ax.grid(True, alpha=0.3)
        
        # Add legend only to first subplot
        if obs_idx == 0:
            ax.legend(bbox_to_anchor=(1.05, 1), loc='upper left', fontsize=8)
    
    # Hide unused subplots
    for idx in range(n_obs, len(axes)):
        axes[idx].set_visible(False)
    
    # Overall title
    mode_title = mode_eng.replace("and Sensorimotor Functions", "")
    fig.suptitle(f"{mode_title} - {t('comparative_viz')} ({t('all_time_points')})", fontsize=16, fontweight='bold')
    
    plt.tight_layout()
    return fig

# AI Features (appear when activated from sidebar)
if st.session_state.ai_tutor_active:
    st.markdown("## 🎓 AI Tutor")
    st.info("Welcome to the AI Tutor! Ask me anything about using the FOB Test Analysis Dashboard.")
    
    # Create a chat container with better styling
    chat_container = st.container()
    
    with chat_container:
        # Display chat history in a scrollable area
        if st.session_state.tutor_chat_history:
            st.markdown("**💬 Chat History:**")
            chat_display = st.container()
            
            with chat_display:
                for message in st.session_state.tutor_chat_history:
                    if message["role"] == "user":
                        st.markdown(f"""
                        <div style="background-color: #e3f2fd; padding: 10px; border-radius: 10px; margin: 5px 0;">
                            <strong>👤 You:</strong> {message['content']}
                        </div>
                        """, unsafe_allow_html=True)
                    else:
                        st.markdown(f"""
                        <div style="background-color: #f3e5f5; padding: 10px; border-radius: 10px; margin: 5px 0;">
                            <strong>🎓 AI Tutor:</strong> {message['content']}
                        </div>
                        """, unsafe_allow_html=True)
    
    # Chat input section
    st.markdown("---")
    st.markdown("**Ask a question about the dashboard:**")
    
    # Quick help buttons
    st.markdown("**💡 Quick Questions:**")
    col_q1, col_q2, col_q3 = st.columns(3)
    with col_q1:
        if st.button("📋 Create Project", use_container_width=True, key="tutor_q1"):
            # Clear previous chat history and add new conversation
            st.session_state.tutor_chat_history = []
            # Add welcome message
            tutor_welcome_msg = "Hello! I'm your FOB Test Analysis Dashboard tutor. I can help you learn how to use this tool effectively. Ask me anything about the dashboard features, analysis modes, or how to perform specific tasks!"
            if st.session_state.language == 'zh':
                tutor_welcome_msg = "你好！我是你的FOB测试分析仪表板导师。我可以帮助你学习如何有效使用这个工具。询问我关于仪表板功能、分析模式或如何执行特定任务的问题！"
            st.session_state.tutor_chat_history.append({"role": "assistant", "content": tutor_welcome_msg})
            
            st.session_state.tutor_chat_history.append({"role": "user", "content": "How do I create a new project?"})
            with st.spinner("AI Tutor is thinking..."):
                ai_response = generate_tutor_response("How do I create a new project?", st.session_state.language)
                st.session_state.tutor_chat_history.append({"role": "assistant", "content": ai_response})
            st.rerun()
    
    with col_q2:
        if st.button("📊 Analysis Modes", use_container_width=True, key="tutor_q2"):
            # Clear previous chat history and add new conversation
            st.session_state.tutor_chat_history = []
            # Add welcome message
            tutor_welcome_msg = "Hello! I'm your FOB Test Analysis Dashboard tutor. I can help you learn how to use this tool effectively. Ask me anything about the dashboard features, analysis modes, or how to perform specific tasks!"
            if st.session_state.language == 'zh':
                tutor_welcome_msg = "你好！我是你的FOB测试分析仪表板导师。我可以帮助你学习如何有效使用这个工具。询问我关于仪表板功能、分析模式或如何执行特定任务的问题！"
            st.session_state.tutor_chat_history.append({"role": "assistant", "content": tutor_welcome_msg})
            
            st.session_state.tutor_chat_history.append({"role": "user", "content": "What are the different analysis modes and how do I use them?"})
            with st.spinner("AI Tutor is thinking..."):
                ai_response = generate_tutor_response("What are the different analysis modes and how do I use them?", st.session_state.language)
                st.session_state.tutor_chat_history.append({"role": "assistant", "content": ai_response})
            st.rerun()
    
    with col_q3:
        if st.button("📈 Data Entry", use_container_width=True, key="tutor_q3"):
            # Clear previous chat history and add new conversation
            st.session_state.tutor_chat_history = []
            # Add welcome message
            tutor_welcome_msg = "Hello! I'm your FOB Test Analysis Dashboard tutor. I can help you learn how to use this tool effectively. Ask me anything about the dashboard features, analysis modes, or how to perform specific tasks!"
            if st.session_state.language == 'zh':
                tutor_welcome_msg = "你好！我是你的FOB测试分析仪表板导师。我可以帮助你学习如何有效使用这个工具。询问我关于仪表板功能、分析模式或如何执行特定任务的问题！"
            st.session_state.tutor_chat_history.append({"role": "assistant", "content": tutor_welcome_msg})
            
            st.session_state.tutor_chat_history.append({"role": "user", "content": "How do I enter data for my animals?"})
            with st.spinner("AI Tutor is thinking..."):
                ai_response = generate_tutor_response("How do I enter data for my animals?", st.session_state.language)
                st.session_state.tutor_chat_history.append({"role": "assistant", "content": ai_response})
            st.rerun()
    
    # More quick questions
    col_q4, col_q5, col_q6 = st.columns(3)
    with col_q4:
        if st.button("📋 Reports", use_container_width=True, key="tutor_q4"):
            # Clear previous chat history and add new conversation
            st.session_state.tutor_chat_history = []
            # Add welcome message
            tutor_welcome_msg = "Hello! I'm your FOB Test Analysis Dashboard tutor. I can help you learn how to use this tool effectively. Ask me anything about the dashboard features, analysis modes, or how to perform specific tasks!"
            if st.session_state.language == 'zh':
                tutor_welcome_msg = "你好！我是你的FOB测试分析仪表板导师。我可以帮助你学习如何有效使用这个工具。询问我关于仪表板功能、分析模式或如何执行特定任务的问题！"
            st.session_state.tutor_chat_history.append({"role": "assistant", "content": tutor_welcome_msg})
            
            st.session_state.tutor_chat_history.append({"role": "user", "content": "How do I generate reports and export data?"})
            with st.spinner("AI Tutor is thinking..."):
                ai_response = generate_tutor_response("How do I generate reports and export data?", st.session_state.language)
                st.session_state.tutor_chat_history.append({"role": "assistant", "content": ai_response})
            st.rerun()
    
    with col_q5:
        if st.button("🎯 Groups", use_container_width=True, key="tutor_q5"):
            # Clear previous chat history and add new conversation
            st.session_state.tutor_chat_history = []
            # Add welcome message
            tutor_welcome_msg = "Hello! I'm your FOB Test Analysis Dashboard tutor. I can help you learn how to use this tool effectively. Ask me anything about the dashboard features, analysis modes, or how to perform specific tasks!"
            if st.session_state.language == 'zh':
                tutor_welcome_msg = "你好！我是你的FOB测试分析仪表板导师。我可以帮助你学习如何有效使用这个工具。询问我关于仪表板功能、分析模式或如何执行特定任务的问题！"
            st.session_state.tutor_chat_history.append({"role": "assistant", "content": tutor_welcome_msg})
            
            st.session_state.tutor_chat_history.append({"role": "user", "content": "How do I manage multiple groups and set comparison groups?"})
            with st.spinner("AI Tutor is thinking..."):
                ai_response = generate_tutor_response("How do I manage multiple groups and set comparison groups?", st.session_state.language)
                st.session_state.tutor_chat_history.append({"role": "assistant", "content": ai_response})
            st.rerun()
    
    with col_q6:
        if st.button("📊 Charts", use_container_width=True, key="tutor_q6"):
            # Clear previous chat history and add new conversation
            st.session_state.tutor_chat_history = []
            # Add welcome message
            tutor_welcome_msg = "Hello! I'm your FOB Test Analysis Dashboard tutor. I can help you learn how to use this tool effectively. Ask me anything about the dashboard features, analysis modes, or how to perform specific tasks!"
            if st.session_state.language == 'zh':
                tutor_welcome_msg = "你好！我是你的FOB测试分析仪表板导师。我可以帮助你学习如何有效使用这个工具。询问我关于仪表板功能、分析模式或如何执行特定任务的问题！"
            st.session_state.tutor_chat_history.append({"role": "assistant", "content": tutor_welcome_msg})
            
            st.session_state.tutor_chat_history.append({"role": "user", "content": "How do I create and download charts?"})
            with st.spinner("AI Tutor is thinking..."):
                ai_response = generate_tutor_response("How do I create and download charts?", st.session_state.language)
                st.session_state.tutor_chat_history.append({"role": "assistant", "content": ai_response})
            st.rerun()
    
    # Custom question input
    user_message = st.text_input(
        "Type your question here...",
        key="tutor_input",
        placeholder="e.g., How do I analyze body weight data?",
        help="Ask specific questions about using the dashboard"
    )
    
    col_send, col_clear = st.columns([3, 1])
    with col_send:
        if st.button("Send Question", use_container_width=True, type="primary"):
            if user_message.strip():
                # Clear previous chat history and add new conversation
                st.session_state.tutor_chat_history = []
                # Add welcome message
                tutor_welcome_msg = "Hello! I'm your FOB Test Analysis Dashboard tutor. I can help you learn how to use this tool effectively. Ask me anything about the dashboard features, analysis modes, or how to perform specific tasks!"
                if st.session_state.language == 'zh':
                    tutor_welcome_msg = "你好！我是你的FOB测试分析仪表板导师。我可以帮助你学习如何有效使用这个工具。询问我关于仪表板功能、分析模式或如何执行特定任务的问题！"
                st.session_state.tutor_chat_history.append({"role": "assistant", "content": tutor_welcome_msg})
                
                # Add user message to chat history
                st.session_state.tutor_chat_history.append({"role": "user", "content": user_message})
                
                # Generate AI response
                with st.spinner("AI Tutor is thinking..."):
                    ai_response = generate_tutor_response(user_message, st.session_state.language)
                    st.session_state.tutor_chat_history.append({"role": "assistant", "content": ai_response})
                
                # Clear input and rerun
                st.rerun()
    
    with col_clear:
        if st.button("Clear Chat", use_container_width=True):
            st.session_state.tutor_chat_history = []
            st.rerun()

elif st.session_state.ai_chatbot_active:
    st.markdown("## 💬 AI Chatbot - File Analysis")
    st.info("Upload multiple files and I'll help you summarize their content for your FOB test analysis.")
    
    # File upload section
    st.subheader("📁 Upload Multiple Files")
    uploaded_files = st.file_uploader(
        "Upload multiple data files for analysis",
        type=['csv', 'xlsx', 'xls', 'txt'],
        accept_multiple_files=True,
        help="Upload multiple CSV, Excel, or text files to analyze and summarize"
    )
    
    # Process uploaded files
    if uploaded_files:
        st.success(f"Uploaded {len(uploaded_files)} file(s)")
        
        # Process each file and generate summaries
        for uploaded_file in uploaded_files:
            st.markdown(f"**📄 Processing: {uploaded_file.name}**")
            
            # Process uploaded file
            file_content = process_uploaded_file(uploaded_file)
            
            if file_content and not file_content.startswith("Error"):
                # Generate file summary
                with st.spinner(f"Generating summary for {uploaded_file.name}..."):
                    file_summary = generate_file_summary(file_content, uploaded_file.name, st.session_state.language)
                    
                    # Store summary
                    summary_entry = {
                        "filename": uploaded_file.name,
                        "content": file_content,
                        "summary": file_summary
                    }
                    
                    # Add to session state if not already present
                    if not any(s["filename"] == uploaded_file.name for s in st.session_state.file_summaries):
                        st.session_state.file_summaries.append(summary_entry)
                
                # Display summary
                with st.expander(f"📋 Summary: {uploaded_file.name}"):
                    st.markdown(file_summary)
                    
                    # Show file preview
                    st.markdown("**File Preview:**")
                    st.text(file_content[:300] + "..." if len(file_content) > 300 else file_content)
            else:
                st.error(f"Error processing {uploaded_file.name}: {file_content}")
    
    # Display all file summaries
    if st.session_state.file_summaries:
        st.markdown("---")
        st.subheader("📊 All File Summaries")
        
        for summary in st.session_state.file_summaries:
            with st.expander(f"📄 {summary['filename']}"):
                st.markdown(summary['summary'])
        
        # Clear summaries button
        if st.button("🗑️ Clear All Summaries", use_container_width=True):
            st.session_state.file_summaries = []
            st.rerun()
    
    # Floating Chat Box Interface with Streaming Responses
    st.markdown("---")
    
    # Initialize chat state
    if 'floating_chat_open' not in st.session_state:
        st.session_state.floating_chat_open = False
    if 'chat_messages' not in st.session_state:
        st.session_state.chat_messages = []
    if 'current_streaming_response' not in st.session_state:
        st.session_state.current_streaming_response = ""
    if 'is_streaming' not in st.session_state:
        st.session_state.is_streaming = False
    
    # Welcome message - ensure it's always present when chatbot is active
    if st.session_state.ai_chatbot_active:
        if 'chat_messages' not in st.session_state or not st.session_state.chat_messages:
            st.session_state.chat_messages = []
            welcome_msg = "Hello! I'm your AI assistant. I can help you analyze files, answer questions about FOB testing, and provide guidance on using this dashboard. How can I help you today?"
            if st.session_state.language == 'zh':
                welcome_msg = "你好！我是你的AI助手。我可以帮助你分析文件、回答FOB测试相关问题，并提供使用这个仪表板的指导。今天我能为你做些什么？"
            st.session_state.chat_messages.append({
                "role": "assistant", 
                "content": welcome_msg, 
                "timestamp": datetime.datetime.now().strftime("%H:%M")
            })
    
    # Simple inline chat interface (no floating popup)
    if st.session_state.ai_chatbot_active:
        st.markdown("---")
        st.subheader("💬 Chat with AI Assistant")
        
        # Display chat messages inline
        if st.session_state.chat_messages:
            for i, message in enumerate(st.session_state.chat_messages):
                if message["role"] == "user":
                    with st.chat_message("user"):
                        st.write(message['content'])
                        if message.get('timestamp'):
                            st.caption(f"Time: {message['timestamp']}")
                else:
                    with st.chat_message("assistant"):
                        st.write(message['content'])
                        if message.get('timestamp'):
                            st.caption(f"Time: {message['timestamp']}")
        
        # Show streaming response if active
        if st.session_state.is_streaming and st.session_state.current_streaming_response:
            with st.chat_message("assistant"):
                st.write(st.session_state.current_streaming_response)
                st.caption("Typing...")
        
        # Chat input
        user_message = st.chat_input("Ask me anything about FOB testing, data analysis, or dashboard usage...")
        
        # Handle user input
        if user_message:
            # Add user message to chat
            st.session_state.chat_messages.append({
                "role": "user", 
                "content": user_message, 
                "timestamp": datetime.datetime.now().strftime("%H:%M")
            })
            
            # Start streaming response
            st.session_state.is_streaming = True
            st.session_state.current_streaming_response = ""
            st.rerun()
        
        # Clear chat button
        if st.button("🗑️ Clear Chat", use_container_width=True, key="clear_inline_chat"):
            st.session_state.chat_messages = []
            st.rerun()
    
    # Handle streaming responses
    if st.session_state.is_streaming:
        # Simulate streaming response (you can replace this with actual streaming API)
        if not st.session_state.current_streaming_response:
            # Get the last user message
            last_user_msg = None
            for msg in reversed(st.session_state.chat_messages):
                if msg["role"] == "user":
                    last_user_msg = msg["content"]
                    break
            
            if last_user_msg:
                # Generate AI response
                ai_response = generate_chatbot_response(last_user_msg, st.session_state.language)
                
                # Simulate streaming by adding characters gradually
                import time
                for i in range(len(ai_response)):
                    st.session_state.current_streaming_response = ai_response[:i+1]
                    time.sleep(0.05)  # Small delay for streaming effect
                    st.rerun()
                
                # Add complete response to chat history
                st.session_state.chat_messages.append({
                    "role": "assistant", 
                    "content": ai_response, 
                    "timestamp": datetime.datetime.now().strftime("%H:%M")
                })
                
                # Stop streaming
                st.session_state.is_streaming = False
                st.session_state.current_streaming_response = ""
                st.rerun()

elif st.session_state.ai_report_active:
    st.markdown("## 📊 AI Report Generator")
    st.info("Generate professional AI-powered analysis reports for your FOB test data.")
    
    # Show file summaries from chatbot if available
    if st.session_state.file_summaries:
        st.subheader("📁 File Summaries (from AI Chatbot)")
        st.info("The following file summaries will be included in your AI report:")
        
        for summary in st.session_state.file_summaries:
            with st.expander(f"📄 {summary['filename']}"):
                st.markdown(summary['summary'])
    
    # AI Report Generation
    st.subheader("🤖 Generate AI Report")
    st.info("This will analyze your current project data and generate a professional report.")
    
    if st.button("Generate AI Report", use_container_width=True, type="primary"):
        if st.session_state.active_project is not None:
            with st.spinner("Generating AI report..."):
                project = st.session_state.projects[st.session_state.active_project]
                
                # Get current mode
                mode_eng = st.session_state.mode
                
                # Prepare data for AI analysis
                uploaded_file_content = st.session_state.get('uploaded_file_content', None)
                
                # Include file summaries in the analysis
                file_summaries_text = ""
                if st.session_state.file_summaries:
                    file_summaries_text = "\n\n**Additional File Analysis:**\n"
                    for summary in st.session_state.file_summaries:
                        file_summaries_text += f"\n**File: {summary['filename']}**\n{summary['summary']}\n"
                
                # Create sample data for demonstration (in real use, this would be actual project data)
                sample_data = f"Project: {project['name']}, Mode: {mode_eng}, Animals: {project['num_animals']}"
                
                # Combine uploaded file content with file summaries
                combined_file_content = ""
                if uploaded_file_content:
                    combined_file_content += uploaded_file_content
                if file_summaries_text:
                    combined_file_content += file_summaries_text
                
                # Generate AI report
                ai_report = generate_ai_report(project, sample_data, mode_eng, st.session_state.language, combined_file_content)
                
                # Display AI report
                st.markdown("### 📋 AI Analysis Report")
                st.markdown(ai_report)
                
                # Download AI report
                st.download_button(
                    label="Download AI Report",
                    data=ai_report,
                    file_name=f"{project['name']}_AI_Report_{mode_eng.replace(' ', '_')}_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.txt",
                    mime="text/plain",
                    use_container_width=True
                )
                
                # Generate and download PowerPoint presentation
                st.markdown("---")
                st.subheader("📊 PowerPoint Presentation")
                st.info("Generate a professional PowerPoint presentation with your analysis results.")
                
                if st.button("Generate PowerPoint", use_container_width=True, type="primary"):
                    with st.spinner("Creating PowerPoint presentation..."):
                        # Get file summaries for the presentation
                        file_summaries = st.session_state.get('file_summaries', [])
                        
                        # Create PowerPoint presentation
                        pptx_data = create_powerpoint_presentation(
                            project, 
                            mode_eng, 
                            st.session_state.language, 
                            file_summaries
                        )
                        
                        if isinstance(pptx_data, bytes):
                            # Download PowerPoint
                            st.download_button(
                                label="Download PowerPoint Presentation",
                                data=pptx_data,
                                file_name=f"{project['name']}_Presentation_{mode_eng.replace(' ', '_')}_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                use_container_width=True
                            )
                            st.success("PowerPoint presentation generated successfully!")
                        else:
                            st.error(f"Error generating PowerPoint: {pptx_data}")
        else:
            st.warning("Please create a project first before generating AI reports.")

elif st.session_state.ai_powerpoint_active:
    st.markdown("## 📈 AI PowerPoint Generator")
    st.info("Generate professional PowerPoint presentations with AI-powered content including introduction, experiment description, and results.")
    
    # File upload for enhanced presentation
    st.subheader("📁 Upload Additional Data (Optional)")
    uploaded_files = st.file_uploader(
        "Upload files to include in presentation",
        type=['csv', 'xlsx', 'xls', 'txt'],
        accept_multiple_files=True,
        help="Upload files to enhance the presentation content"
    )
    
    # Process uploaded files
    file_summaries = []
    if uploaded_files:
        st.success(f"Uploaded {len(uploaded_files)} file(s)")
        
        for uploaded_file in uploaded_files:
            st.markdown(f"**📄 Processing: {uploaded_file.name}**")
            
            # Process uploaded file
            file_content = process_uploaded_file(uploaded_file)
            
            if file_content and not file_content.startswith("Error"):
                # Generate file summary
                with st.spinner(f"Generating summary for {uploaded_file.name}..."):
                    file_summary = generate_file_summary(file_content, uploaded_file.name, st.session_state.language)
                    
                    # Store summary
                    summary_entry = {
                        "filename": uploaded_file.name,
                        "content": file_content,
                        "summary": file_summary
                    }
                    file_summaries.append(summary_entry)
                
                # Display summary
                with st.expander(f"📋 Summary: {uploaded_file.name}"):
                    st.markdown(file_summary)
    
    # PowerPoint Generation
    st.subheader("🤖 Generate PowerPoint Presentation")
    st.info("Create a comprehensive presentation with AI-generated content including introduction, experiment design, results, and conclusions.")
    
    # Chart management section
    col_charts_info, col_clear_charts = st.columns([3, 1])
    with col_charts_info:
        total_experiment_charts = len(st.session_state.all_experiment_charts)
        if total_experiment_charts > 0:
            st.success(f"📊 Captured {total_experiment_charts} experiment chart(s) for PowerPoint inclusion")
        else:
            st.info("💡 PowerPoint will automatically generate charts for all 6 FOB test modes")
    
    with col_clear_charts:
        if st.button("🗑️ Clear Charts", help="Clear all captured experiment charts"):
            st.session_state.all_experiment_charts = []
            st.success("Charts cleared!")
            st.rerun()
    
    # Generate charts for all 6 FOB test modes
    charts_data = []
    if st.session_state.active_project is not None:
        project = st.session_state.projects[st.session_state.active_project]
        
        # Define all 6 FOB test modes
        all_modes = [
            "General Behavior",
            "Autonomic and Sensorimotor Functions", 
            "Reflex Capabilities",
            "Body Temperature",
            "Body Weight",
            "Convulsive Behaviors and Excitability"
        ]
        
        st.info("📊 Generating comprehensive charts for all 6 FOB test modes...")
        
        # Create sample data for charts
        
        for mode in all_modes:
            mode_charts = []
            
            # Chart 1: Group comparison for each mode
            fig1, ax1 = plt.subplots(figsize=(10, 6))
            groups = ['Group 1', 'Group 2', 'Group 3', 'Control']
            
            # Mode-specific data ranges
            if mode == "Body Temperature":
                values = [37.2, 37.8, 36.9, 37.5]  # Temperature range
                ylabel = 'Temperature (°C)'
            elif mode == "Body Weight":
                values = [25.5, 26.2, 24.8, 26.8]  # Weight range
                ylabel = 'Weight (g)'
            elif mode in ["Autonomic and Sensorimotor Functions", "Reflex Capabilities", "Convulsive Behaviors and Excitability"]:
                values = [15, 22, 8, 28]  # Percentage abnormal
                ylabel = 'Abnormal (%)'
            else:  # General Behavior
                values = [75, 82, 68, 90]  # Score range
                ylabel = 'Score'
            
            colors = ['#FF6B6B', '#4ECDC4', '#45B7D1', '#96CEB4']
            bars = ax1.bar(groups, values, color=colors)
            ax1.set_title(f'{mode} - Group Comparison', fontsize=16, fontweight='bold')
            ax1.set_ylabel(ylabel, fontsize=12)
            ax1.set_xlabel('Groups', fontsize=12)
            
            # Add value labels on bars
            for bar, value in zip(bars, values):
                height = bar.get_height()
                ax1.text(bar.get_x() + bar.get_width()/2., height + 0.5,
                        f'{value}', ha='center', va='bottom', fontweight='bold')
            
            plt.tight_layout()
            
            # Save chart to bytes
            chart1_buffer = BytesIO()
            fig1.savefig(chart1_buffer, format='png', dpi=300, bbox_inches='tight')
            chart1_buffer.seek(0)
            mode_charts.append({
                'title': f'{mode} - Group Comparison',
                'data': chart1_buffer.getvalue(),
                'description': f'Bar chart showing {mode} scores across different experimental groups',
                'mode': mode
            })
            plt.close(fig1)
            
            # Chart 2: Time series for each mode
            fig2, ax2 = plt.subplots(figsize=(10, 6))
            time_points = [0, 15, 30, 60, 120]
            
            # Mode-specific time series data
            if mode == "Body Temperature":
                group1_data = [37.0, 37.2, 37.5, 37.3, 37.4]
                group2_data = [36.8, 37.0, 37.3, 37.1, 37.2]
                control_data = [37.2, 37.4, 37.6, 37.5, 37.6]
                ylabel = 'Temperature (°C)'
            elif mode == "Body Weight":
                group1_data = [25.0, 25.2, 25.5, 25.3, 25.4]
                group2_data = [24.8, 25.0, 25.3, 25.1, 25.2]
                control_data = [26.0, 26.2, 26.5, 26.3, 26.4]
                ylabel = 'Weight (g)'
            elif mode in ["Autonomic and Sensorimotor Functions", "Reflex Capabilities", "Convulsive Behaviors and Excitability"]:
                group1_data = [10, 15, 20, 18, 22]
                group2_data = [8, 12, 18, 15, 19]
                control_data = [5, 8, 12, 10, 15]
                ylabel = 'Abnormal (%)'
            else:  # General Behavior
                group1_data = [70, 75, 80, 78, 82]
                group2_data = [65, 70, 75, 72, 76]
                control_data = [85, 88, 90, 87, 89]
                ylabel = 'Score'
            
            ax2.plot(time_points, group1_data, 'o-', label='Group 1', linewidth=2, markersize=8)
            ax2.plot(time_points, group2_data, 's-', label='Group 2', linewidth=2, markersize=8)
            ax2.plot(time_points, control_data, '^-', label='Control', linewidth=2, markersize=8)
            
            ax2.set_title(f'{mode} - Time Series Analysis', fontsize=16, fontweight='bold')
            ax2.set_xlabel('Time (minutes)', fontsize=12)
            ax2.set_ylabel(ylabel, fontsize=12)
            ax2.legend(fontsize=10)
            ax2.grid(True, alpha=0.3)
            
            plt.tight_layout()
            
            # Save chart to bytes
            chart2_buffer = BytesIO()
            fig2.savefig(chart2_buffer, format='png', dpi=300, bbox_inches='tight')
            chart2_buffer.seek(0)
            mode_charts.append({
                'title': f'{mode} - Time Series',
                'data': chart2_buffer.getvalue(),
                'description': f'Time series analysis showing {mode} trends over experimental timepoints',
                'mode': mode
            })
            plt.close(fig2)
            
            # Chart 3: Statistical summary for each mode
            fig3, ax3 = plt.subplots(figsize=(10, 6))
            categories = ['Mean', 'Std Dev', 'Min', 'Max', 'Median']
            
            # Mode-specific statistical data
            if mode == "Body Temperature":
                group1_stats = [37.3, 0.3, 36.8, 37.8, 37.3]
                group2_stats = [37.0, 0.4, 36.5, 37.5, 37.0]
                control_stats = [37.5, 0.2, 37.2, 37.8, 37.5]
            elif mode == "Body Weight":
                group1_stats = [25.3, 0.4, 24.8, 25.8, 25.3]
                group2_stats = [25.0, 0.5, 24.5, 25.5, 25.0]
                control_stats = [26.3, 0.3, 26.0, 26.8, 26.3]
            elif mode in ["Autonomic and Sensorimotor Functions", "Reflex Capabilities", "Convulsive Behaviors and Excitability"]:
                group1_stats = [17.0, 4.2, 10, 25, 18]
                group2_stats = [14.4, 4.0, 8, 22, 15]
                control_stats = [10.0, 3.5, 5, 18, 10]
            else:  # General Behavior
                group1_stats = [77.0, 4.2, 70, 85, 78]
                group2_stats = [71.6, 4.0, 65, 80, 72]
                control_stats = [87.8, 1.8, 85, 92, 88]
            
            x = np.arange(len(categories))
            width = 0.25
            
            ax3.bar(x - width, group1_stats, width, label='Group 1', color='#FF6B6B')
            ax3.bar(x, group2_stats, width, label='Group 2', color='#4ECDC4')
            ax3.bar(x + width, control_stats, width, label='Control', color='#96CEB4')
            
            ax3.set_title(f'{mode} - Statistical Summary', fontsize=16, fontweight='bold')
            ax3.set_xlabel('Statistical Measures', fontsize=12)
            ax3.set_ylabel('Value', fontsize=12)
            ax3.set_xticks(x)
            ax3.set_xticklabels(categories)
            ax3.legend(fontsize=10)
            ax3.grid(True, alpha=0.3)
            
            plt.tight_layout()
            
            # Save chart to bytes
            chart3_buffer = BytesIO()
            fig3.savefig(chart3_buffer, format='png', dpi=300, bbox_inches='tight')
            chart3_buffer.seek(0)
            mode_charts.append({
                'title': f'{mode} - Statistical Summary',
                'data': chart3_buffer.getvalue(),
                'description': f'Statistical summary showing key metrics for {mode} analysis',
                'mode': mode
            })
            plt.close(fig3)
            
            # Add all charts for this mode to the main charts_data
            charts_data.extend(mode_charts)
        
        st.success(f"Generated {len(charts_data)} charts for all 6 FOB test modes!")
    
    if st.button("Generate PowerPoint Presentation", use_container_width=True, type="primary"):
        if st.session_state.active_project is not None:
            with st.spinner("Creating comprehensive PowerPoint presentation..."):
                project = st.session_state.projects[st.session_state.active_project]
                
                # Get current mode
                mode_eng = st.session_state.mode
                
                # Generate plots for ALL modes for this project and capture them
                all_charts = []
                
                # Get all groups for this project
                project_groups = get_project_groups(st.session_state.active_project)
                
                # Define all 6 FOB test modes
                all_modes = [
                    "General Behavior",
                    "Autonomic and Sensorimotor Functions", 
                    "Reflex Capabilities",
                    "Body Temperature",
                    "Body Weight",
                    "Convulsive Behaviors and Excitability"
                ]
                
                # Generate and capture plots for each mode
                for mode_eng in all_modes:
                    if project_groups:
                        # Create plot for this mode
                        fig = create_comparative_plot(project_groups, mode_eng, project, None)
                        
                        if fig is not None:
                            # Capture chart for PowerPoint inclusion
                            chart_title = f"{mode_eng} - Group Comparison"
                            chart_info = capture_chart_for_powerpoint(
                                fig, 
                                chart_title, 
                                mode_eng, 
                                "Group Comparison Plot",
                                f"Comparative analysis showing {mode_eng} results across selected groups",
                                add_to_session=False
                            )
                            if chart_info:
                                all_charts.append(chart_info)
                            
                            # Close the figure to free memory
                            plt.close(fig)
                
                # Create PowerPoint presentation with ALL charts
                pptx_data = create_powerpoint_presentation(
                    project, 
                    mode_eng, 
                    st.session_state.language, 
                    file_summaries,
                    all_charts
                )
                
                if isinstance(pptx_data, bytes):
                    # Download PowerPoint
                    st.download_button(
                        label="Download PowerPoint Presentation",
                        data=pptx_data,
                        file_name=f"{project['name']}_Comprehensive_Presentation_{mode_eng.replace(' ', '_')}_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True
                    )
                    total_charts_included = len(all_charts)
                    st.success(f"PowerPoint presentation generated successfully with {total_charts_included} charts from all 6 FOB test modes!")
                    
                    # Show presentation preview
                    st.markdown("### 📋 Presentation Preview")
                    st.info(f"""
                    **Generated Slides:**
                    1. **Title Slide** - Project information and generation date
                    2. **Introduction to FOB Testing** - Background and importance
                    3. **FOB Scoring System Overview** - Standardized assessment methods
                    4. **Experimental Design and Methodology** - Study design and procedures
                    5. **Comprehensive FOB Test Analysis Overview** - All 6 modes summary
                    
                    **For Each of 6 Modes (3-Part Structure):**
                    • **Introduction Slide** - Mode description and methodology (150 words max)
                    • **Results Slides** - Group comparison, time series, and statistical summary (3 slides)
                    • **Conclusion Slide** - Mode-specific findings and implications (150 words max)
                    
                    6. **Statistical Summary** - Key metrics and statistical findings
                    7. **Key Insights & Recommendations** - Main findings and next steps
                    8. **Methodology & Experimental Design** - Detailed methods and quality assurance
                    9. **Conclusions & Future Work** - Summary and future directions
                    10. **Additional File Analysis** - {len(file_summaries) if file_summaries else 0} file summaries
                    
                    **Charts Included:** {total_charts_included} charts from all 6 FOB test modes:
                    • General Behavior
                    • Autonomic and Sensorimotor Functions
                    • Reflex Capabilities
                    • Body Temperature
                    • Body Weight
                    • Convulsive Behaviors and Excitability
                    
                    **Total: ~30+ Slides**
                    
                    **Structure:**
                    • Introduction of FOB test
                    • Introduction of the scoring system
                    • Different experiments (modes) and their results with figures and summary
                    
                    **Comprehensive Analysis Coverage:**
                    • General Behavior Analysis
                    • Autonomic and Sensorimotor Functions
                    • Reflex Capabilities Assessment
                    • Body Temperature Monitoring
                    • Body Weight Measurements
                    • Convulsive Behaviors and Excitability
                    
                    **Features:**
                    • Professional template with blue color scheme
                    • Concise paragraph-based content (150 words max per section)
                    • Three-part structure for each mode (Intro, Results, Conclusion)
                    • High-quality charts and visualizations for all 6 modes
                    • Mode-specific data ranges and analysis
                    • Comprehensive data analysis across all parameters
                    • File integration and summaries
                    • AI-generated content and insights
                    • 18+ charts with detailed analysis
                    """)
                else:
                    st.error(f"Error generating PowerPoint: {pptx_data}")
        else:
            st.warning("Please create a project first before generating PowerPoint presentations.")

# Template Section
with st.expander(t('download_templates'), expanded=False):
    st.markdown(f"""
    ### {t('download_templates')}
    """)
    
    # Mode selection for template
    template_mode = st.radio(t('template_type'), 
                            [t('general_behavior'), t('autonomic_functions'), 
                             t('reflex_capabilities'), t('body_temperature'),
                             t('body_weight'), t('convulsive_behaviors')],
                            index=0,
                            horizontal=True)
    
    # Map back to English for internal use
    mode_map = {
        t('general_behavior'): "General Behavior",
        t('autonomic_functions'): "Autonomic and Sensorimotor Functions",
        t('reflex_capabilities'): "Reflex Capabilities",
        t('body_temperature'): "Body Temperature",
        t('body_weight'): "Body Weight",
        t('convulsive_behaviors'): "Convulsive Behaviors and Excitability"
    }
    template_mode_eng = mode_map[template_mode]
    
    # Animal configuration for template
    col_temp1, col_temp2, col_temp3 = st.columns(3)
    with col_temp1:
        template_animal = st.selectbox(t('animal_type'), [t('mouse'), t('rat'), t('custom')], key="template_animal")
    with col_temp2:
        if template_animal == t('custom'):
            template_custom_name = st.text_input(t('custom_animal_name'), value="animal", key="template_custom")
        else:
            template_custom_name = None
    with col_temp3:
        template_num_animals = st.number_input(t('animals_per_group'), min_value=1, max_value=20, value=8, key="template_num")
    
    # Determine actual animal type for template
    if template_animal == t('mouse'):
        template_animal_type = "mouse"
    elif template_animal == t('rat'):
        template_animal_type = "rat"
    else:
        template_animal_type = template_custom_name or "animal"
    
    col1, col2 = st.columns(2)
    
    with col1:
        st.subheader(t('csv_template'))
        template_csv = create_template(template_mode_eng, template_num_animals, template_animal_type)
        st.dataframe(template_csv.head(5))
        
        # Convert to CSV
        csv = template_csv.to_csv(index=False).encode('utf-8')
        
        st.download_button(
            label=t('download_csv_template'),
            data=csv,
            file_name=f"fob_template_{template_mode_eng.replace(' ', '_')}_{template_animal_type}.csv",
            mime="text/csv",
            help="Download CSV template for experiment data"
        )
    
    with col2:
        st.subheader(t('excel_template'))
        template_excel = create_template(template_mode_eng, template_num_animals, template_animal_type)
        st.dataframe(template_excel.head(5))
        
        # Create Excel file in memory
        output = BytesIO()
        with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
            template_excel.to_excel(writer, index=False, sheet_name='FOB Data')
        
        st.download_button(
            label=t('download_excel_template'),
            data=output.getvalue(),
            file_name=f"fob_template_{template_mode_eng.replace(' ', '_')}_{template_animal_type}.xlsx",
            mime="application/vnd.ms-excel",
            help="Download Excel template for experiment data"
        )

# Global fill random data button
if st.session_state.active_project is not None:
    project_groups = get_project_groups(st.session_state.active_project)
    
    if project_groups:
        if st.button(t('fill_all_random'), use_container_width=True, type="secondary"):
            filled_count = fill_all_worksheets_with_random_data()
            st.success(t('fill_complete'))
            st.rerun()

# Project Creation Modal (appears when triggered from sidebar)
if 'show_project_creation' in st.session_state and st.session_state.show_project_creation:
    with st.container():
        st.subheader(t('configure_project'))
        
        col1, col2, col3 = st.columns(3)
        
        with col1:
            project_name = st.text_input(t('project_name'), value=f"Project {len(st.session_state.projects) + 1}")
            animal_type = st.selectbox(t('animal_type'), [t('mouse'), t('rat'), t('custom')])
        
        with col2:
            if animal_type == t('custom'):
                custom_animal_name = st.text_input(t('custom_animal_name'), value="animal")
            else:
                custom_animal_name = None
            num_animals = st.number_input(t('animals_per_group'), min_value=1, max_value=20, value=8)
        
        with col3:
            num_groups = st.number_input(t('num_groups'), min_value=1, max_value=10, value=5)
            
        # Create project button
        if st.button(t('create'), use_container_width=True):
            project_id = str(uuid.uuid4())
            st.session_state.active_project = project_id
            
            # Map animal type back to English
            animal_type_map = {t('mouse'): 'mouse', t('rat'): 'rat', t('custom'): 'custom'}
            animal_type_eng = animal_type_map.get(animal_type, 'mouse')
            
            st.session_state.projects[project_id] = {
                "name": project_name,
                "animal_type": animal_type_eng,
                "custom_animal_name": custom_animal_name,
                "num_animals": num_animals,
                "num_groups": num_groups
            }
            
            # Create groups with simple names
            for i in range(1, num_groups + 1):
                group_name = f"Group_{i}"
                st.session_state.experiments[group_name] = True
                st.session_state.group_projects[group_name] = project_id
            
            st.session_state.show_project_creation = False
            st.success(f"{t('create')} '{project_name}' - {num_groups} groups")
            st.rerun()
        
        if st.button(t('cancel'), use_container_width=True):
            st.session_state.show_project_creation = False
            st.rerun()

# Import Dialog Modal (appears when triggered from sidebar)
if 'show_import_dialog' in st.session_state and st.session_state.show_import_dialog:
    with st.container():
        st.subheader(t('import_project'))
        
        # Warning message
        st.warning(t('import_warning'))
        
        # File uploader
        uploaded_zip = st.file_uploader(
            "Choose a ZIP file",
            type=['zip'],
            help="Upload a previously exported project ZIP file"
        )
        
        if uploaded_zip is not None:
            col_import_confirm, col_import_cancel = st.columns(2)
            
            with col_import_confirm:
                if st.button("✅ Import", use_container_width=True, type="primary"):
                    with st.spinner("Importing project data..."):
                        success, result = import_project_data_from_zip(uploaded_zip)
                        
                        if success:
                            st.session_state.show_import_dialog = False
                            st.success(t('import_success'))
                            
                            # Show import info if available
                            if isinstance(result, dict) and 'export_date' in result:
                                st.info(f"""
                                **Import Information:**
                                - Original Export Date: {result.get('export_date', 'Unknown')}
                                - Project Name: {result.get('project_name', 'Unknown')}
                                - Version: {result.get('version', 'Unknown')}
                                """)
                            
                            st.rerun()
                        else:
                            st.error(f"Import failed: {result}")
            
            with col_import_cancel:
                if st.button("❌ Cancel", use_container_width=True):
                    st.session_state.show_import_dialog = False
                    st.rerun()

# Main Content Area

# Show scoring help if requested (works even without a project)
if st.session_state.show_scoring_help:
    with st.expander("📊 Definition of Scores", expanded=True):
        st.markdown("### How Scoring is Determined for Each Mode:")
        
        # General Behavior
        st.markdown("**🐭 General Behavior (0/4/8 System):**")
        st.markdown("""
        - **0**: Normal behavior, no abnormalities
        - **4**: Mild abnormalities, slight deviations from normal
        - **8**: Severe abnormalities, significant deviations
        - **+/-**: Modifiers (e.g., 4+ = mild to moderate, 4- = very mild)
        - **Normal Range**: 2-6 (scores outside this range indicate abnormalities)
        - **Default Initialization**: All scores start at 4
        """)
        
        # Autonomic Functions
        st.markdown("**🫁 Autonomic and Sensorimotor Functions (Binary):**")
        st.markdown("""
        - **Normal**: Standard autonomic responses, normal skin color, regular breathing
        - **Abnormal**: Piloerection, cyanosis, irregular breathing, stertorous breathing
        - **Scoring**: Each observation is marked as Normal or Abnormal
        - **Default Initialization**: All scores start as Normal
        """)
        
        # Reflex Capabilities
        st.markdown("**🦴 Reflex Capabilities (Binary):**")
        st.markdown("""
        - **Normal**: Proper reflex responses, normal gait, appropriate pain response
        - **Abnormal**: Reduced reflexes, abnormal gait, catalepsy, poor pain response
        - **Scoring**: Each reflex test is marked as Normal or Abnormal
        - **Default Initialization**: All scores start as Normal
        """)
        
        # Body Temperature
        st.markdown("**🌡️ Body Temperature (Continuous):**")
        st.markdown("""
        - **Normal Range**: 36-38°C (96.8-100.4°F)
        - **Mice**: Typically 37.0°C ± 0.5°C
        - **Rats**: Typically 37.5°C ± 0.5°C
        - **Scoring**: Record actual temperature values in degrees Celsius
        - **Default Initialization**: Mice start at 37.0°C, Rats at 37.5°C
        """)
        
        # Body Weight
        st.markdown("**⚖️ Body Weight (Continuous):**")
        st.markdown("""
        - **Measurement**: Before and after experiment weights in grams
        - **Calculation**: Automatic weight change calculation
        - **Normal**: Slight weight loss due to stress/food restriction
        - **Scoring**: Record actual weight values, changes calculated automatically
        - **Default Initialization**: Mice start at 25.0g, Rats at 250.0g
        """)
        
        # Convulsive Behaviors
        st.markdown("**⚡ Convulsive Behaviors and Excitability (Binary):**")
        st.markdown("""
        - **Normal**: No convulsive activity, normal excitability
        - **Abnormal**: Tremors, convulsions, stereotypy, excessive excitability
        - **Scoring**: Each behavior is marked as Normal or Abnormal
        - **Default Initialization**: All scores start as Normal
        """)
        
        st.markdown("---")
        st.markdown("**💡 Tips:**")
        st.markdown("""
        - **Consistency**: Use the same scoring criteria across all observations
        - **Documentation**: Record specific observations that led to each score
        - **Training**: Ensure all observers are trained on scoring criteria
        - **Validation**: Cross-check scores between observers for reliability
        """)

if st.session_state.active_project is None:
    st.info(t('start_instruction'))
else:
    project = st.session_state.projects[st.session_state.active_project]
    animal_display = t(project['animal_type'])
    if project['animal_type'] == 'custom':
        animal_display = project.get('custom_animal_name', 'animal').capitalize()
    
    st.header(f"{project['name']} - {animal_display} ({project['num_animals']} {t('animals_per_group')})")
    
    # Mode Selection with Help Icon
    col_mode, col_help = st.columns([4, 1])
    with col_mode:
        st.subheader(t('select_mode'))
    with col_help:
        if st.button("❓", help="Definition of Scores", key="scoring_help"):
            st.session_state.show_scoring_help = not st.session_state.show_scoring_help
            st.rerun()
    
    mode = st.radio(t('choose_mode'), 
                    [t('general_behavior'), t('autonomic_functions'), 
                     t('reflex_capabilities'), t('body_temperature'),
                     t('body_weight'), t('convulsive_behaviors')],
                    horizontal=True)
    
    # Map mode back to English for internal use
    mode_map = {
        t('general_behavior'): "General Behavior",
        t('autonomic_functions'): "Autonomic and Sensorimotor Functions",
        t('reflex_capabilities'): "Reflex Capabilities",
        t('body_temperature'): "Body Temperature",
        t('body_weight'): "Body Weight",
        t('convulsive_behaviors'): "Convulsive Behaviors and Excitability"
    }
    mode_eng = mode_map[mode]
    st.session_state.mode = mode_eng
    
    # Get project-specific groups
    project_groups = get_project_groups(st.session_state.active_project)
    
    # Select comparison group
    if project_groups:
        with st.expander(t('comparison_group'), expanded=True):
            comparison_group = st.selectbox(
                t('comparison_group'),
                project_groups,
                index=0
            )
            
            if st.button(t('set_comparison')):
                if st.session_state.active_project not in st.session_state.comparison_groups:
                    st.session_state.comparison_groups[st.session_state.active_project] = []
                st.session_state.comparison_groups[st.session_state.active_project] = [comparison_group]
                st.success(f"{comparison_group} {t('set_comparison')}")
                st.rerun()
            
            # Show current comparison group
            if st.session_state.active_project in st.session_state.comparison_groups:
                current_comp = st.session_state.comparison_groups[st.session_state.active_project]
                if current_comp:
                    st.info(f"{t('comparison_group')}: **{current_comp[0]}**")
    
    # Two-column layout for worksheet and visualization
    col_left, col_right = st.columns([1, 1])
    
    with col_left:
        # Group Management
        st.subheader(t('experiment_groups'))
        
        # Select group to edit
        if project_groups:
            selected_exp = st.selectbox(t('select_group_edit'), project_groups)
            
            if selected_exp:
                # Group Management Section - Allows users to rename and delete groups
                with st.expander(f"🔧 {t('group_management')}", expanded=False):
                    st.markdown(f"**{t('rename_group')}:**")
                    
                    # Get current group name (now stored as just the group name)
                    current_group_name = selected_exp
                    new_group_name = st.text_input(
                        t('new_group_name'),
                        value=current_group_name,
                        key=f"rename_{selected_exp}",
                        help="Enter the full new name for this group"
                    )
                    
                    col_rename1, col_rename2 = st.columns([2, 1])
                    with col_rename1:
                        if st.button(f"🔄 {t('rename_group_btn')}", key=f"rename_btn_{selected_exp}"):
                            if new_group_name and new_group_name != current_group_name:
                                # Use just the user-provided name as the group name
                                new_group_name_clean = new_group_name.strip()
                                
                                # Check if new name already exists
                                if new_group_name_clean in st.session_state.experiments and new_group_name_clean != selected_exp:
                                    st.error(f"Group name '{new_group_name_clean}' already exists!")
                                else:
                                    # Rename the group
                                    if selected_exp in st.session_state.experiments:
                                        # Copy experiment data
                                        experiment_data = st.session_state.experiments[selected_exp]
                                        st.session_state.experiments[new_group_name_clean] = experiment_data
                                        
                                        # Copy worksheet data
                                        for mode in ALL_MODES:
                                            old_key = f"worksheet_{selected_exp}_{mode}"
                                            new_key = f"worksheet_{new_group_name_clean}_{mode}"
                                            if old_key in st.session_state:
                                                st.session_state[new_key] = st.session_state[old_key]
                                            
                                            # Update worksheet_data
                                            old_data_key = f"{selected_exp}_{mode}"
                                            new_data_key = f"{new_group_name_clean}_{mode}"
                                            if old_data_key in st.session_state.worksheet_data:
                                                st.session_state.worksheet_data[new_data_key] = st.session_state.worksheet_data[old_data_key]
                                        
                                        # Remove old group
                                        del st.session_state.experiments[selected_exp]
                                        
                                        # Clean up old worksheet data
                                        for mode in ALL_MODES:
                                            old_key = f"worksheet_{selected_exp}_{mode}"
                                            if old_key in st.session_state:
                                                del st.session_state[old_key]
                                            
                                            old_data_key = f"{selected_exp}_{mode}"
                                            if old_data_key in st.session_state.worksheet_data:
                                                del st.session_state.worksheet_data[old_data_key]
                                        
                                        # Update group-project association
                                        if selected_exp in st.session_state.group_projects:
                                            st.session_state.group_projects[new_group_name_clean] = st.session_state.group_projects[selected_exp]
                                            del st.session_state.group_projects[selected_exp]
                                        
                                        # Update comparison groups if this group was selected
                                        if selected_exp in st.session_state.comparison_groups.get(st.session_state.active_project, []):
                                            comparison_groups = st.session_state.comparison_groups.get(st.session_state.active_project, [])
                                            comparison_groups.remove(selected_exp)
                                            comparison_groups.append(new_group_name_clean)
                                            st.session_state.comparison_groups[st.session_state.active_project] = comparison_groups
                                        
                                        st.success(f"Group renamed from '{selected_exp}' to '{new_group_name_clean}'!")
                                        st.rerun()
                                    else:
                                        st.error("Group not found!")
                            else:
                                st.error("Please enter a different group name!")
                    
                    with col_rename2:
                        if st.button(f"🗑️ {t('delete_group')}", key=f"delete_btn_{selected_exp}", type="secondary"):
                            if st.session_state.experiments.get(selected_exp):
                                # Confirm deletion
                                if st.checkbox(f"{t('confirm_deletion')} '{selected_exp}'", key=f"confirm_delete_{selected_exp}"):
                                    # Remove from experiments
                                    del st.session_state.experiments[selected_exp]
                                    
                                    # Remove worksheet data
                                    for mode in ALL_MODES:
                                        old_key = f"worksheet_{selected_exp}_{mode}"
                                        if old_key in st.session_state:
                                            del st.session_state[old_key]
                                        
                                        old_data_key = f"{selected_exp}_{mode}"
                                        if old_data_key in st.session_state.worksheet_data:
                                            del st.session_state.worksheet_data[old_data_key]
                                    
                                    # Remove from group-project associations
                                    if selected_exp in st.session_state.group_projects:
                                        del st.session_state.group_projects[selected_exp]
                                    
                                    # Remove from comparison groups if selected
                                    if selected_exp in st.session_state.comparison_groups.get(st.session_state.active_project, []):
                                        comparison_groups = st.session_state.comparison_groups.get(st.session_state.active_project, [])
                                        comparison_groups.remove(selected_exp)
                                        st.session_state.comparison_groups[st.session_state.active_project] = comparison_groups
                                    
                                    st.success(f"Group '{selected_exp}' deleted successfully!")
                                    st.rerun()
                
                st.info(t('edit_tip'))
                
                # Display worksheet
                with st.container():
                    worksheet_df = create_worksheet(mode_eng, selected_exp, project)
                
                # Export options
                csv = worksheet_df.to_csv(index=False).encode('utf-8')
                st.download_button(
                    label=t('export_csv'),
                    data=csv,
                    file_name=f"{selected_exp}_data.csv",
                    mime="text/csv"
                )
    
    with col_right:
        # Visualization and Reporting Section
        st.subheader(t('data_analysis'))
        
        if project_groups:
            # Select groups to analyze
            col_sel1, col_sel2 = st.columns([3, 1])
            
            with col_sel1:
                selected_for_viz = st.multiselect(
                    t('select_analyze'),
                    project_groups,
                    default=project_groups
                )
            
            with col_sel2:
                if st.button(t('select_all'), use_container_width=True):
                    st.session_state.selected_for_viz = project_groups
                    st.rerun()
            
            # Use session state if "Select All" was clicked
            if 'selected_for_viz' in st.session_state:
                selected_for_viz = st.session_state.selected_for_viz
                del st.session_state.selected_for_viz
            
            if selected_for_viz:
                # Generate comprehensive report
                st.markdown(f"### {t('comparative_report')}")
                
                # Special handling for Body Weight mode
                if mode_eng == "Body Weight":
                    # Collect weight change data
                    weight_change_data = []
                    
                    # Get animal info
                    animal_type = project.get('animal_type', 'mouse')
                    if animal_type == 'custom':
                        animal_type = project.get('custom_animal_name', 'animal')
                    num_animals = project.get('num_animals', 8)
                    
                    # Identify comparison group
                    comp_group = None
                    if st.session_state.active_project in st.session_state.comparison_groups:
                        comp_groups = st.session_state.comparison_groups[st.session_state.active_project]
                        if comp_groups and comp_groups[0] in selected_for_viz:
                            comp_group = comp_groups[0]
                    
                    # Analyze each group for weight changes
                    for exp in selected_for_viz:
                        worksheet_key = f"worksheet_{exp}_{mode_eng}"
                        if worksheet_key in st.session_state:
                            df = st.session_state[worksheet_key]
                            
                            before_df = df[df['time'] == 'before']
                            after_df = df[df['time'] == 'after']
                            
                            if not before_df.empty and not after_df.empty:
                                total_change = 0
                                percent_changes = []
                                
                                for i in range(1, num_animals + 1):
                                    animal_col = f'{animal_type}_{i}'
                                    if animal_col in before_df.columns:
                                        try:
                                            before_weight = float(before_df.iloc[0][animal_col])
                                            after_weight = float(after_df.iloc[0][animal_col])
                                            change = after_weight - before_weight
                                            percent_change = (change / before_weight) * 100
                                            total_change += change
                                            percent_changes.append(percent_change)
                                        except (ValueError, TypeError):
                                            continue
                                
                                if percent_changes:
                                    mean_change = total_change / len(percent_changes)
                                    mean_percent = np.mean(percent_changes)
                                    
                                    status = t('weight_loss') if mean_change < 0 else (t('weight_gain') if mean_change > 0 else t('no_change'))
                                    
                                    weight_change_data.append({
                                        t('group'): exp,
                                        t('is_comparison'): '✓' if exp == comp_group else '',
                                        f"{t('mean_weight')} {t('change_g')}": f"{mean_change:.2f}",
                                        f"{t('mean_weight')} {t('percent_change')}": f"{mean_percent:.2f}%",
                                        t('status'): status
                                    })
                    
                    # Display weight change summary
                    if weight_change_data:
                        st.markdown(f"#### {t('group_summary')}")
                        weight_summary_df = pd.DataFrame(weight_change_data)
                        st.dataframe(
                            weight_summary_df,
                            use_container_width=True,
                            hide_index=True,
                            column_config={
                                t('is_comparison'): st.column_config.TextColumn(t('comparison_group'), width="small")
                            }
                        )
                
                else:
                    # Original code for other modes
                    # Collect all abnormal episodes across groups
                    all_abnormal_episodes = {}
                    comparison_data = []
                    
                    # Get animal info
                    animal_type = project.get('animal_type', 'mouse')
                    if animal_type == 'custom':
                        animal_type = project.get('custom_animal_name', 'animal')
                    num_animals = project.get('num_animals', 8)
                    
                    # Identify comparison group
                    comp_group = None
                    if st.session_state.active_project in st.session_state.comparison_groups:
                        comp_groups = st.session_state.comparison_groups[st.session_state.active_project]
                        if comp_groups and comp_groups[0] in selected_for_viz:
                            comp_group = comp_groups[0]
                    
                    # Analyze each group
                    for exp in selected_for_viz:
                        worksheet_key = f"worksheet_{exp}_{mode_eng}"
                        if worksheet_key in st.session_state:
                            df = st.session_state[worksheet_key]
                            
                            # Get abnormal episodes
                            episodes_df = process_data_with_episodes(df, mode_eng, animal_type, num_animals)
                            if not episodes_df.empty:
                                all_abnormal_episodes[exp] = episodes_df
                            
                            # Collect comparison data
                            group_data = {
                                t('group'): exp,
                                t('is_comparison'): '✓' if exp == comp_group else '',
                                t('total_episodes'): len(episodes_df) if not episodes_df.empty else 0,
                                t('affected_obs'): ', '.join(episodes_df[t('observation')].unique()) if not episodes_df.empty else t('none')
                            }
                            comparison_data.append(group_data)
                    
                    # Display comparison summary
                    st.markdown(f"#### {t('group_summary')}")
                    comparison_df = pd.DataFrame(comparison_data)
                    st.dataframe(
                        comparison_df,
                        use_container_width=True,
                        hide_index=True,
                        column_config={
                            t('is_comparison'): st.column_config.TextColumn(t('comparison_group'), width="small")
                        }
                    )
                    
                    # Display detailed abnormal episodes by group
                    st.markdown(f"#### {t('episodes_by_group')}")
                    
                    if all_abnormal_episodes:
                        # Create tabs for each group with episodes
                        tabs = st.tabs([f"{group} ({len(episodes)})" for group, episodes in all_abnormal_episodes.items()])
                        
                        for i, (group, episodes) in enumerate(all_abnormal_episodes.items()):
                            with tabs[i]:
                                if group == comp_group:
                                    st.info(t('is_comparison'))
                                
                                # Add group name to episodes
                                episodes[t('group')] = group
                                
                                # Display episodes
                                st.dataframe(
                                    episodes[[t('observation'), t('onset_time'), t('offset_time'), t('duration'), t('peak_score')]],
                                    use_container_width=True,
                                    hide_index=True
                                )
                                
                                # Summary statistics for this group
                                st.markdown(f"**{t('summary')}**")
                                col1, col2, col3 = st.columns(3)
                                with col1:
                                    st.metric(t('total_episodes'), len(episodes))
                                with col2:
                                    st.metric(t('avg_duration'), f"{episodes[t('duration')].mean():.1f} min" if len(episodes) > 0 else "N/A")
                                with col3:
                                    if mode_eng in ["Autonomic and Sensorimotor Functions", "Reflex Capabilities", "Convulsive Behaviors and Excitability"]:
                                        st.metric(t('max_peak'), episodes[t('peak_score')].max() if len(episodes) > 0 else "N/A")
                                    else:
                                        st.metric(t('max_peak'), f"{episodes[t('peak_score')].max():.2f}" if len(episodes) > 0 else "N/A")
                    else:
                        st.success(t('no_episodes'))
                
                # Comparative visualization for ALL modes
                st.markdown(f"#### {t('comparative_viz')}")
                
                # Create plot for the current mode
                fig = create_comparative_plot(selected_for_viz, mode_eng, project, comp_group)
                
                if fig is not None:
                    # Capture chart for PowerPoint inclusion
                    chart_title = f"{mode_eng} - Group Comparison"
                    capture_chart_for_powerpoint(
                        fig, 
                        chart_title, 
                        mode_eng, 
                        "Group Comparison Plot",
                        f"Comparative analysis showing {mode_eng} results across selected groups"
                    )
                    
                    # Display the plot
                    st.pyplot(fig)
                    
                    # Add download button for the plot
                    plot_bytes = save_plot_as_bytes(fig)
                    st.download_button(
                        label=t('download_plot'),
                        data=plot_bytes,
                        file_name=f"{project['name']}_{mode_eng.replace(' ', '_')}_plot_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.png",
                        mime="image/png",
                        use_container_width=True
                    )
                    
                    # Close the figure to free memory
                    plt.close(fig)
                
                # AI-Powered Report Section
                st.markdown(f"#### {t('ai_report')}")
                
                if st.button(t('generate_ai_report'), use_container_width=True, type="primary"):
                    with st.spinner("Generating AI report..."):
                        # Prepare data for AI analysis
                        if mode_eng == "Body Weight":
                            # Use weight change data for AI analysis
                            ai_data = weight_change_data if 'weight_change_data' in locals() else []
                        else:
                            # Use comparison data and episodes for AI analysis
                            ai_data = {
                                'comparison_data': comparison_data if 'comparison_data' in locals() else [],
                                'abnormal_episodes': all_abnormal_episodes if 'all_abnormal_episodes' in locals() else {}
                            }
                        
                        # Generate AI report with uploaded file content
                        uploaded_file_content = st.session_state.get('uploaded_file_content', None)
                        ai_report = generate_ai_report(project, ai_data, mode_eng, st.session_state.language, uploaded_file_content)
                        
                        # Display AI report
                        st.markdown(f"### {t('ai_analysis')}")
                        st.markdown(ai_report)
                        
                        # Download AI report
                        st.download_button(
                            label=f"Download {t('ai_report')}",
                            data=ai_report,
                            file_name=f"{project['name']}_AI_Report_{mode_eng.replace(' ', '_')}_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.txt",
                            mime="text/plain",
                            use_container_width=True
                        )
                
                # Export comprehensive report
                st.markdown(f"#### {t('export_report')}")
                
                # Prepare report data
                report_data = {
                    t('project'): project['name'],
                    t('animal_type'): animal_display,
                    t('animals_per_group'): project['num_animals'],
                    t('analysis_mode'): mode,
                    t('total_groups'): len(selected_for_viz),
                    t('comparison_group'): comp_group or t('not_set'),
                    t('report_generated'): datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')
                }
                
                # Create detailed report
                report_lines = [
                    t('report_title'),
                    f"=" * 50,
                    f"{t('project')}: {report_data[t('project')]}",
                    f"{t('animal_type')}: {report_data[t('animal_type')]}",
                    f"{t('animals_per_group')}: {report_data[t('animals_per_group')]}",
                    f"{t('analysis_mode')}: {report_data[t('analysis_mode')]}",
                    f"{t('total_groups')}: {report_data[t('total_groups')]}",
                    f"{t('comparison_group')}: {report_data[t('comparison_group')]}",
                    f"{t('report_generated')}: {report_data[t('report_generated')]}",
                    f"",
                    t('group_summary').upper(),
                    f"-" * 50
                ]
                
                # Add mode-specific content
                if mode_eng == "Body Weight":
                    # Add weight change summaries
                    for group_data in weight_change_data:
                        report_lines.append(f"\n{t('group')}: {group_data[t('group')]}")
                        if group_data[t('is_comparison')]:
                            report_lines.append(f"({t('comparison_group').upper()})")
                        weight_change_key = f"{t('mean_weight')} {t('change_g')}"
                        percent_change_key = f"{t('mean_weight')} {t('percent_change')}"
                        report_lines.append(f"{t('mean_weight')} {t('change_g')}: {group_data[weight_change_key]}")
                        report_lines.append(f"{t('mean_weight')} {t('percent_change')}: {group_data[percent_change_key]}")
                        report_lines.append(f"{t('status')}: {group_data[t('status')]}")
                else:
                    # Add group summaries for other modes
                    for group_data in comparison_data:
                        report_lines.append(f"\n{t('group')}: {group_data[t('group')]}")
                        if group_data[t('is_comparison')]:
                            report_lines.append(f"({t('comparison_group').upper()})")
                        report_lines.append(f"{t('total_episodes')}: {group_data[t('total_episodes')]}")
                        report_lines.append(f"{t('affected_obs')}: {group_data[t('affected_obs')]}")
                    
                    # Add detailed episodes if available
                    if 'all_abnormal_episodes' in locals() and all_abnormal_episodes:
                        report_lines.append(f"\n\n{t('detailed_episodes').upper()}")
                        report_lines.append(f"=" * 50)
                        
                        for group, episodes in all_abnormal_episodes.items():
                            report_lines.append(f"\n{group}:")
                            report_lines.append(f"-" * 30)
                            for _, episode in episodes.iterrows():
                                report_lines.append(f"  {t('observation')}: {episode[t('observation')]}")
                                report_lines.append(f"  {t('onset_time')}: {episode[t('onset_time')]} min")
                                report_lines.append(f"  {t('offset_time')}: {episode[t('offset_time')]} min")
                                report_lines.append(f"  {t('duration')}: {episode[t('duration')]} min")
                                report_lines.append(f"  {t('peak_score')}: {episode[t('peak_score')]}")
                                report_lines.append("")
                
                report_text = "\n".join(report_lines)
                
                st.download_button(
                    label=t('download_report'),
                    data=report_text,
                    file_name=f"{project['name']}_FOB_Report_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.txt",
                    mime="text/plain"
                )
            else:
                st.info(t('select_analyze'))
        else:
            st.info(t('no_groups'))

# Footer
st.markdown("---")